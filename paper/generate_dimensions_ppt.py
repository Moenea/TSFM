"""Generate PPT discussing the macroscopic dimensions of time series forecasting.

User's view: good forecasting needs (a) information completeness + (b) model
representation capacity. We argue this is correct but incomplete; we propose
a 6-dimension framework, then identify on which dimensions TSFM actually
makes its contribution.

Slide map (22 slides):
  P1   Title
  P2-3 User's two-axis framework — what is correct
  P4-9 What is missing — six macro dimensions
  P10-11 The unified 6-axis macro framework
  P12-17 Where TSFM actually makes its contribution
  P18-19 What TSFM does NOT solve
  P20-22 Engineering implication + final answer
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import math

# ---------- Style ----------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
DEEP = RGBColor(0x14, 0x2C, 0x52)
ACCENT = RGBColor(0xE8, 0x8B, 0x1A)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
GREY = RGBColor(0x4A, 0x55, 0x68)
DGREY = RGBColor(0x33, 0x3A, 0x47)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xDC, 0xE3, 0xF0)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)
PURPLE = RGBColor(0x6B, 0x3F, 0xA0)
TEAL = RGBColor(0x16, 0x7E, 0x8A)
ORANGE = RGBColor(0xE0, 0x6E, 0x18)
GOLD = RGBColor(0xD4, 0xA0, 0x2A)

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def _i(v):
    return int(v) if v is not None else v


def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_oval(slide, x, y, w, h, fill, line=None):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def rich(slide, x, y, w, h, blocks, line_spacing=1.3, anchor=MSO_ANCHOR.TOP):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, runs in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        for txt, size, color, bold in runs:
            r = p.add_run()
            r.text = txt
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return tb


def line(slide, x1, y1, x2, y2, color, w=1.8):
    x1, y1, x2, y2 = _i(x1), _i(y1), _i(x2), _i(y2)
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(w)
    return ln


def arrow(slide, x1, y1, x2, y2, color=ACCENT, w=2.0):
    x1, y1, x2, y2 = _i(x1), _i(y1), _i(x2), _i(y2)
    ln = slide.shapes.add_connector(2, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(w)
    return ln


def page(slide, idx, total, title, subtitle=None, section=None):
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_rect(slide, 0, 0, SW, Inches(0.08), ACCENT)
    add_rect(slide, 0, Inches(0.08), Inches(0.18), SH, NAVY)
    if section:
        add_text(slide, Inches(0.5), Inches(0.18), Inches(11.5), Inches(0.3),
                 section, size=11, bold=True, color=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.6),
             title, size=24, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(12.0), Inches(0.4),
                 subtitle, size=12, color=GREY, italic=True)
    add_rect(slide, Inches(0.5), Inches(1.28), Inches(0.6), Emu(28000), ACCENT)
    add_text(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
             f"{idx} / {total}", size=10, color=GREY, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.3),
             "时序预测的宏观维度 — TSFM 究竟在哪里发力 · Jinsong Zhao", size=10, color=GREY)


TOTAL = 22

# ============================================================
# Slide 1 — Title
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, 0, SW, Inches(0.15), ACCENT)
# decoration
for i in range(8):
    add_rect(s, Inches(0.5 + i * 1.6), Inches(0.4), Inches(0.6), Emu(15000), ACCENT)

add_text(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.6),
         "时序预测的宏观维度", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.5),
         "The Macro Dimensions of Time Series Forecasting",
         size=22, color=ACCENT, italic=True)

add_rect(s, Inches(0.8), Inches(3.15), Inches(0.7), Emu(30000), ACCENT)

rich(s, Inches(0.8), Inches(3.35), Inches(11.5), Inches(2.8), [
    [("用户的两轴模型: ", 18, ACCENT, True),
     ("信息完备性 × 模型表征能力", 18, WHITE, True)],
    [("→ ", 16, ACCENT, True),
     ("这两个维度都是对的, 但远远不够 — 还需要至少 4 个被忽略的宏观维度",
      16, SOFT, False)],
    [("", 8, WHITE, False)],
    [("TSFM 不是单一维度的胜利, 而是在 ", 16, SOFT, False),
     ("4 个维度", 16, ACCENT, True),
     (" 上的协同性创新", 16, SOFT, False)],
    [("但在另外 ", 16, SOFT, False),
     ("2 个维度", 16, RED, True),
     (" 上, TSFM 完全没有帮助 — 这就是它的合理使用边界", 16, SOFT, False)],
])

add_rect(s, Inches(0.8), Inches(6.5), Inches(11.5), Emu(15000), ACCENT)
add_text(s, Inches(0.8), Inches(6.65), Inches(11.5), Inches(0.4),
         "Jinsong Zhao    ·    清华大学    ·    2026-05-07",
         size=14, color=ACCENT)

# ============================================================
# Slide 2 — User's framework verbatim
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 2, TOTAL, "用户的两轴框架", "我们先把用户的认知精确写下来",
     section="PART I · 用户的认知")

# Quote box
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.6), LIGHT)
add_rect(s, Inches(0.5), Inches(1.55), Inches(0.18), Inches(1.6), ACCENT)
add_text(s, Inches(0.85), Inches(1.65), Inches(11.8), Inches(0.4),
         "用户原话:", size=13, bold=True, color=ACCENT)
rich(s, Inches(0.85), Inches(2.0), Inches(11.8), Inches(1.1), [
    [("「好的、准确的时序预测需要 ", 14, NAVY, False),
     ("足够完备的信息", 14, ACCENT, True),
     (" — 即足够多、足够完备的上下文; 另一方面, 还需要 ", 14, NAVY, False),
     ("模型的表征能力", 14, ACCENT, True),
     (" 足够强。」", 14, NAVY, False)],
])

# Two pillars
add_rect(s, Inches(0.7), Inches(3.5), Inches(5.9), Inches(3.3), DEEP)
add_rect(s, Inches(0.7), Inches(3.5), Inches(5.9), Inches(0.55), ACCENT)
add_text(s, Inches(0.7), Inches(3.5), Inches(5.9), Inches(0.55),
         "AXIS 1 · 信息完备性", size=15, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(0.95), Inches(4.2), Inches(5.5), Inches(2.5), [
    [("• ", 14, ACCENT, True), ("可观测性 (observability)", 14, WHITE, True)],
    [("  系统状态能否从观测中『还原』", 11, SOFT, False)],
    [("• ", 14, ACCENT, True), ("充分统计量 (sufficient statistic)", 14, WHITE, True)],
    [("  上下文 X 是否包含了 Y 的所有信息: H(Y|X) = H(Y|状态)", 11, SOFT, False)],
    [("• ", 14, ACCENT, True), ("协变量 / 外生变量", 14, WHITE, True)],
    [("  受迫系统中, 驱动信号必须可观测", 11, SOFT, False)],
])

add_rect(s, Inches(6.9), Inches(3.5), Inches(5.9), Inches(3.3), DEEP)
add_rect(s, Inches(6.9), Inches(3.5), Inches(5.9), Inches(0.55), ACCENT)
add_text(s, Inches(6.9), Inches(3.5), Inches(5.9), Inches(0.55),
         "AXIS 2 · 模型表征能力", size=15, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(7.15), Inches(4.2), Inches(5.5), Inches(2.5), [
    [("• ", 14, ACCENT, True), ("函数空间的覆盖能力", 14, WHITE, True)],
    [("  ∃ θ ∈ Θ, 使得 f_θ ≈ f* (universal approximation)", 11, SOFT, False)],
    [("• ", 14, ACCENT, True), ("分辨率 / 容量", 14, WHITE, True)],
    [("  参数量 N、深度、注意力头数 — 决定能拟合多复杂的关系", 11, SOFT, False)],
    [("• ", 14, ACCENT, True), ("非线性 / 多尺度建模", 14, WHITE, True)],
    [("  能否表达从秒到月的耦合", 11, SOFT, False)],
])

add_text(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
         "→ 这两个维度都是必要的; 但单凭这两个, 你解释不了为什么 TSFM 比传统方法更好。",
         size=12, italic=True, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 3 — Why these two are right (formal)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 3, TOTAL, "你说的两轴 · 形式化论证为什么是对的",
     "信息完备性 ↔ 不可约误差; 表征能力 ↔ 近似误差", section="PART I · 用户的认知")

# Decompose forecast error
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.7), LIGHT)
add_text(s, Inches(0.85), Inches(1.7), Inches(11.8), Inches(0.4),
         "预测误差的经典分解 (bias-variance + irreducible)",
         size=14, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(2.15), Inches(11.8), Inches(1.0), [
    [("E[(Y − Ŷ)²]  =  ", 16, NAVY, True),
     ("σ²(Y|X)", 16, RED, True),
     ("              +              ", 16, NAVY, False),
     ("‖f* − f̂‖²", 16, ACCENT, True),
     ("              +              ", 16, NAVY, False),
     ("Var[f̂]", 16, PURPLE, True)],
    [("           ", 12, NAVY, False),
     ("不可约 (信息不完备)", 12, RED, True),
     ("        ", 12, NAVY, False),
     ("近似误差 (表征不够)", 12, ACCENT, True),
     ("        ", 12, NAVY, False),
     ("方差 (数据 / 优化)", 12, PURPLE, True)],
])

# Two interpretation panes
add_rect(s, Inches(0.5), Inches(3.45), Inches(6.0), Inches(3.55), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(3.45), Inches(6.0), Inches(0.5), RED)
add_text(s, Inches(0.5), Inches(3.45), Inches(6.0), Inches(0.5),
         "信息完备性  ↔  σ²(Y|X)", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(0.7), Inches(4.05), Inches(5.7), Inches(2.85), [
    [("• 如果上下文 X 不包含决定 Y 的所有信息,",
      12, DGREY, False)],
    [("   ", 12, NAVY, False),
     ("σ²(Y|X) > 0 — 这是『硬下界』", 12, NAVY, True)],
    [("", 6, NAVY, False)],
    [("• ", 12, RED, True),
     ("再强的模型也无法跨越这个下界", 12, RED, True)],
    [("", 6, NAVY, False)],
    [("• 反应器例子: 只观测 T(t), 不观测组分 c(t),",
      12, DGREY, False)],
    [("   σ²(T(t+h) | T_past) ≥ σ²(T | T_past, c_past)",
      11, GREY, False)],
    [("", 6, NAVY, False)],
    [("• 这是用户两轴中第一轴的精确含义",
      12, NAVY, True)],
])

add_rect(s, Inches(6.8), Inches(3.45), Inches(6.0), Inches(3.55), WHITE, line=SOFT)
add_rect(s, Inches(6.8), Inches(3.45), Inches(6.0), Inches(0.5), ACCENT)
add_text(s, Inches(6.8), Inches(3.45), Inches(6.0), Inches(0.5),
         "表征能力  ↔  ‖f* − f̂‖²", size=14, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(7.0), Inches(4.05), Inches(5.7), Inches(2.85), [
    [("• 模型族 F = {f_θ : θ ∈ Θ} 必须 ",
      12, DGREY, False),
     ("能逼近真函数 f*", 12, NAVY, True)],
    [("• 否则即使 X 完备, 你也 fit 不出来",
      12, DGREY, False)],
    [("", 6, NAVY, False)],
    [("• 经典案例: 线性回归 vs 强非线性动力学",
      12, DGREY, False)],
    [("   线性 model 永远抓不住 chaotic 系统",
      12, GREY, False)],
    [("", 6, NAVY, False)],
    [("• Universal Approximation: 足够宽 / 深的 NN ",
      12, DGREY, False)],
    [("   ⇒ 任意连续函数, 任意精度 (Cybenko 1989)",
      12, GREY, False)],
    [("• 这是用户两轴中第二轴的精确含义",
      12, NAVY, True)],
])

# ============================================================
# Slide 4 — But this framework is incomplete — overview
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 4, TOTAL, "但这个两轴框架并不完整",
     "信息完备 + 容量足够 ⟹ 好的预测? 这是一个常见的误解",
     section="PART II · 缺失的维度")

# big claim
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.4), DEEP)
add_text(s, Inches(0.85), Inches(1.7), Inches(11.8), Inches(0.5),
         "反例 (反直觉): 信息完备 + 容量足够 ≠ 一定能学到好的预测",
         size=15, bold=True, color=ACCENT)
rich(s, Inches(0.85), Inches(2.2), Inches(11.8), Inches(0.7), [
    [("假设你给一个 10 亿参数的 MLP 喂入完整的 Lorenz 系统观测 (信息完备 ✓ , 容量充足 ✓), ",
      13, WHITE, False)],
    [("它能否学会预测? 在合理算力 / 数据下 ", 13, WHITE, False),
     ("多半不能", 13, ACCENT, True),
     (" — 因为缺少『序列结构 / 局部性 / 多尺度』这些归纳偏置。",
      13, WHITE, False)],
])

# 4 missing axes
items = [
    ("Axis 3", "归纳偏置 / 先验",
     "No Free Lunch — 没有合适的先验, 即使有无限数据 + 无限容量, 学习效率也会塌陷。\n例: CNN 对图像、Transformer 对语言、Patching+RevIN 对时序。"),
    ("Axis 4", "训练数据规模 + 多样性",
     "和『per-instance 上下文』完全不同: 这是训练分布的覆盖度。\n泛化界: ε ≤ ε_train + O(√(C/N)) — N 是训练数据量, 不是单条样本上下文。"),
    ("Axis 5", "泛化机制 / 迁移机制",
     "怎么从 train 迁到 test? IID vs OOD vs zero-shot — 不同迁移机制有不同样本复杂度。\nICL (in-context learning) 是新近发现的高效泛化机制。"),
    ("Axis 6", "可识别性 (Identifiability)",
     "比『信息完备』更上游: 即使观测变量看似齐全, 系统的某些参数 / 状态可能在观测下不可分。\n例: 两个不同的隐状态产生相同的观测序列 — 任何模型都无能为力。"),
]
gx = Inches(0.5); gy = Inches(3.15); cw = Inches(6.0); ch = Inches(1.85)
for i, (tag, title, body) in enumerate(items):
    col = i % 2; row = i // 2
    x = gx + col * (cw + Inches(0.3))
    y = gy + row * (ch + Inches(0.15))
    add_rect(s, x, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, y, Inches(0.18), ch, ACCENT)
    add_text(s, x + Inches(0.3), y + Inches(0.1), Inches(1.5), Inches(0.35),
             tag, size=12, bold=True, color=ACCENT)
    add_text(s, x + Inches(1.7), y + Inches(0.1), cw - Inches(1.9), Inches(0.4),
             title, size=14, bold=True, color=NAVY)
    add_text(s, x + Inches(0.3), y + Inches(0.55), cw - Inches(0.5), Inches(1.25),
             body, size=11, color=DGREY)

# ============================================================
# Slide 5 — Axis 3: Inductive bias / prior
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 5, TOTAL, "缺失轴 3 · 归纳偏置 (Inductive Bias)",
     "No Free Lunch 定理: 没有合适的先验, 学习就是在沙子上写字",
     section="PART II · 缺失的维度")

# Top: NFL theorem
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.5), LIGHT)
add_rect(s, Inches(0.5), Inches(1.55), Inches(0.18), Inches(1.5), ACCENT)
add_text(s, Inches(0.85), Inches(1.65), Inches(11.8), Inches(0.4),
         "No Free Lunch 定理 (Wolpert 1996)", size=14, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(2.05), Inches(11.8), Inches(0.95), [
    [("对所有可能的目标函数取平均, 任何两个学习算法的期望性能 ", 12, DGREY, False),
     ("完全相同", 12, RED, True), (".", 12, DGREY, False)],
    [("→ ", 13, ACCENT, True),
     ("学习算法的『好』完全来自于它的归纳偏置和真实问题分布的匹配。",
      12, NAVY, True)],
    [("→ ", 13, ACCENT, True),
     ("没有先验 = 没有泛化。容量再大、数据再多, 也只是在记忆。",
      12, NAVY, True)],
])

# Bottom: examples in 3 columns
cols = [
    ("视觉 / CNN",
     "卷积 = 平移不变性 + 局部性",
     "ImageNet 之所以可学, 不是因为数据多,\n而是因为 CNN 的先验匹配了图像的统计结构。\n用 MLP 在 ImageNet 上同样数据量下\n准确率掉 30+ 点。",
     PURPLE),
    ("语言 / Transformer",
     "self-attention = 任意距离的关联",
     "语言的远程依赖 + 排列性 ↔ Transformer 的归纳偏置。\nRNN (顺序先验) 在长上下文上落败,\n不是因为容量, 而是先验不对。",
     TEAL),
    ("时序 / TSFM",
     "Patching + Channel-Indep + RevIN",
     "Patching = 多尺度 + BPE 类压缩\nChannel-Indep = 每条序列独立 (放大有效数据)\nRevIN = 仪器尺度无关\n这些先验是 TSFM 设计者刻意挑出来的。",
     ACCENT),
]
gx = Inches(0.5); gy = Inches(3.25); cw = Inches(4.05); ch = Inches(3.7)
for i, (title, sub, body, col) in enumerate(cols):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.65), col)
    add_text(s, x, gy, cw, Inches(0.65), title, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), gy + Inches(0.8), cw - Inches(0.4), Inches(0.5),
             sub, size=12, bold=True, color=NAVY, italic=True)
    add_text(s, x + Inches(0.2), gy + Inches(1.4), cw - Inches(0.4), Inches(2.2),
             body, size=11, color=DGREY)

add_text(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
         "→ 信息完备 + 容量充足 是 必要条件; 但归纳偏置才决定 你能否高效逼近真函数。",
         size=12, italic=True, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 6 — Axis 4: training data scale & diversity
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 6, TOTAL, "缺失轴 4 · 训练数据规模 + 多样性",
     "用户混淆了两个 X: per-instance 上下文 vs 训练分布的覆盖",
     section="PART II · 缺失的维度")

# Distinction
add_rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(2.4), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.55), DEEP)
add_text(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.55),
         "X_context  (用户的轴 1)", size=14, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(0.7), Inches(2.25), Inches(5.7), Inches(1.6), [
    [("• ", 12, ACCENT, True), ("单个样本", 12, NAVY, True),
     (" 可见的过去窗口 + 协变量。", 12, DGREY, False)],
    [("• 例: 给我看过去 96 个温度点 + 当前组分。",
      12, DGREY, False)],
    [("• 决定 ", 12, DGREY, False),
     ("不可约误差", 12, RED, True),
     (" 的下界 σ²(Y|X)。", 12, DGREY, False)],
    [("• 你需要 ", 12, DGREY, False),
     ("在推理时", 12, NAVY, True),
     (" 拥有这些信息。", 12, DGREY, False)],
])

add_rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(2.4), WHITE, line=SOFT)
add_rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.55), ACCENT)
add_text(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.55),
         "D_train  (被忽略的轴 4)", size=14, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(7.0), Inches(2.25), Inches(5.7), Inches(1.6), [
    [("• ", 12, ACCENT, True), ("整个训练分布", 12, NAVY, True),
     (" — 跨多少域 / 频率 / 模式。", 12, DGREY, False)],
    [("• 例: TSFM 在 LOTSA 27B 时序点上预训练, 跨 9 大域。",
      12, DGREY, False)],
    [("• 决定 ", 12, DGREY, False),
     ("泛化误差", 12, RED, True),
     (" 的范围: |R(f̂) − R̂(f̂)| ≤ O(√(C/N)).",
      12, DGREY, False)],
    [("• 你需要 ", 12, DGREY, False),
     ("在训练时", 12, NAVY, True),
     (" 见过足够多样的世界。", 12, DGREY, False)],
])

# Empirical scaling: lookup table
add_rect(s, Inches(0.5), Inches(4.15), Inches(12.3), Inches(2.4), LIGHT)
add_text(s, Inches(0.85), Inches(4.3), Inches(11.8), Inches(0.4),
         "训练数据规模的实证: TSFM 把这一轴向上推了 5-7 个数量级",
         size=14, bold=True, color=NAVY)

# Mini bar chart
labels = ["ETT", "M4", "Monash", "LOTSA", "GIFT-Pretrain"]
counts = ["7×10⁴", "1×10⁵", "5×10⁹", "2.7×10¹⁰", "2.3×10¹¹"]
colors_b = [GREY, GREY, TEAL, ACCENT, GREEN]
heights = [0.15, 0.18, 0.55, 0.85, 1.0]
bx = Inches(1.0); by = Inches(6.35); bw = Inches(1.8); bh_max = Inches(1.4)
for i, (lbl, cnt, c, h) in enumerate(zip(labels, counts, colors_b, heights)):
    x = bx + i * (bw + Inches(0.5))
    bh = int(bh_max * h)
    add_rect(s, x, by - bh, bw, bh, c)
    add_text(s, x, by - bh - Inches(0.4), bw, Inches(0.35),
             cnt, size=11, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_text(s, x, by + Inches(0.05), bw, Inches(0.3),
             lbl, size=11, color=DGREY, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3),
         "→ TSFM 之前: ~10⁵ 数据点 (per-dataset);  TSFM 之后: ~10¹⁰ 数据点 (cross-domain). 这是质变。",
         size=12, italic=True, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 7 — Axis 5: generalization mechanism
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 7, TOTAL, "缺失轴 5 · 泛化机制 (Transfer Mechanism)",
     "训练好的模型怎么从 train 用到 test? 不同机制的样本复杂度差几个量级",
     section="PART II · 缺失的维度")

# Four mechanisms
mechs = [
    ("Mechanism A", "IID 监督学习",
     "train / test 同分布",
     "需要 N 大;\n经典 PAC 界 O(VC/ε²);\n换数据集 = 重新训练。",
     "传统 LSTM / TFT", GREY),
    ("Mechanism B", "迁移学习 (transfer)",
     "预训练 + 下游微调",
     "ImageNet → 下游任务范式;\n下游需要数百-数千标签;\n仍需训练阶段。",
     "BERT-style fine-tune", TEAL),
    ("Mechanism C", "元学习 (meta-learning)",
     "学习『如何快速学习』",
     "MAML / Reptile;\n下游需要少量梯度步;\n仍需训练阶段。",
     "MAML, Reptile", PURPLE),
    ("Mechanism D", "in-context learning (ICL)",
     "推理时直接用上下文学习",
     "GPT-3 涌现;\n下游零梯度步;\n是机制层面的质变。",
     "TSFM zero-shot", ACCENT),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(3.05); ch = Inches(4.0)
for i, (tag, title, sub, body, ex, col) in enumerate(mechs):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.6), col)
    add_text(s, x, gy, cw, Inches(0.6), tag, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.15), gy + Inches(0.7), cw - Inches(0.3), Inches(0.5),
             title, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), gy + Inches(1.2), cw - Inches(0.3), Inches(0.4),
             sub, size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), gy + Inches(1.7), cw - Inches(0.3), Inches(1.6),
             body, size=11, color=DGREY)
    add_rect(s, x + Inches(0.15), gy + Inches(3.4), cw - Inches(0.3), Inches(0.45), LIGHT)
    add_text(s, x + Inches(0.15), gy + Inches(3.4), cw - Inches(0.3), Inches(0.45),
             ex, size=10, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Bottom: punchline
add_rect(s, Inches(0.5), Inches(5.75), Inches(12.3), Inches(1.2), DEEP)
rich(s, Inches(0.85), Inches(5.85), Inches(11.8), Inches(1.0), [
    [("→ ", 14, ACCENT, True),
     ("机制 A → D 是『样本效率』的指数提升:", 13, WHITE, True)],
    [("    A: 10⁴-10⁵ 标签;  B: 10²-10³ 标签;  C: 10¹-10² 标签;  D: ", 12, SOFT, False),
     ("0 标签 (zero-shot) 或 0 训练步", 12, ACCENT, True)],
    [("→ ", 14, ACCENT, True),
     ("用户的两轴框架默认是机制 A。但 ICL 是一个全新的轴 — 你不能用 A 的直觉评估 D。",
      13, WHITE, True)],
])

# ============================================================
# Slide 8 — Axis 6: identifiability
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 8, TOTAL, "缺失轴 6 · 可识别性 (Identifiability)",
     "比信息完备性更上游 — 你的观测语义本身能否区分不同的因果机制?",
     section="PART II · 缺失的维度")

# Concrete example
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(2.4), LIGHT)
add_text(s, Inches(0.85), Inches(1.7), Inches(11.8), Inches(0.4),
         "形式定义", size=14, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(2.1), Inches(11.8), Inches(1.8), [
    [("如果存在两组不同的参数 / 隐状态 (θ₁, h₁) ≠ (θ₂, h₂), 但它们诱导出 ", 12, DGREY, False),
     ("完全相同的观测分布", 12, ACCENT, True),
     (" P(Y | θ₁, h₁) = P(Y | θ₂, h₂),", 12, DGREY, False)],
    [("→ ", 13, ACCENT, True),
     ("那么从观测中 不可能 区分这两种情况, 不论模型多强、数据多多。",
      12, NAVY, True)],
    [("", 8, NAVY, False)],
    [("• 这是 σ²(Y|X) 之外的另一种 ", 12, DGREY, False),
     ("结构性下界", 12, RED, True),
     (": 不是『噪声』, 而是『歧义』。", 12, DGREY, False)],
    [("• 即使你给 TSFM 海量数据 + 完美架构, 在不可识别区域它依然 50/50。",
      12, DGREY, False)],
])

# Two examples
add_rect(s, Inches(0.5), Inches(4.15), Inches(6.0), Inches(2.85), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(4.15), Inches(0.18), Inches(2.85), RED)
add_text(s, Inches(0.85), Inches(4.25), Inches(5.6), Inches(0.4),
         "案例 1 · 隐马尔可夫的标签置换",
         size=13, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(4.65), Inches(5.6), Inches(2.3), [
    [("• HMM 的两个隐状态 {A, B} 互换名称 ⇒ ", 11, DGREY, False),
     ("观测概率不变", 11, RED, True), (".", 11, DGREY, False)],
    [("• 不论数据多少, 你 ", 11, DGREY, False),
     ("无法区分", 11, RED, True),
     (" 哪个是 A 哪个是 B (label-switching).", 11, DGREY, False)],
    [("• 这不是信息缺失, 是参数本身的对称性。",
      11, DGREY, False)],
    [("• TSFM 也无能为力 — 这是观测语义的限制。",
      11, NAVY, True)],
])

add_rect(s, Inches(6.8), Inches(4.15), Inches(6.0), Inches(2.85), WHITE, line=SOFT)
add_rect(s, Inches(6.8), Inches(4.15), Inches(0.18), Inches(2.85), RED)
add_text(s, Inches(7.15), Inches(4.25), Inches(5.6), Inches(0.4),
         "案例 2 · 反应器场景 (前 PPT)",
         size=13, bold=True, color=NAVY)
rich(s, Inches(7.15), Inches(4.65), Inches(5.6), Inches(2.3), [
    [("• 反应器只测 T(t), 不测组分 c(t).", 11, DGREY, False)],
    [("• 多种 (c(t), T) 组合可能产生相同的 T 历史 ⇒ ",
      11, DGREY, False)],
    [("  ", 11, DGREY, False),
     ("c(t) 在 T 的观测下不可识别。", 11, RED, True)],
    [("• 即使 TSFM 再强大, 它无法回答 ", 11, DGREY, False),
     ("『下一秒组分会突变到多少』", 11, RED, True),
     (" — 因为这个信息在观测序列中 物理上 不存在。", 11, DGREY, False)],
    [("→ ", 11, ACCENT, True),
     ("解决方法: 加传感器, 不是加模型。",
      11, NAVY, True)],
])

# ============================================================
# Slide 9 — Axis 7 + summary of missing
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 9, TOTAL, "次要维度 · 优化可行性 + 校准",
     "工程上同样关键, 但相对来说是『可解决』的 — 不是宏观瓶颈",
     section="PART II · 缺失的维度")

# two columns: optimization + calibration
add_rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.45), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.6), DEEP)
add_text(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.6),
         "Axis 7 · 优化 / 可训练性", size=14, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(0.7), Inches(2.3), Inches(5.7), Inches(4.65), [
    [("• ", 13, ACCENT, True), ("即使最优 θ* 存在, SGD 能否找到?",
      13, NAVY, True)],
    [("", 5, NAVY, False)],
    [("非凸损失景观, 局部极小, 鞍点。",
      11, DGREY, False)],
    [("LayerNorm / RMSNorm / RevIN 不仅是先验, 也是优化稳定剂。",
      11, DGREY, False)],
    [("学习率调度 / warmup / weight decay。",
      11, DGREY, False)],
    [("", 8, NAVY, False)],
    [("• ", 13, ACCENT, True), ("为什么这是『次要』维度?",
      13, NAVY, True)],
    [("过去十年沉淀: Adam, AdamW, Lion, μP, ZeRO, FlashAttention.",
      11, DGREY, False)],
    [("对 TSFM 来说: 已经有成熟的工程栈, 不再是创新瓶颈。",
      11, DGREY, False)],
    [("→ ", 12, ACCENT, True),
     ("但你不能假装这一轴不存在 — 它决定『纸上能跑』和『实际能学』的差距。",
      12, NAVY, True)],
])

add_rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(5.45), WHITE, line=SOFT)
add_rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.6), DEEP)
add_text(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.6),
         "Axis 8 · 不确定性 / 校准", size=14, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(7.0), Inches(2.3), Inches(5.7), Inches(4.65), [
    [("• ", 13, ACCENT, True),
     ("点预测 ŷ 不够 — 工业场景需要 ",
      13, NAVY, True),
     ("置信区间", 13, ACCENT, True), (".", 13, NAVY, True)],
    [("", 5, NAVY, False)],
    [("Quantile loss → P10 / P50 / P90.", 11, DGREY, False)],
    [("Conformal prediction → 提供『有限样本可证』的覆盖率保证。",
      11, DGREY, False)],
    [("Bayesian posterior approximation.", 11, DGREY, False)],
    [("", 8, NAVY, False)],
    [("• ", 13, ACCENT, True), ("和用户两轴的关系?",
      13, NAVY, True)],
    [("用户的『信息完备』决定 σ²(Y|X) 大小, 但不决定你能否 ",
      11, DGREY, False),
     ("正确地报告", 11, RED, True),
     (" 这个不确定性。", 11, DGREY, False)],
    [("校准 ↔ 模型对自己的『不确定性』有自知之明。",
      11, DGREY, False)],
    [("→ ", 12, ACCENT, True),
     ("现代 TSFM 默认输出概率分布 (Chronos token sampling, Moirai mixture).",
      12, NAVY, True)],
])

# ============================================================
# Slide 10 — The unified 6-axis framework
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 10, TOTAL, "完整的宏观维度 · 一张图",
     "时序预测 = 6 轴 + 2 个次要轴的协同博弈",
     section="PART III · 完整框架")

# Hexagon-style layout: center + 6 surrounding
cx = Inches(6.65); cy = Inches(4.4)
R = Inches(2.3)
axes = [
    (0,    -1, "Axis 1\n信息完备性", "X 是否包含 Y 的所有信息", ACCENT, "用户提的"),
    (0.95, -0.3, "Axis 2\n表征能力", "f_θ 能否逼近 f*", ACCENT, "用户提的"),
    (0.6,   0.85, "Axis 3\n归纳偏置", "学习算法的『先验形状』", GREEN, "TSFM 创新"),
    (-0.6,  0.85, "Axis 4\n训练数据规模", "D_train 的覆盖度", GREEN, "TSFM 创新"),
    (-0.95, -0.3, "Axis 5\n泛化机制", "train→test 怎么传递", GREEN, "TSFM 创新"),
    (0,     0.0, "(中心)\n好的预测", "", NAVY, ""),  # placeholder for center
]
# We'll place 6 around (skip center placeholder differently). Use 6 directions:
dirs = [
    (0.0,  -1.0, "Axis 1", "信息完备性", ACCENT, "[用户]"),
    (0.866, -0.5, "Axis 2", "表征能力", ACCENT, "[用户]"),
    (0.866,  0.5, "Axis 3", "归纳偏置", GREEN, "[TSFM 发力]"),
    (0.0,   1.0, "Axis 4", "训练数据规模", GREEN, "[TSFM 发力]"),
    (-0.866, 0.5, "Axis 5", "泛化机制", GREEN, "[TSFM 发力]"),
    (-0.866, -0.5, "Axis 6", "可识别性", RED, "[TSFM 不解决]"),
]

# center
add_oval(s, cx - Inches(0.9), cy - Inches(0.7), Inches(1.8), Inches(1.4), NAVY)
add_text(s, cx - Inches(0.9), cy - Inches(0.7), Inches(1.8), Inches(1.4),
         "好的\n预测", size=16, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

bw = Inches(2.6); bh = Inches(1.1)
for dx, dy, tag, name, col, badge in dirs:
    x = cx + int(R * dx) - bw // 2
    y = cy + int(R * dy) - bh // 2
    add_rect(s, x, y, bw, bh, WHITE, line=col, line_w=Pt(2))
    add_rect(s, x, y, bw, Inches(0.32), col)
    add_text(s, x, y, bw, Inches(0.32), tag, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, y + Inches(0.4), bw, Inches(0.4),
             name, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.78), bw, Inches(0.3),
             badge, size=10, italic=True, color=col, align=PP_ALIGN.CENTER)
    # connecting line to center
    line(s, cx + int(R * dx * 0.3), cy + int(R * dy * 0.3),
         cx + int(R * dx * 0.7) - int(bw * dx * 0.15), cy + int(R * dy * 0.7) - int(bh * dy * 0.15),
         SOFT, w=1.2)

# Legend on left
add_rect(s, Inches(0.5), Inches(1.55), Inches(2.6), Inches(2.6), LIGHT)
add_text(s, Inches(0.65), Inches(1.65), Inches(2.4), Inches(0.35),
         "颜色含义", size=12, bold=True, color=NAVY)
rich(s, Inches(0.65), Inches(2.05), Inches(2.4), Inches(2.0), [
    [("■ ", 14, ACCENT, True), ("用户提的轴", 11, DGREY, True)],
    [("  必要, 但非充分", 10, GREY, False)],
    [("", 6, NAVY, False)],
    [("■ ", 14, GREEN, True), ("TSFM 发力点", 11, DGREY, True)],
    [("  这就是 TSFM 的真正创新", 10, GREY, False)],
    [("", 6, NAVY, False)],
    [("■ ", 14, RED, True), ("TSFM 无能为力", 11, DGREY, True)],
    [("  上游 / 观测端的问题", 10, GREY, False)],
])

# Right side: macro identity
add_rect(s, Inches(10.0), Inches(1.55), Inches(2.85), Inches(5.45), DEEP)
add_text(s, Inches(10.15), Inches(1.7), Inches(2.6), Inches(0.4),
         "宏观恒等式", size=12, bold=True, color=ACCENT)
rich(s, Inches(10.15), Inches(2.15), Inches(2.6), Inches(4.7), [
    [("好的预测  =  ", 11, ACCENT, True)],
    [("  f(轴1, ..., 轴6)", 11, WHITE, True)],
    [("", 6, WHITE, False)],
    [("不是 ", 11, SOFT, False),
     ("max", 11, ACCENT, True),
     (" 而是 ", 11, SOFT, False),
     ("min", 11, RED, True), (".", 11, SOFT, False)],
    [("", 4, WHITE, False)],
    [("即: 最弱的那一轴 决定整体表现 (木桶效应).",
      10, SOFT, False)],
    [("", 8, WHITE, False)],
    [("→ TSFM 强化了 4 个轴, 但如果你的瓶颈在轴 1 或轴 6, ",
      10, SOFT, False)],
    [("  TSFM ", 11, SOFT, False),
     ("救不了你", 11, ACCENT, True),
     (".", 11, SOFT, False)],
])

# ============================================================
# Slide 11 — Axes interact
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 11, TOTAL, "维度间不是独立的 · 它们相互制约和放大",
     "理解 TSFM 必须理解轴间耦合, 不能逐一评估",
     section="PART III · 完整框架")

# Five interactions
inters = [
    ("容量 × 数据规模",
     "Chinchilla 缩放律: 给定算力 C, 最优 N ∝ C^0.5, D ∝ C^0.5.",
     "TSFM 把数据从 10⁵ 推到 10¹⁰ ⇒ 解锁了 100M-1B 参数模型的潜力。",
     ACCENT),
    ("归纳偏置 × 数据规模",
     "更好的先验 ⇒ 同样精度需要更少数据。",
     "Channel-Indep 让有效数据放大 c 倍 (c = 通道数);\nPatching 让序列长度缩短 16×.",
     GREEN),
    ("信息完备 × 表征能力",
     "若 X 不完备, 容量再大也撞 σ²(Y|X) 下界 ⇒ 投资容量的边际回报为 0.",
     "反应器场景: 单变量下 200M / 1B / 10B 模型 MAE 几乎相同。",
     RED),
    ("泛化机制 × 数据规模",
     "ICL 是『数据规模触发的涌现现象』 (Olsson 2022, Chan 2022).",
     "在小数据上 ICL 不出现; 必须有足够数据多样性 ⇒ 这是新轴的 prerequisites.",
     PURPLE),
    ("可识别性 × 一切",
     "若不可识别, 任何一轴上的进步都白搭.",
     "这是『硬墙』: 必须从观测端 (传感器 / 实验设计) 解决.",
     RED),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(12.3); ch = Inches(1.05)
for i, (title, body1, body2, col) in enumerate(inters):
    y = gy + i * (ch + Inches(0.05))
    add_rect(s, gx, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, gx, y, Inches(0.2), ch, col)
    add_text(s, gx + Inches(0.35), y + Inches(0.08), Inches(3.0), Inches(0.4),
             title, size=13, bold=True, color=col)
    add_text(s, gx + Inches(0.35), y + Inches(0.5), Inches(11.7), Inches(0.3),
             body1, size=11, color=DGREY)
    add_text(s, gx + Inches(0.35), y + Inches(0.78), Inches(11.7), Inches(0.3),
             body2, size=11, italic=True, color=GREY)

# ============================================================
# Slide 12 — Where TSFM hits — overview
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 12, TOTAL, "TSFM 究竟在哪里发力?",
     "答案: 协同性地推动了 4 个维度, 而非单一维度的胜利",
     section="PART IV · TSFM 的发力点")

# Big visual: 6 axes with TSFM contribution
axes = [
    ("Axis 1 · 信息完备性",     0,    "—",      "—",                       GREY),
    ("Axis 2 · 表征能力",        85,   "↑↑",    "scaled transformer 100M-1B", ACCENT),
    ("Axis 3 · 归纳偏置",        95,   "↑↑↑",   "Patching + RevIN + Channel-Indep + Decoder-only", GREEN),
    ("Axis 4 · 训练数据规模",    100,  "↑↑↑",   "LOTSA / GIFT-Pretrain — 5-7 量级跃升",        GREEN),
    ("Axis 5 · 泛化机制",        100,  "↑↑↑",   "ICL / induction heads — 机制层面质变",        GREEN),
    ("Axis 6 · 可识别性",        0,    "—",     "TSFM 完全不参与 (观测端问题)",                 RED),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(12.3); rh = Inches(0.88)
add_rect(s, gx, gy, cw, Inches(0.45), DEEP)
add_text(s, gx + Inches(0.2), gy, Inches(2.5), Inches(0.45),
         "维度", size=12, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, gx + Inches(3.5), gy, Inches(2.0), Inches(0.45),
         "TSFM 贡献度", size=12, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, gx + Inches(5.7), gy, Inches(1.5), Inches(0.45),
         "强度", size=12, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, gx + Inches(7.4), gy, Inches(4.7), Inches(0.45),
         "代表创新", size=12, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

for i, (name, score, sym, body, col) in enumerate(axes):
    y = gy + Inches(0.5) + i * rh
    add_rect(s, gx, y, cw, rh - Inches(0.05), WHITE, line=SOFT)
    add_text(s, gx + Inches(0.2), y + Inches(0.15), Inches(3.2), Inches(0.5),
             name, size=13, bold=True, color=NAVY)
    # bar
    bar_x = gx + Inches(3.5)
    bar_full = Inches(2.0)
    add_rect(s, bar_x, y + Inches(0.3), bar_full, Inches(0.25), LIGHT)
    if score > 0:
        add_rect(s, bar_x, y + Inches(0.3), int(bar_full * score / 100), Inches(0.25), col)
    add_text(s, gx + Inches(5.7), y + Inches(0.15), Inches(1.5), Inches(0.5),
             sym, size=14, bold=True, color=col,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, gx + Inches(7.4), y + Inches(0.15), Inches(4.7), Inches(0.5),
             body, size=11, color=DGREY, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.5), LIGHT)
add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.5),
         "→  TSFM 是 4 个轴的同时进步; 这种『协同推进』是它和『更大的 LSTM』的本质区别。",
         size=13, italic=True, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 13 — Axis 4: data scale (deep dive)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 13, TOTAL, "TSFM 维度 1 · 训练数据规模 + 多样性",
     "把 TS 训练分布从 10⁵ 推到 10¹⁰ — 这是质变, 不是量变",
     section="PART IV · TSFM 的发力点")

# Three subpanels
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(2.3), LIGHT)
add_text(s, Inches(0.85), Inches(1.7), Inches(11.8), Inches(0.4),
         "为什么数据规模是『新维度』, 不只是更多数据",
         size=14, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(2.15), Inches(11.8), Inches(1.7), [
    [("• ", 13, ACCENT, True),
     ("跨域多样性 (cross-domain diversity)", 13, NAVY, True),
     (": LOTSA 包含能源、交通、金融、气象、医疗、销售、Web 等 9 大域;",
      12, DGREY, False)],
    [("  M4 / Monash 加起来还在单一基准内. ", 12, DGREY, False),
     ("→ 跨域是『让模型学到时序的共同语义』的关键.",
      12, NAVY, True)],
    [("• ", 13, ACCENT, True),
     ("跨频率覆盖", 13, NAVY, True),
     (": 从秒级 (网络流量) 到年级 (经济周期), 频率跨 8 个数量级 ⇒ ",
      12, DGREY, False)],
    [("  模型见过的『时间尺度词汇表』 与传统模型差几个量级.",
      12, DGREY, False)],
    [("• ", 13, ACCENT, True),
     ("缩放律涌现", 13, NAVY, True),
     (": Loss(N, D) ≈ A·N^-α + B·D^-β  (α≈0.05, β≈0.10 for TSFM, GIFT-Eval 实证)",
      12, DGREY, False)],
])

# Quantitative comparison
add_rect(s, Inches(0.5), Inches(4.05), Inches(6.0), Inches(2.95), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(4.05), Inches(6.0), Inches(0.55), GREY)
add_text(s, Inches(0.5), Inches(4.05), Inches(6.0), Inches(0.55),
         "传统范式 (per-dataset)", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(0.7), Inches(4.7), Inches(5.7), Inches(2.25), [
    [("• 训练数据: ", 12, DGREY, False),
     ("10⁴ - 10⁵ 时间点", 12, NAVY, True),
     (" (单一数据集)", 12, DGREY, False)],
    [("• 跨域: ", 12, DGREY, False), ("不可行", 12, RED, True)],
    [("• 训练成本: ", 12, DGREY, False), ("每个新场景重新训练", 12, RED, True)],
    [("• 数据效率取决于: ", 12, DGREY, False),
     ("当前数据集的样本量", 12, DGREY, False)],
    [("• 极限性能 (best-in-class):", 12, DGREY, False)],
    [("   TFT / N-BEATS 经过精调 → SOTA-on-domain",
      11, GREY, False)],
])

add_rect(s, Inches(6.8), Inches(4.05), Inches(6.0), Inches(2.95), WHITE, line=SOFT)
add_rect(s, Inches(6.8), Inches(4.05), Inches(6.0), Inches(0.55), GREEN)
add_text(s, Inches(6.8), Inches(4.05), Inches(6.0), Inches(0.55),
         "TSFM 范式 (cross-domain)", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(7.0), Inches(4.7), Inches(5.7), Inches(2.25), [
    [("• 训练数据: ", 12, DGREY, False),
     ("10⁹ - 10¹¹ 时间点", 12, NAVY, True),
     (" (跨 9+ 域)", 12, DGREY, False)],
    [("• 跨域: ", 12, DGREY, False), ("是设计目标", 12, GREEN, True)],
    [("• 训练成本: ", 12, DGREY, False), ("一次预训练, 永久复用", 12, GREEN, True)],
    [("• 数据效率: ", 12, DGREY, False),
     ("零样本 + 少样本 (跨域积累)", 12, GREEN, True)],
    [("• 极限性能:", 12, DGREY, False)],
    [("   零样本对齐或胜过领域 SOTA",
      11, GREY, False)],
])

# ============================================================
# Slide 14 — Axis 3: inductive bias deep dive
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 14, TOTAL, "TSFM 维度 2 · 归纳偏置的精心选择",
     "TSFM 不是『更大的 transformer』, 而是装上了 4 个时序专属先验",
     section="PART IV · TSFM 的发力点")

priors = [
    ("Patching",
     "把 16-64 个时间点合并为 1 个 token",
     "类比 BPE: 把『字符』升级为『词』.\n• 序列长度 ÷16  ⇒ 长 context 下注意力可承受.\n• 多尺度先验: 邻近时间点必然相关.\n• 减少高频噪声主导.",
     ACCENT),
    ("Channel-Independence",
     "每条序列独立处理, 不混合通道",
     "反直觉但实证证明强:\n• 多变量数据被『展开』为多倍单变量样本 ⇒ 有效数据放大.\n• 避免不同变量噪声相互污染.\n• 与跨域预训练完美契合.",
     GREEN),
    ("Reversible Norm (RevIN / InstanceNorm)",
     "推理前减均值除标准差, 推理后还原",
     "尺度不变先验:\n• 不论温度是 K 还是 ℃, 模型看到的都是 z-score.\n• 解决跨域『单位 / 量级不同』的根本问题.\n• 是 zero-shot 能 work 的前提之一.",
     PURPLE),
    ("Decoder-only / Causal AR",
     "只看过去, 自回归生成",
     "时序的因果先验:\n• 不能用未来信息预测当前 ⇒ 强约束.\n• 与 LLM 同样的训练范式 ⇒ 利用 LLM 工具栈.\n• 天然支持任意长度的 horizon.",
     TEAL),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(6.05); ch = Inches(2.7)
for i, (title, sub, body, col) in enumerate(priors):
    col_idx = i % 2; row_idx = i // 2
    x = gx + col_idx * (cw + Inches(0.2))
    y = gy + row_idx * (ch + Inches(0.15))
    add_rect(s, x, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, y, cw, Inches(0.55), col)
    add_text(s, x + Inches(0.2), y + Inches(0.05), cw - Inches(0.4), Inches(0.45),
             title, size=14, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), y + Inches(0.65), cw - Inches(0.4), Inches(0.4),
             sub, size=11, italic=True, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(1.1), cw - Inches(0.4), Inches(1.5),
             body, size=11, color=DGREY)

# ============================================================
# Slide 15 — Axis 5: generalization mechanism (ICL)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 15, TOTAL, "TSFM 维度 3 · 泛化机制 — ICL / Induction Heads",
     "这是真正的『新轴』 — 而非旧轴上更进一步",
     section="PART IV · TSFM 的发力点")

# Top: definition
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.6), DEEP)
rich(s, Inches(0.85), Inches(1.7), Inches(11.8), Inches(1.4), [
    [("In-Context Learning (ICL):", 14, ACCENT, True)],
    [("把 ", 12, WHITE, False),
     ("『新任务的少量示例』", 12, ACCENT, True),
     (" 直接放在 prompt 里, ", 12, WHITE, False),
     ("不微调", 12, ACCENT, True),
     (" 直接得出该任务上的预测.", 12, WHITE, False)],
    [("→ 对于 TSFM, 输入序列的过去 = ", 13, SOFT, False),
     ("『示例』", 13, ACCENT, True),
     (";  模型在推理时『看出』这是什么样的序列, 然后预测后续.",
      13, SOFT, False)],
    [("数学上: ICL 等价于 ", 12, WHITE, False),
     ("隐式贝叶斯推断", 12, ACCENT, True),
     (" P(y | x, context) = ∫ P(y | x, θ) · P(θ | context) dθ.",
      12, WHITE, False)],
])

# Bottom: 3 columns
mech = [
    ("Induction Heads",
     "Olsson et al. 2022",
     "Transformer 在第二层 (≥3 层) 涌现的注意力模式:\n搜索『过去出现过的相似模式』, 然后 复制 它的后续.\n→ 这是 ICL 的电路级机制.",
     PURPLE),
    ("Implicit Bayesian Inference",
     "Xie et al. 2022",
     "GPT-3 在 prompt 上的行为可以严格解释为:\n模型对 P(θ | examples) 的近似贝叶斯更新.\n→ 解释了『更多 example ⇒ 更准』的规律.",
     TEAL),
    ("Phase Transition",
     "Chan et al. 2022",
     "ICL 在数据规模 / 模型规模到达某阈值时 涌现.\n小数据 / 小模型不出现.\n→ 这就是为什么 TSFM 必须先突破数据规模.",
     ACCENT),
]
gx = Inches(0.5); gy = Inches(3.3); cw = Inches(4.05); ch = Inches(3.6)
for i, (title, ref, body, col) in enumerate(mech):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.55), col)
    add_text(s, x, gy, cw, Inches(0.55), title, size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), gy + Inches(0.65), cw - Inches(0.4), Inches(0.4),
             ref, size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), gy + Inches(1.1), cw - Inches(0.4), Inches(2.4),
             body, size=11, color=DGREY)

add_text(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
         "→ ICL 不是用户两轴中任何一轴的延伸 — 它是机制层面的 emergence, 由数据规模 + 容量 + 偏置 共同触发.",
         size=11, italic=True, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 16 — Axis 2: representation capacity (directed)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 16, TOTAL, "TSFM 维度 4 · 表征能力 — 但是『定向』提升",
     "不是单纯加大模型, 而是把容量花在正确的位置",
     section="PART IV · TSFM 的发力点")

# Two columns: naive vs TSFM
add_rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(5.45), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.55), GREY)
add_text(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.55),
         "天真做法 · 单纯放大 LSTM/MLP", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(0.7), Inches(2.25), Inches(5.7), Inches(4.6), [
    [("• 假设: 容量越大 ⇒ 性能越好",
      12, NAVY, True)],
    [("• 实证: 在 TS 上, 增大 LSTM 收益快速饱和 (Salinas 2020).",
      12, DGREY, False)],
    [("", 6, NAVY, False)],
    [("瓶颈不在容量, 而在:", 12, NAVY, True)],
    [("  • 长上下文下注意力 / 记忆受限",
      11, GREY, False)],
    [("  • 单数据集训练 ⇒ 容量被噪声占满",
      11, GREY, False)],
    [("  • 没有『多任务可迁移结构』可学",
      11, GREY, False)],
    [("", 6, NAVY, False)],
    [("→ 结果:", 12, NAVY, True)],
    [("  100M LSTM 在 ETT 上 比 10M 强不了多少,",
      11, DGREY, False)],
    [("  但在 LOTSA 上, 100M Timer-XL 远胜 10M.",
      11, DGREY, False)],
    [("  容量必须配合其他三轴才能兑现.",
      11, ACCENT, True)],
])

add_rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(5.45), WHITE, line=SOFT)
add_rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.55), GREEN)
add_text(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.55),
         "TSFM 做法 · 容量定向投入", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(7.0), Inches(2.25), Inches(5.7), Inches(4.6), [
    [("容量分配的工程巧思:", 12, NAVY, True)],
    [("", 6, NAVY, False)],
    [("• ", 12, ACCENT, True),
     ("Patching 减序列长度", 12, NAVY, True),
     (" ⇒ attention 复杂度从 O(L²) 降到 O((L/p)²)",
      11, DGREY, False)],
    [("  ⇒ 同算力下能放大模型.",
      11, GREY, False)],
    [("• ", 12, ACCENT, True),
     ("Channel-Indep 让 batch 放大", 12, NAVY, True),
     (" ⇒ 每个 GPU 可同时处理几十个序列",
      11, DGREY, False)],
    [("  ⇒ 训练吞吐 ↑.",
      11, GREY, False)],
    [("• ", 12, ACCENT, True),
     ("RoPE 解锁长 context", 12, NAVY, True),
     (" ⇒ 不需要重新训练就能延长输入历史.",
      11, DGREY, False)],
    [("• ", 12, ACCENT, True),
     ("Mixture of frequencies", 12, NAVY, True),
     (" ⇒ 一个模型同时学秒级和年级.",
      11, DGREY, False)],
    [("", 6, NAVY, False)],
    [("→ 容量从 10M 推到 1B 的可行路径,",
      12, GREEN, True)],
    [("   而不是『硬加层数』.", 12, GREEN, True)],
])

# ============================================================
# Slide 17 — Bonus: scaling laws discovered
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 17, TOTAL, "副产品 · 时序的缩放律被发现",
     "TSFM 的另一个深远贡献: 让 TS 进入了『可预测改进』时代",
     section="PART IV · TSFM 的发力点")

# Top: equation
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.6), LIGHT)
add_text(s, Inches(0.85), Inches(1.7), Inches(11.8), Inches(0.4),
         "Time-Series Scaling Laws (Edwards 2024, GIFT-Eval, Chronos team)",
         size=14, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(2.15), Inches(11.8), Inches(0.95), [
    [("        Loss(N, D)  ≈  A · N^(-α)  +  B · D^(-β)  +  L_∞",
      18, NAVY, True)],
    [("        ", 12, NAVY, False),
     ("α ≈ 0.05 ", 14, ACCENT, True),
     ("(参数缩放指数);  ", 12, NAVY, True),
     ("β ≈ 0.10 ", 14, ACCENT, True),
     ("(数据缩放指数);  ", 12, NAVY, True),
     ("L_∞", 14, RED, True),
     (" = 不可约误差", 12, NAVY, True)],
])

# Bottom: implications
imps = [
    ("可预测",
     "知道扩大模型 / 数据 ⇒ 性能改善的精确数字",
     "在工业上意味着可以 ROI 评估: 多花 X 美元算力, 误差能降到 Y%.",
     ACCENT),
    ("有明确目标",
     "α / β 越小 ⇒ 收益越快边际递减 ⇒ 知道何时停手",
     "TSFM 的 α/β 比 LLM 大 ⇒ 还有较大空间; 但 β 远小于 α ⇒ 数据是更稀缺的资源.",
     GREEN),
    ("有不可越过的下界",
     "L_∞ 体现了信息完备性 (轴 1) 的最终墙",
     "对应用户的轴 1: 即使 N → ∞, D → ∞, 也只能逼近 L_∞.\n打破 L_∞ ⇒ 加传感器 / 加协变量, 不是加模型.",
     RED),
]
gx = Inches(0.5); gy = Inches(3.3); cw = Inches(4.05); ch = Inches(3.55)
for i, (title, sub, body, col) in enumerate(imps):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.55), col)
    add_text(s, x, gy, cw, Inches(0.55), title, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), gy + Inches(0.7), cw - Inches(0.4), Inches(0.5),
             sub, size=12, italic=True, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), gy + Inches(1.3), cw - Inches(0.4), Inches(2.2),
             body, size=11, color=DGREY)

# ============================================================
# Slide 18 — TSFM does NOT solve information completeness
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 18, TOTAL, "TSFM 的边界 (1) · 信息完备性问题",
     "用户的轴 1 是 TSFM 完全不解决的问题 — 对应上一份 PPT 的反应器场景",
     section="PART V · TSFM 的边界")

add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.5), DEEP)
rich(s, Inches(0.85), Inches(1.65), Inches(11.8), Inches(1.4), [
    [("一句话:", 14, ACCENT, True)],
    [("TSFM 让你『更高效地用 X 预测 Y』, ", 13, SOFT, False),
     ("但它不能让 X 包含原本不在 X 里的信息", 13, ACCENT, True),
     (".", 13, SOFT, False)],
    [("", 6, SOFT, False)],
    [("→ 如果你的瓶颈是 σ²(Y|X) (轴 1), TSFM ", 13, SOFT, False),
     ("零贡献", 13, RED, True), (".", 13, SOFT, False)],
])

# Two scenarios
add_rect(s, Inches(0.5), Inches(3.2), Inches(6.0), Inches(3.4), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(3.2), Inches(0.18), Inches(3.4), RED)
add_text(s, Inches(0.85), Inches(3.3), Inches(5.6), Inches(0.4),
         "场景 A · 反应器组分突变 (前 PPT)", size=13, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(3.7), Inches(5.6), Inches(2.85), [
    [("• 你只观测 T(t).", 12, DGREY, False)],
    [("• 真实系统: dT/dt = f(T, c(t), F, ...).",
      12, DGREY, False)],
    [("• 组分 c(t) 突变 ⇒ T 的轨迹突变.",
      12, DGREY, False)],
    [("• ", 12, DGREY, False),
     ("Stark 1999 定理: c 不能从 T 序列恢复",
      12, RED, True), (".", 12, DGREY, False)],
    [("• 不论 TSFM 多大, σ²(T(t+h) | T_past) > 0",
      12, DGREY, False)],
    [("• → 解决: 加组分传感器 (硬件), 不是换模型.",
      12, ACCENT, True)],
])

add_rect(s, Inches(6.8), Inches(3.2), Inches(6.0), Inches(3.4), WHITE, line=SOFT)
add_rect(s, Inches(6.8), Inches(3.2), Inches(0.18), Inches(3.4), RED)
add_text(s, Inches(7.15), Inches(3.3), Inches(5.6), Inches(0.4),
         "场景 B · 黑天鹅 / 罕见事件",
         size=13, bold=True, color=NAVY)
rich(s, Inches(7.15), Inches(3.7), Inches(5.6), Inches(2.85), [
    [("• 训练数据中没出现过类似事件",
      12, DGREY, False),
     (" (例: 一个新冠级别的需求冲击).",
      12, DGREY, False)],
    [("• TSFM 的 ICL 也无法 ", 12, DGREY, False),
     ("从无中生有", 12, RED, True), (".", 12, DGREY, False)],
    [("• 实证: TSFM 在『分布外极端事件』上误差激增 (Goerg 2024).",
      12, DGREY, False)],
    [("• 这本质上也是信息完备性问题:",
      12, DGREY, False)],
    [("  上下文 = 历史观测;  历史从未见过的事件 = 信息缺口.",
      11, GREY, False)],
    [("• → 解决: 引入外生信号 / 异常机制, 不是更大模型.",
      12, ACCENT, True)],
])

add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
         "→ 用户的轴 1 是『硬墙』; TSFM 只能让你贴墙更近, 不能穿墙.",
         size=12, italic=True, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 19 — TSFM does NOT solve identifiability
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 19, TOTAL, "TSFM 的边界 (2) · 可识别性问题",
     "比信息完备性更上游 — TSFM 同样无能为力, 因为问题在观测语义",
     section="PART V · TSFM 的边界")

# Top: definition refresher
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.4), LIGHT)
rich(s, Inches(0.85), Inches(1.7), Inches(11.8), Inches(1.2), [
    [("不可识别 = 多个本质不同的真实机制, 在你的观测下产生 ",
      13, NAVY, False),
     ("相同的统计行为", 13, RED, True), (".", 13, NAVY, False)],
    [("→ 这不是『噪声大』而是『歧义』 — 任何 ML 算法 (包括 TSFM) 在歧义点上都退化为随机猜.",
      12, DGREY, False)],
    [("→ 解决方法: 改变 ", 12, ACCENT, True),
     ("观测设计", 12, NAVY, True),
     (" — 加传感器 / 主动激励 / 干预实验.", 12, ACCENT, True)],
])

# Examples
exs = [
    ("控制系统",
     "稳态下 (ẋ = 0), 多个 (A, b) 给出相同观测",
     "需要主动施加扰动 (PRBS) 才能识别参数.\n→ TSFM 在被动观测下无能为力.",
     RED),
    ("生物 / 流行病",
     "感染率 + 报告率 不可分",
     "如果只看新增病例数, 无法分离 真实感染率 和 报告偏差.\n→ 必须有血清学采样作为额外观测.",
     PURPLE),
    ("金融",
     "因果效应 vs 相关",
     "时序中两个变量同向变化, 可能因果反过来或共同被第三因素驱动.\n→ TSFM 只学相关, 不能给出因果.",
     ORANGE),
]
gx = Inches(0.5); gy = Inches(3.05); cw = Inches(4.05); ch = Inches(3.85)
for i, (title, sub, body, col) in enumerate(exs):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.55), col)
    add_text(s, x, gy, cw, Inches(0.55), title, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), gy + Inches(0.7), cw - Inches(0.4), Inches(0.7),
             sub, size=11, italic=True, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), gy + Inches(1.55), cw - Inches(0.4), Inches(2.2),
             body, size=11, color=DGREY)

# ============================================================
# Slide 20 — Engineering implication: match bottleneck to axis
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 20, TOTAL, "工程启示 · 把瓶颈对应到轴, 再选工具",
     "不要用 TSFM 解决它不擅长的问题; 也不要忽视它能解决的问题",
     section="PART VI · 启示与决策")

# Decision matrix table
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.5), DEEP)
add_text(s, Inches(0.7), Inches(1.55), Inches(3.0), Inches(0.5),
         "瓶颈所在轴", size=12, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(4.5), Inches(1.55), Inches(2.0), Inches(0.5),
         "TSFM 帮得上吗?", size=12, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.0), Inches(1.55), Inches(5.6), Inches(0.5),
         "正确的下一步", size=12, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

rows = [
    ("Axis 1 · 信息完备性",  "✗ 不能", "加传感器 / 加协变量 / 改善观测设计", RED),
    ("Axis 2 · 表征能力",     "✓ 能",   "选 100M-1B 级 TSFM (TimesFM / Chronos / Timer-XL)", GREEN),
    ("Axis 3 · 归纳偏置",     "✓✓ 强",  "直接用 TSFM (内置 Patching+RevIN+CI 等先验)",   GREEN),
    ("Axis 4 · 训练数据",     "✓✓✓ 极强", "TSFM zero-shot / few-shot / LoRA",          GREEN),
    ("Axis 5 · 泛化机制",     "✓✓✓ 极强", "TSFM ICL — 不需要每个任务重训",            GREEN),
    ("Axis 6 · 可识别性",     "✗ 不能",  "改变观测策略 (主动扰动 / 干预实验)",        RED),
    ("Axis 7 · 优化",          "○ 中性",  "成熟工程栈, 选择已知稳定的 base model",      GREY),
    ("Axis 8 · 校准",          "✓ 帮", "TSFM 多数自带概率输出; 后续 conformal 校准",    GREEN),
]
y = Inches(2.05); rh = Inches(0.55)
for i, (axis, ok, action, col) in enumerate(rows):
    add_rect(s, Inches(0.5), y, Inches(12.3), rh - Inches(0.05),
             WHITE if i % 2 == 0 else LIGHT, line=SOFT)
    add_text(s, Inches(0.7), y + Inches(0.13), Inches(3.7), Inches(0.3),
             axis, size=11, bold=True, color=NAVY)
    add_text(s, Inches(4.5), y + Inches(0.13), Inches(2.0), Inches(0.3),
             ok, size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.0), y + Inches(0.13), Inches(5.6), Inches(0.3),
             action, size=11, color=DGREY)
    y = y + rh

# ============================================================
# Slide 21 — Where does TSFM stand among forecasting methods?
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 21, TOTAL, "把 TSFM 放回时序方法的全景",
     "不同方法在 6 轴上的相对位置 — 一目了然",
     section="PART VI · 启示与决策")

# Methods comparison: 4 methods × 6 axes
# rows = methods, cols = axes
methods = ["传统 ARIMA / ETS", "深度学习 (LSTM/TCN/TFT)", "本任务监督 + 协变量", "TSFM (zero-shot)", "TSFM + LoRA + 物理"]
scores = [
    [3, 1, 2, 1, 2, 0],   # ARIMA
    [3, 4, 2, 2, 2, 0],   # DL
    [4, 3, 3, 3, 2, 0],   # supervised + cov
    [2, 4, 5, 5, 5, 0],   # TSFM zero-shot
    [5, 5, 5, 4, 4, 0],   # TSFM + LoRA + physics
]
ax_names = ["轴1\n信息", "轴2\n容量", "轴3\n偏置", "轴4\n数据", "轴5\n泛化", "轴6\n识别"]

# Header
hdr_x = Inches(3.7); col_w = Inches(1.55)
add_rect(s, Inches(0.5), Inches(1.55), Inches(3.0), Inches(0.65), DEEP)
add_text(s, Inches(0.7), Inches(1.55), Inches(2.6), Inches(0.65),
         "方法", size=13, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
for j, ax in enumerate(ax_names):
    x = hdr_x + j * col_w
    add_rect(s, x, Inches(1.55), col_w - Inches(0.03), Inches(0.65), DEEP)
    add_text(s, x, Inches(1.55), col_w, Inches(0.65),
             ax, size=10, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Rows
for i, (name, row) in enumerate(zip(methods, scores)):
    y = Inches(2.25) + i * Inches(0.85)
    add_rect(s, Inches(0.5), y, Inches(3.0), Inches(0.8),
             LIGHT if i % 2 == 0 else WHITE, line=SOFT)
    add_text(s, Inches(0.7), y, Inches(2.6), Inches(0.8),
             name, size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    for j, score in enumerate(row):
        x = hdr_x + j * col_w
        add_rect(s, x, y, col_w - Inches(0.03), Inches(0.8),
                 LIGHT if i % 2 == 0 else WHITE, line=SOFT)
        # filled circles to indicate score (0-5)
        cx_d = x + col_w / 2 - Inches(0.55)
        cy_d = y + Inches(0.32)
        for k in range(5):
            cc = ACCENT if k < score else SOFT
            add_oval(s, cx_d + k * Inches(0.22), cy_d, Inches(0.18), Inches(0.18), cc)

# Bottom note
add_rect(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.5), LIGHT)
add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.5),
         "→  TSFM 的优势集中在 轴 3-5; 受迫系统中, 必须搭配协变量补足 轴 1; 轴 6 永远在观测端解决.",
         size=12, italic=True, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 22 — Final answer to user's question
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 22, TOTAL, "最终回答你的问题",
     "总结: 你的两轴框架是对的, 但不完整; TSFM 在 4 个被忽略的轴上做了文章",
     section="结论")

# Three answer blocks
answers = [
    ("Q1", "你对时序预测的两轴认知准确吗?",
     "准确 — 但只是必要条件, 不是充分.",
     "信息完备性 ↔ 不可约误差下界 σ²(Y|X);  表征能力 ↔ 近似误差 ‖f*−f̂‖².\n两者必须满足, 但还要加上至少 4 个被你的框架忽略的宏观维度.",
     ACCENT),
    ("Q2", "在宏观上还需要补充什么?",
     "至少 4 个: 归纳偏置, 训练数据规模 (≠per-instance 上下文), 泛化机制, 可识别性. 加上次要的 优化 + 校准.",
     "其中 可识别性 比信息完备性更上游 (观测语义层面); 数据规模 + 泛化机制 共同决定 ICL 这种新型迁移能否涌现.",
     GREEN),
    ("Q3", "TSFM 在哪些维度上做了文章?",
     "在 4 个轴上同时发力: 数据规模, 归纳偏置, 泛化机制, 表征容量. 副产品: 缩放律.",
     "重要的是 协同性 — 把容量 / 偏置 / 数据 / 机制 同时推到了新水平; 不是单一维度的胜利.\n它 不解决 信息完备性 + 可识别性 — 这两个是观测端问题, 工程上必须另想办法.",
     PURPLE),
]
gy = Inches(1.55)
for i, (q, sub, head, body, col) in enumerate(answers):
    y = gy + i * Inches(1.85)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.7), WHITE, line=col, line_w=Pt(2))
    add_rect(s, Inches(0.5), y, Inches(0.85), Inches(1.7), col)
    add_text(s, Inches(0.5), y, Inches(0.85), Inches(1.7),
             q, size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.5), y + Inches(0.1), Inches(11.0), Inches(0.4),
             sub, size=12, bold=True, color=NAVY, italic=True)
    add_text(s, Inches(1.5), y + Inches(0.5), Inches(11.0), Inches(0.45),
             head, size=14, bold=True, color=col)
    add_text(s, Inches(1.5), y + Inches(1.0), Inches(11.0), Inches(0.65),
             body, size=11, color=DGREY)

# Save
out = "/home/aicode/sherwin/TSFM/paper/TSFM_时序预测的宏观维度.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
