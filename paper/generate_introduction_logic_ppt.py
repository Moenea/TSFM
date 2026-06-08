"""Introduction logic-derivation PPT for the ASProNet manuscript.

Walks through the logical chain that *justifies* studying abnormal situation
prognosis (ASP), step by step, faithfully following the manuscript Introduction
in TSFM/paper/manuscript_abstract_introduction.docx.

Logical chain (5 parts):
  PART I   · 安全压力        → 重大事故和经济损失数字
  PART II  · 行业现状        → FDD 主导, 反应式 ASM
  PART III · 范式跃迁        → Prognosis = FDD 的前向延伸
  PART IV  · 障碍 (失配)     → 数据失衡 → 传统 forecasting 偏向正常
  PART V   · 方法论          → Transformer + MoE → ASProNet

Style mirrors the existing TSFM_*.pptx files (navy + accent orange,
Microsoft YaHei, 16:9, 22pt page titles, dense bullet content).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY   = RGBColor(0x0B, 0x1F, 0x3A)
DEEP   = RGBColor(0x14, 0x2C, 0x52)
ACCENT = RGBColor(0xE8, 0x8B, 0x1A)
LIGHT  = RGBColor(0xF4, 0xF6, 0xFA)
GREY   = RGBColor(0x4A, 0x55, 0x68)
DGREY  = RGBColor(0x33, 0x3A, 0x47)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SOFT   = RGBColor(0xDC, 0xE3, 0xF0)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
RED    = RGBColor(0xC0, 0x39, 0x2B)
PURPLE = RGBColor(0x6B, 0x3F, 0xA0)
TEAL   = RGBColor(0x16, 0x7E, 0x8A)
ORANGE = RGBColor(0xE0, 0x6E, 0x18)
GOLD   = RGBColor(0xD4, 0xA0, 0x2A)
SAND   = RGBColor(0xFB, 0xF1, 0xDF)

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _i(v):
    return int(v) if v is not None else v


def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_round(slide, x, y, w, h, fill, line=None):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_oval(slide, x, y, w, h, fill, line=None):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def rich(slide, x, y, w, h, blocks, line_spacing=1.3, anchor=MSO_ANCHOR.TOP):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, runs in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        for txt, size, color, bold in runs:
            r = p.add_run(); r.text = txt
            r.font.name = FONT; r.font.size = Pt(size)
            r.font.bold = bold; r.font.color.rgb = color
    return tb


def line(slide, x1, y1, x2, y2, color, w=1.8):
    x1, y1, x2, y2 = _i(x1), _i(y1), _i(x2), _i(y2)
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color; ln.line.width = Pt(w)
    return ln


def arrow(slide, x1, y1, x2, y2, color=ACCENT, w=2.0):
    x1, y1, x2, y2 = _i(x1), _i(y1), _i(x2), _i(y2)
    ln = slide.shapes.add_connector(2, x1, y1, x2, y2)
    ln.line.color.rgb = color; ln.line.width = Pt(w)
    return ln


# -- Page chrome --------------------------------------------------------------
SUBTITLE_BAR = "ASProNet · Introduction 逻辑推导 · Jinsong Zhao"
TOTAL = 16


def page(slide, idx, total, title, subtitle=None, section=None):
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_rect(slide, 0, 0, SW, Inches(0.08), ACCENT)
    add_rect(slide, 0, Inches(0.08), Inches(0.18), SH, NAVY)
    if section:
        add_text(slide, Inches(0.5), Inches(0.18), Inches(11.5), Inches(0.3),
                 section, size=11, bold=True, color=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.6),
             title, size=22, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(12.0), Inches(0.4),
                 subtitle, size=12, color=GREY, italic=True)
    add_rect(slide, Inches(0.5), Inches(1.28), Inches(0.6), Emu(28000), ACCENT)
    add_text(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
             f"{idx} / {total}", size=10, color=GREY, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.3),
             SUBTITLE_BAR, size=10, color=GREY)


def section_divider(num, label, en_label, idx):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, 0, SW, Inches(0.15), ACCENT)
    add_rect(s, Inches(0.8), Inches(2.6), Inches(0.7), Emu(30000), ACCENT)
    add_text(s, Inches(0.8), Inches(2.0), Inches(8), Inches(0.6),
             f"PART {num}", size=22, bold=True, color=ACCENT)
    add_text(s, Inches(0.8), Inches(2.85), Inches(11.5), Inches(0.9),
             label, size=40, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(4.05), Inches(11.5), Inches(0.5),
             en_label, size=18, color=SOFT, italic=True)
    add_text(s, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
             f"{idx} / {TOTAL}", size=10, color=SOFT, align=PP_ALIGN.RIGHT)
    return s


# =============================================================================
# Slide 1 · Title
# =============================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, 0, SW, Inches(0.15), ACCENT)
for i in range(8):
    add_rect(s, Inches(0.5 + i * 1.6), Inches(0.4), Inches(0.6), Emu(15000), ACCENT)

add_text(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8),
         "为什么要做异常工况预警?",
         size=36, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.5),
         "Why Abnormal Situation Prognosis — Introduction 逻辑推导",
         size=18, color=ACCENT, italic=True)
add_rect(s, Inches(0.8), Inches(2.95), Inches(0.7), Emu(30000), ACCENT)

rich(s, Inches(0.8), Inches(3.15), Inches(11.5), Inches(3.4), [
    [("基于 ", 16, SOFT, False),
     ("ASProNet 手稿 Introduction", 16, ACCENT, True),
     (" 的逐句逻辑展开 ", 16, SOFT, False),
     ("(manuscript_abstract_introduction.docx)", 14, GREY, False)],
    [("", 8, SOFT, False)],
    [("一句话提炼:", 16, ACCENT, True)],
    [("过程安全 ", 14, SOFT, False),
     ("⇒ 现行 FDD 是反应式 ", 14, WHITE, False),
     ("⇒ 预警是 FDD 的前向延伸 ", 14, WHITE, False),
     ("⇒ 但数据失衡使常规 forecasting 失效 ", 14, WHITE, False),
     ("⇒ 需要 MoE Transformer 专门化", 14, ACCENT, True)],
    [("", 6, SOFT, False)],
    [("• ", 14, ACCENT, True),
     ("PART I · 安全后果", 14, WHITE, True),
     (": 事故、伤亡、年 $47B 经济损失 — ASM 必须存在", 14, SOFT, False)],
    [("• ", 14, ACCENT, True),
     ("PART II · 现状极限", 14, WHITE, True),
     (": FDD 主导, 但 reactive + alarm flood — 解决得不够早", 14, SOFT, False)],
    [("• ", 14, ACCENT, True),
     ("PART III · 范式跃迁", 14, WHITE, True),
     (": Prognosis 把决策点 提前 到 alarm 触发 ", 14, SOFT, False),
     ("之前", 14, ACCENT, True)],
    [("• ", 14, ACCENT, True),
     ("PART IV · 障碍 (失配)", 14, WHITE, True),
     (": normal-dominant 数据 + global-MSE 训练 ⇒ 模型偏向正常态", 14, SOFT, False)],
    [("• ", 14, ACCENT, True),
     ("PART V · 方法论", 14, WHITE, True),
     (": MoE-Transformer + tail-aware = ASProNet", 14, ACCENT, True)],
])

add_rect(s, Inches(0.8), Inches(6.7), Inches(11.5), Emu(15000), ACCENT)
add_text(s, Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.4),
         "Jinsong Zhao   ·   清华大学   ·   2026-05-15",
         size=14, color=ACCENT)


# =============================================================================
# Slide 2 · The full logic chain (preview)
# =============================================================================
s = prs.slides.add_slide(BLANK)
page(s, 2, TOTAL, "整篇 Introduction 的逻辑骨架",
     "5 个 PART · 9 步推理 · 一张总图",
     section="ROADMAP · 逻辑全景")

chain_y = Inches(1.85)
chain_h = Inches(4.5)
add_rect(s, Inches(0.5), chain_y, Inches(12.4), chain_h, LIGHT)

steps = [
    ("①", "过程事故 ≈ 来自异常工况", "CCPS leading/lagging 数据"),
    ("②", "异常工况 → 重大伤亡 + 巨额损失", "$20B US, $27B UK · 4088 fatalities (CN)"),
    ("③", "传统做法 = FDD (Fault Detection & Diagnosis)", "ASM 的分析核心 · 但 decision point 在异常已可观测之后"),
    ("④", "FDD 限制 = reactive + alarm flood", "操作员认知超载 · 干预为时已晚"),
    ("⑤", "Prognosis = 在 alarm 触发之前预测", "FDD 框架的 forward extension"),
    ("⑥", "数据极度失衡: normal ≫ abnormal", "abnormal 观测稀缺"),
    ("⑦", "Global-MSE forecasting 会偏向 normal", "结构性 mismatch — 在最关键的区域反而最弱"),
    ("⑧", "模型必须: 抓时序 + 对工况敏感", "Transformer 是自然起点, 但单模型仍偏正常态"),
    ("⑨", "MoE 让不同 expert 专精不同 regime",  "→ ASProNet (MoE + 双分辨率 + tail-aware)"),
]

# Layout 9 steps in a 3x3 grid
cols, rows = 3, 3
ox, oy = Inches(0.7), Inches(1.95)
cw, ch = Inches(4.0), Inches(1.40)
gx, gy = Inches(0.10), Inches(0.10)

for k, (num, head, sub) in enumerate(steps):
    r, c = divmod(k, cols)
    x = ox + c * (cw + gx)
    y = oy + r * (ch + gy)
    add_round(s, x, y, cw, ch, WHITE, line=NAVY)
    add_rect(s, x, y, Inches(0.55), ch, NAVY)
    add_text(s, x, y + Inches(0.05), Inches(0.55), Inches(0.55),
             num, size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    rich(s, x + Inches(0.65), y + Inches(0.12), cw - Inches(0.7), ch - Inches(0.2), [
        [(head, 12, NAVY, True)],
        [(sub, 10, GREY, False)],
    ], line_spacing=1.15)

add_text(s, Inches(0.5), Inches(6.55), Inches(12.4), Inches(0.4),
         "下面 13 张片 = 把每一步用一页讲清楚 (PART I → PART V).",
         size=13, color=DGREY, italic=True)


# =============================================================================
# PART I divider
# =============================================================================
section_divider("I", "安全后果", "Why Abnormal Situations Matter — Safety & Economic Cost", 3)


# =============================================================================
# Slide 4 · Step ① — ASM 是过程安全的核心 (CCPS)
# =============================================================================
s = prs.slides.add_slide(BLANK)
page(s, 4, TOTAL, "Step ①  几乎每一起重大事故 = 来自异常工况",
     "ASM (Abnormal Situation Management) 是化工安全不可缺少的核心",
     section="PART I · 安全后果")

# Left: CCPS quote box
add_round(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(3.4), LIGHT, line=NAVY)
add_text(s, Inches(0.75), Inches(1.85), Inches(5.5), Inches(0.45),
         "CCPS (AIChE Center for Chemical Process Safety) 关键观察",
         size=14, bold=True, color=NAVY)
rich(s, Inches(0.75), Inches(2.35), Inches(5.5), Inches(2.7), [
    [("『 ", 14, ACCENT, True),
     ("according to CCPS leading and lagging metrics", 14, DGREY, False),
     (" 』", 14, ACCENT, True)],
    [("", 6, NAVY, False)],
    [("→  ", 16, ACCENT, True),
     ("almost every major incident", 16, NAVY, True),
     (" is caused by an", 16, NAVY, False)],
    [("    ", 16, NAVY, False),
     ("abnormal situation.", 18, RED, True)],
    [("", 8, NAVY, False)],
    [("含义 : ", 13, NAVY, True),
     ("安全管理的胜负 ", 13, DGREY, False),
     ("=", 13, ACCENT, True),
     (" 异常工况能否被及时识别和介入. ", 13, DGREY, False)],
    [("        ASM 不是 nice-to-have, ", 13, DGREY, False),
     ("而是 industrial safety 的核心组件", 13, NAVY, True),
     (".", 13, DGREY, False)],
])

# Right: hierarchy "ASM = analytical core" diagram
add_text(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.4),
         "ASM 在化工安全管理中的位置", size=14, bold=True, color=NAVY)

add_round(s, Inches(7.0), Inches(2.15), Inches(5.8), Inches(0.7), NAVY)
add_text(s, Inches(7.0), Inches(2.25), Inches(5.8), Inches(0.5),
         "Process Safety", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

arrow(s, Inches(9.9), Inches(2.85), Inches(9.9), Inches(3.15), color=ACCENT)

add_round(s, Inches(7.0), Inches(3.15), Inches(5.8), Inches(0.7), DEEP)
add_text(s, Inches(7.0), Inches(3.25), Inches(5.8), Inches(0.5),
         "Abnormal Situation Management (ASM)", size=15, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)

arrow(s, Inches(9.9), Inches(3.85), Inches(9.9), Inches(4.15), color=ACCENT)

add_round(s, Inches(7.0), Inches(4.15), Inches(2.85), Inches(0.7), SAND)
add_text(s, Inches(7.0), Inches(4.25), Inches(2.85), Inches(0.5),
         "FDD (今日主流)", size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_round(s, Inches(9.95), Inches(4.15), Inches(2.85), Inches(0.7), GOLD)
add_text(s, Inches(9.95), Inches(4.25), Inches(2.85), Inches(0.5),
         "Prognosis (新范式)", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, Inches(7.0), Inches(5.05), Inches(5.8), Inches(0.4),
         "→ 全文的故事: 从左侧的 FDD 推到右侧的 Prognosis.",
         size=12, color=DGREY, italic=True)

# Bottom: reference tag
add_rect(s, Inches(0.5), Inches(6.05), Inches(12.4), Emu(15000), ACCENT)
add_text(s, Inches(0.5), Inches(6.20), Inches(12.4), Inches(0.4),
         "Manuscript 引用 : [1-3] ASM 综述 ·  [4] CCPS leading/lagging metrics ·  [16,20] FDD 综述",
         size=12, color=NAVY)


# =============================================================================
# Slide 5 · Step ② — 后果数字
# =============================================================================
s = prs.slides.add_slide(BLANK)
page(s, 5, TOTAL, "Step ②  异常工况的代价: 伤亡 + 巨额经济损失",
     "数字告诉我们: 这不是学术 toy problem",
     section="PART I · 安全后果")

# Three big number cards
cards = [
    (NAVY, "4,088", "fatalities", "中国 20 年间化工事故死亡总数",
     "  最严重一起 62 人遇难, 直接经济损失 $100M", RED),
    (DEEP, "$20B / yr", "U.S. petroleum + chem",
     "异常工况导致的直接损失", "  英国为 $27B / yr — 合计 ~$47B / yr", ORANGE),
    (PURPLE, "off-spec · shutdown · damage · emissions",
     "minor incidents (everyday)",
     "异常工况不仅引发大事故, 还日常造成:",
     "  • 不合规产品   • 非计划停车  • 设备损坏  • 排放超标",
     GOLD),
]

cw = Inches(4.0); gx = Inches(0.15); ox = Inches(0.55); oy = Inches(1.85); ch = Inches(4.0)
for k, (bg, big, lab, head, body, tone) in enumerate(cards):
    x = ox + k * (cw + gx)
    add_round(s, x, oy, cw, ch, bg)
    add_text(s, x + Inches(0.25), oy + Inches(0.20), cw - Inches(0.5), Inches(0.4),
             "Big Number" if k < 2 else "Minor But Recurrent",
             size=11, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.25), oy + Inches(0.65), cw - Inches(0.5), Inches(1.4),
             big, size=26 if k < 2 else 14, bold=True, color=WHITE)
    add_text(s, x + Inches(0.25), oy + Inches(1.95) if k < 2 else oy + Inches(1.6),
             cw - Inches(0.5), Inches(0.4),
             lab, size=12, bold=True, color=tone, italic=True)
    add_text(s, x + Inches(0.25), oy + Inches(2.35) if k < 2 else oy + Inches(2.0),
             cw - Inches(0.5), Inches(1.3),
             head, size=12, color=SOFT)
    add_text(s, x + Inches(0.25), oy + Inches(2.9) if k < 2 else oy + Inches(2.5),
             cw - Inches(0.5), Inches(1.0),
             body, size=11, color=SAND)

add_rect(s, Inches(0.5), Inches(6.1), Inches(12.4), Inches(0.7), LIGHT)
add_text(s, Inches(0.7), Inches(6.20), Inches(12.0), Inches(0.4),
         "推理桥梁:", size=13, bold=True, color=NAVY)
add_text(s, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.4),
         "既然代价如此巨大, 那么 ASM 必须工作; 看看现状 ASM 是怎么做的 → PART II.",
         size=13, color=DGREY, italic=True)


# =============================================================================
# PART II divider
# =============================================================================
section_divider("II", "现状的极限", "Industrial Practice = FDD = Reactive ASM", 6)


# =============================================================================
# Slide 7 · Step ③ — FDD 是 ASM 的分析核心
# =============================================================================
s = prs.slides.add_slide(BLANK)
page(s, 7, TOTAL, "Step ③  现行做法: Fault Detection & Diagnosis (FDD)",
     "工业 ASM 二十年来的分析支柱 — 但它的工作点在哪里?",
     section="PART II · 现状极限")

# FDD definition box
add_round(s, Inches(0.5), Inches(1.7), Inches(7.0), Inches(2.4), LIGHT, line=NAVY)
add_text(s, Inches(0.75), Inches(1.85), Inches(6.7), Inches(0.45),
         "FDD 的定义 (manuscript 原话)", size=14, bold=True, color=NAVY)
rich(s, Inches(0.75), Inches(2.35), Inches(6.7), Inches(1.7), [
    [("① ", 14, ACCENT, True),
     ("Detection: ", 14, NAVY, True),
     ("identifies whether the current process state has", 13, DGREY, False)],
    [("    ", 13, DGREY, False),
     ("departed from normal operation", 13, NAVY, True),
     (".", 13, DGREY, False)],
    [("", 6, NAVY, False)],
    [("② ", 14, ACCENT, True),
     ("Diagnosis: ", 14, NAVY, True),
     ("determines the ", 13, DGREY, False),
     ("likely cause / fault class", 13, NAVY, True),
     (" of the deviation.", 13, DGREY, False)],
    [("", 6, NAVY, False)],
    [("=> 它回答 ", 13, DGREY, False),
     ("『现在是不是异常 & 是哪种异常』", 13, NAVY, True),
     (".", 13, DGREY, False)],
])

# Right: timeline showing FDD's decision point
add_text(s, Inches(7.85), Inches(1.7), Inches(5.0), Inches(0.4),
         "FDD 的 decision point 在时间轴上的位置", size=14, bold=True, color=NAVY)

tl_y = Inches(2.5)
add_rect(s, Inches(7.85), tl_y, Inches(5.0), Inches(0.04), NAVY)
# events on timeline
for (xfrac, lbl, col, dy) in [(0.05, "Normal", GREEN, -0.55),
                               (0.45, "Precursor", GOLD, -0.55),
                               (0.65, "Abnormality\nobservable", ORANGE, -0.85),
                               (0.85, "FDD fires", RED, -0.55)]:
    cx = Inches(7.85) + Inches(5.0) * xfrac
    add_oval(s, cx - Inches(0.08), tl_y - Inches(0.08), Inches(0.16), Inches(0.16), col)
    add_text(s, cx - Inches(0.9), tl_y + Inches(dy), Inches(1.8), Inches(0.7),
             lbl, size=10, bold=True, color=col, align=PP_ALIGN.CENTER)

# Arrow showing missed window
add_text(s, Inches(7.85), Inches(3.55), Inches(5.0), Inches(0.4),
         "↑ FDD 只能在异常 已经可观测 后才会响应", size=11, color=RED, italic=True)
add_text(s, Inches(7.85), Inches(3.85), Inches(5.0), Inches(0.4),
         "= 干预窗口被错过了", size=12, bold=True, color=RED)

# Bottom: takeaway box
add_round(s, Inches(0.5), Inches(4.4), Inches(12.4), Inches(1.6), SAND)
rich(s, Inches(0.75), Inches(4.55), Inches(12.0), Inches(1.4), [
    [("一句话: ", 16, ACCENT, True),
     ("FDD 是必要的, 但它的工作点是 ", 14, NAVY, False),
     ("『异常已经发生』", 14, RED, True),
     (" — 不是 ", 14, NAVY, False),
     ("『异常即将发生』", 14, GREEN, True),
     (".", 14, NAVY, False)],
    [("", 8, NAVY, False)],
    [("引用: ", 12, DGREY, True),
     ("[2,16,20] FDD 是 ASM 的核心分析组件 ; ", 12, DGREY, False),
     ("[20,22] FDD 的 decision point tied to observable abnormality.",
      12, DGREY, False)],
])

add_text(s, Inches(0.5), Inches(6.30), Inches(12.4), Inches(0.4),
         "下一页: 看看 FDD 在异常 ‘已发生后’ 这件事, 还会带来什么次生问题.",
         size=12, color=DGREY, italic=True)


# =============================================================================
# Slide 8 · Step ④ — FDD 的代价: alarm flood + reactive workflow
# =============================================================================
s = prs.slides.add_slide(BLANK)
page(s, 8, TOTAL, "Step ④  FDD 的结构性缺陷 — Reactive + Alarm Flood",
     "等到异常浮现, 操作员已经被淹没在告警里",
     section="PART II · 现状极限")

# Left: alarm flood visualization
add_text(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
         "Alarm flood: 异常事件发生瞬间, 告警从几条变成几十条", size=13, bold=True, color=NAVY)

# Visual: time axis with alarm bars
av_y = Inches(2.25)
add_rect(s, Inches(0.5), av_y + Inches(2.0), Inches(6.0), Inches(0.04), NAVY)
import random
random.seed(7)
for k in range(80):
    if k < 50:
        height = Inches(0.25 + 0.4 * random.random())  # baseline noise
    else:
        height = Inches(0.7 + 1.2 * random.random())   # flood
    x = Inches(0.5) + Inches(6.0) * (k / 80)
    col = GREEN if k < 50 else (GOLD if k < 58 else RED)
    add_rect(s, x, av_y + Inches(2.0) - height, Inches(0.06), height, col)

add_text(s, Inches(0.5), Inches(4.45), Inches(2.5), Inches(0.4),
         "正常运行 ⤺", size=11, color=GREEN, italic=True)
add_text(s, Inches(4.0), Inches(4.45), Inches(2.5), Inches(0.4),
         "异常浮现 = 告警洪流", size=11, color=RED, italic=True, bold=True)

# Right: consequences box
add_round(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(3.3), LIGHT, line=NAVY)
add_text(s, Inches(7.05), Inches(1.85), Inches(5.7), Inches(0.45),
         "→ 已知的连锁后果 (alarm-management 文献)",
         size=14, bold=True, color=NAVY)
rich(s, Inches(7.05), Inches(2.35), Inches(5.7), Inches(2.6), [
    [("• ", 14, ACCENT, True),
     ("Operator overwhelm", 13, NAVY, True),
     (": 短时间内告警数 ↑↑, ", 12, DGREY, False)],
    [("    认知带宽撑爆.", 12, DGREY, False)],
    [("", 6, NAVY, False)],
    [("• ", 14, ACCENT, True),
     ("Critical interventions delayed", 13, NAVY, True),
     (": 即使知道异常, ", 12, DGREY, False)],
    [("    也来不及在合适时刻去做对的事.", 12, DGREY, False)],
    [("", 6, NAVY, False)],
    [("• ", 14, ACCENT, True),
     ("ASM 工作流呈现 ", 13, NAVY, True),
     ("reactive 形态", 13, RED, True),
     (": ", 12, DGREY, False)],
    [("    corrective action only after the process has entered an ", 12, DGREY, False),
     ("abnormal", 12, RED, True),
     (".", 12, DGREY, False)],
])

# Bottom: bridge to PART III
add_round(s, Inches(0.5), Inches(5.4), Inches(12.4), Inches(1.3), SAND)
rich(s, Inches(0.75), Inches(5.55), Inches(12.0), Inches(1.1), [
    [("逻辑出口: ", 16, ACCENT, True),
     ("既然 ‘等异常浮现再响应’ 的代价是 alarm flood + reactive, ",
      14, NAVY, False)],
    [("                     那么应该把 ASM 的 ", 14, NAVY, False),
     ("decision point", 14, ACCENT, True),
     (" 提前到 ", 14, NAVY, False),
     ("『alarm 触发之前』", 14, GREEN, True),
     (" — 这就是 Prognosis. ", 14, NAVY, False)],
    [("引用 : ", 12, DGREY, True),
     ("[23,24] alarm system / management 综述 ;  [3,16,20] ASM 现状的 reactive 性质.",
      12, DGREY, False)],
])


# =============================================================================
# PART III divider
# =============================================================================
section_divider("III", "范式跃迁", "Paradigm Shift: From Reactive FDD → Proactive Prognosis", 9)


# =============================================================================
# Slide 10 · Step ⑤ — Prognosis = FDD 的 forward extension
# =============================================================================
s = prs.slides.add_slide(BLANK)
page(s, 10, TOTAL, "Step ⑤  把决策点提前: Abnormal Situation Prognosis",
     "Prognosis 不是替代 FDD, 而是 FDD 框架的 forward extension",
     section="PART III · 范式跃迁")

# Top: definition strip
add_round(s, Inches(0.5), Inches(1.7), Inches(12.4), Inches(1.1), LIGHT, line=NAVY)
rich(s, Inches(0.75), Inches(1.85), Inches(12.0), Inches(0.9), [
    [("Prognosis (manuscript 定义): ", 14, ACCENT, True),
     ("predict future trends of process variables to determine whether",
      13, NAVY, False)],
    [("an abnormal situation is developing ", 13, NAVY, True),
     ("before an alarm triggers", 13, RED, True),
     (".", 13, NAVY, False)],
    [("=>  ", 13, ACCENT, True),
     ("把 ASM 的 decision point 从 ‘异常已发生’ 移到 ‘异常正在形成’; ",
      13, NAVY, True),
     ("提供 critical intervention window.", 13, NAVY, False)],
])

# Two side-by-side comparison cards: FDD vs Prognosis
cy = Inches(3.0); ch2 = Inches(2.7); cw2 = Inches(6.0); gx2 = Inches(0.4)

# FDD card (left, reactive)
add_round(s, Inches(0.5), cy, cw2, ch2, SAND)
add_rect(s, Inches(0.5), cy, cw2, Inches(0.45), DEEP)
add_text(s, Inches(0.5), cy + Inches(0.05), cw2, Inches(0.4),
         "传统 FDD — Reactive", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rich(s, Inches(0.75), cy + Inches(0.6), cw2 - Inches(0.5), ch2 - Inches(0.7), [
    [("•  ", 14, NAVY, True),
     ("输入:", 13, NAVY, True),
     (" 当前 + 历史 process state", 13, DGREY, False)],
    [("•  ", 14, NAVY, True),
     ("决策:", 13, NAVY, True),
     (" 现在是否异常 ? 是哪种异常 ?", 13, DGREY, False)],
    [("•  ", 14, NAVY, True),
     ("时间点:", 13, NAVY, True),
     (" 异常已可观测 之后", 13, RED, True)],
    [("•  ", 14, NAVY, True),
     ("典型代价:", 13, NAVY, True),
     (" alarm flood, intervention delay", 13, RED, False)],
    [("•  ", 14, NAVY, True),
     ("操作哲学:", 13, NAVY, True),
     (" 出事再补救", 13, DGREY, False)],
])

# Prognosis card (right, proactive)
add_round(s, Inches(0.5) + cw2 + gx2, cy, cw2, ch2, RGBColor(0xDA, 0xEE, 0xDC))
add_rect(s, Inches(0.5) + cw2 + gx2, cy, cw2, Inches(0.45), GREEN)
add_text(s, Inches(0.5) + cw2 + gx2, cy + Inches(0.05), cw2, Inches(0.4),
         "Prognosis — Proactive", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rich(s, Inches(0.5) + cw2 + gx2 + Inches(0.25), cy + Inches(0.6), cw2 - Inches(0.5), ch2 - Inches(0.7), [
    [("•  ", 14, NAVY, True),
     ("输入:", 13, NAVY, True),
     (" 当前 + 历史 process state", 13, DGREY, False)],
    [("•  ", 14, NAVY, True),
     ("决策:", 13, NAVY, True),
     (" 未来 H 步内 alarm 是否会触发 ?", 13, DGREY, False)],
    [("•  ", 14, NAVY, True),
     ("时间点:", 13, NAVY, True),
     (" alarm 触发 之前", 13, GREEN, True)],
    [("•  ", 14, NAVY, True),
     ("收益:", 13, NAVY, True),
     (" critical intervention window", 13, GREEN, True)],
    [("•  ", 14, NAVY, True),
     ("操作哲学:", 13, NAVY, True),
     (" 防患于未然", 13, NAVY, False)],
])

# Bottom: bridge — but it sounds easy, what's hard?
add_round(s, Inches(0.5), Inches(5.9), Inches(12.4), Inches(0.85), LIGHT)
rich(s, Inches(0.75), Inches(6.0), Inches(12.0), Inches(0.7), [
    [("听起来很合理 — 那为什么 prognosis 还没成为标准做法 ?  ",
      14, NAVY, True),
     ("PART IV 给出答案: ", 14, DGREY, False),
     ("数据本身的失衡使得『直接套 forecasting 模型』会失败. ",
      14, ACCENT, True)],
    [("引用 : ", 12, DGREY, True),
     ("[22,25-27] prognosis 把 ASM 从 reactive 推向 proactive ;  [3,16,22,28] critical intervention window.",
      12, DGREY, False)],
])


# =============================================================================
# PART IV divider
# =============================================================================
section_divider("IV", "障碍 (失配)",
                 "The Imbalance Trap — Why Naïve Forecasting Fails", 11)


# =============================================================================
# Slide 12 · Step ⑥ — 数据失衡的事实
# =============================================================================
s = prs.slides.add_slide(BLANK)
page(s, 12, TOTAL, "Step ⑥  问题的本质: 数据极度失衡",
     "Normal 工况几乎占满所有观测; abnormal 是 极小尾部",
     section="PART IV · 障碍 (失配)")

# Left: bar / pie style imbalance visual
add_text(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
         "工业过程数据的典型构成", size=13, bold=True, color=NAVY)

bar_y = Inches(2.3); bar_h = Inches(1.0)
# Normal — huge bar
add_rect(s, Inches(0.5), bar_y, Inches(5.7), bar_h, GREEN)
add_text(s, Inches(0.5), bar_y + Inches(0.15), Inches(5.7), Inches(0.6),
         "Normal operating situations", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# Abnormal — tiny sliver
add_rect(s, Inches(6.2), bar_y, Inches(0.3), bar_h, RED)
add_text(s, Inches(6.6), bar_y + Inches(0.3), Inches(2.0), Inches(0.4),
         "← 几乎不可见的 abnormal", size=11, color=RED, italic=True, bold=True)

add_text(s, Inches(0.5), Inches(3.5), Inches(6.0), Inches(0.4),
         "→ Manuscript: ‘process data are overwhelmingly dominated by normal operating situations,",
         size=11, color=DGREY, italic=True)
add_text(s, Inches(0.5), Inches(3.8), Inches(6.0), Inches(0.4),
         "       whereas observations capturing abnormal situations are exceedingly scarce.’",
         size=11, color=DGREY, italic=True)

# Right: implications block
add_round(s, Inches(7.0), Inches(1.7), Inches(5.9), Inches(3.5), LIGHT, line=NAVY)
add_text(s, Inches(7.25), Inches(1.85), Inches(5.5), Inches(0.4),
         "对建模意味着什么", size=14, bold=True, color=NAVY)
rich(s, Inches(7.25), Inches(2.35), Inches(5.5), Inches(2.8), [
    [("• ", 14, ACCENT, True),
     ("样本数极不均衡:", 13, NAVY, True),
     (" abnormal 占比常为 ", 12, DGREY, False),
     ("< 1%", 12, RED, True),
     (" 甚至 ", 12, DGREY, False),
     ("< 0.1%", 12, RED, True),
     (".", 12, DGREY, False)],
    [("", 8, NAVY, False)],
    [("• ", 14, ACCENT, True),
     ("梯度被正常样本主导:", 13, NAVY, True),
     ("  loss landscape 几乎只由 normal 数据塑形.",
      12, DGREY, False)],
    [("", 8, NAVY, False)],
    [("• ", 14, ACCENT, True),
     ("RMSE 等全局指标 ‘看起来很好’:", 13, NAVY, True),
     ("  但在 abnormal 边界的预测精度并未被衡量.",
      12, DGREY, False)],
    [("", 8, NAVY, False)],
    [("• ", 14, ACCENT, True),
     ("结果:  ", 13, NAVY, True),
     ("最关键的少数样本被模型忽略", 13, RED, True),
     (".", 13, DGREY, False)],
])

# Bottom: bridge
add_rect(s, Inches(0.5), Inches(5.5), Inches(12.4), Emu(15000), ACCENT)
rich(s, Inches(0.5), Inches(5.65), Inches(12.4), Inches(1.2), [
    [("→ 下一步要追问 : 如果直接拿这种失衡数据训练常规的 time-series forecasting 模型, ",
      13, NAVY, True),
     ("会发生什么 ?", 13, ACCENT, True)],
    [("", 6, NAVY, False)],
    [("引用 : ", 12, DGREY, True),
     ("[22,23,30] abnormal 数据极度稀缺 ;  [30,37] 类别不平衡对深度学习的影响 ;  [31,32] 现有 forecasting 综述均默认 global error.",
      12, DGREY, False)],
])


# =============================================================================
# Slide 13 · Step ⑦ — Global-MSE 的结构性失配
# =============================================================================
s = prs.slides.add_slide(BLANK)
page(s, 13, TOTAL, "Step ⑦  常规 forecasting 在这里 ‘训练目标’ 就错了",
     "minimizing global error ≠ minimizing the error that matters",
     section="PART IV · 障碍 (失配)")

# Big formula card
add_round(s, Inches(0.5), Inches(1.7), Inches(12.4), Inches(1.4), LIGHT, line=NAVY)
rich(s, Inches(0.75), Inches(1.85), Inches(12.0), Inches(1.1), [
    [("常规 forecasting 训练目标 :   ", 15, NAVY, True),
     ("min  ", 16, ACCENT, True),
     ("E[ (y_pred − y_true)² ]", 16, DGREY, True),
     ("        (Global MSE over all timesteps)", 12, GREY, True)],
    [("", 6, NAVY, False)],
    [("隐含假设 : ", 14, ACCENT, True),
     ("Normal 和 Abnormal 两类时刻 ", 14, NAVY, False),
     ("有相同的运行意义", 14, RED, True),
     (".  ", 14, NAVY, False),
     ("⇒ 在工业现实中这一假设是错的.", 14, NAVY, True)],
])

# Left: failure illustration (cartoon error)
add_text(s, Inches(0.5), Inches(3.4), Inches(6.0), Inches(0.4),
         "模型的能力分配 (示意)", size=13, bold=True, color=NAVY)

# A "capacity bar" split between Normal and Abnormal
cap_y = Inches(3.95)
add_rect(s, Inches(0.5), cap_y, Inches(5.5), Inches(0.6), LIGHT, line=NAVY)
add_rect(s, Inches(0.5), cap_y, Inches(5.10), Inches(0.6), GREEN)  # ~93%
add_text(s, Inches(0.5), cap_y + Inches(0.1), Inches(5.10), Inches(0.4),
         "Capacity to model Normal  (~93%)", size=12, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(5.55), cap_y + Inches(0.1), Inches(0.5), Inches(0.4),
         "← 缺 ↓", size=10, color=RED, bold=True)

add_text(s, Inches(0.5), cap_y + Inches(0.85), Inches(6.0), Inches(0.4),
         "→ 在 abnormal 边界的预测精度 = ‘运行最关键的那部分’ = 模型给的最少.",
         size=11, color=DGREY, italic=True)

# Right: takeaway box — the structural mismatch
add_round(s, Inches(7.0), Inches(3.4), Inches(5.9), Inches(3.0), SAND)
add_text(s, Inches(7.25), Inches(3.55), Inches(5.5), Inches(0.45),
         "结构性失配 (Structural Mismatch)", size=15, bold=True, color=NAVY)
rich(s, Inches(7.25), Inches(4.05), Inches(5.5), Inches(2.4), [
    [("行业需要的: ", 14, ACCENT, True),
     ("在 abnormal 边界 ", 13, NAVY, False),
     ("准确预测", 13, GREEN, True),
     (".", 13, NAVY, False)],
    [("", 6, NAVY, False)],
    [("模型优化的: ", 14, ACCENT, True),
     ("在 ", 13, NAVY, False),
     ("normal regime ", 13, RED, True),
     ("准确预测.", 13, NAVY, False)],
    [("", 6, NAVY, False)],
    [("两者方向 ", 14, NAVY, True),
     ("不一致", 14, RED, True),
     (" ⇒  ", 14, NAVY, True),
     ("conventional approaches frequently",
      12, DGREY, True)],
    [("    fail to deliver reliable abnormal", 12, DGREY, True)],
    [("    situation prognosis.", 12, DGREY, True)],
])

# Bottom: bridge to PART V
add_rect(s, Inches(0.5), Inches(6.6), Inches(12.4), Emu(15000), ACCENT)
add_text(s, Inches(0.5), Inches(6.75), Inches(12.4), Inches(0.4),
         "推理出口: 要解决失配, 模型必须 (a) 抓时序依赖 + (b) 对工况敏感 → PART V.",
         size=13, color=NAVY, italic=True)


# =============================================================================
# PART V divider
# =============================================================================
section_divider("V", "方法论",
                 "From Imbalance to ASProNet — MoE Transformer with Tail-Aware Training", 14)


# =============================================================================
# Slide 15 · Step ⑧ + ⑨ — Transformer → MoE → ASProNet
# =============================================================================
s = prs.slides.add_slide(BLANK)
page(s, 15, TOTAL, "Step ⑧ + ⑨  从 ‘双重要求’ 到 ASProNet 的设计",
     "(a) 时序依赖 + (b) 工况敏感  ⇒  Transformer + Mixture-of-Experts",
     section="PART V · 方法论")

# Top: dual requirement
add_round(s, Inches(0.5), Inches(1.7), Inches(12.4), Inches(0.95), LIGHT, line=NAVY)
rich(s, Inches(0.75), Inches(1.82), Inches(12.0), Inches(0.8), [
    [("解决 mismatch 的两条 ‘必要条件’ :   ",
      14, NAVY, True)],
    [("  (a) ", 14, ACCENT, True),
     ("representing temporal dependencies in chemical process data", 13, NAVY, True),
     (",   (b) ", 14, ACCENT, True),
     ("remaining sensitive to operating situations", 13, NAVY, True),
     (".", 13, NAVY, False)],
])

# Three-step build-up: Transformer → +MoE → +Tail-aware
step_y = Inches(2.85)
step_h = Inches(3.5)
sw = Inches(4.0); sgx = Inches(0.13); sox = Inches(0.55)

step_specs = [
    ("Transformer", DEEP, WHITE, [
        ("自然起点", "在 FDD / prognosis 文献中 attention-based 模型已被验证 [22,35,36]."),
        ("能力", "全局 attention 捕获长程依赖,  适合化工 multivariable 时序."),
        ("缺陷", "single forecasting model 仍然把 表征容量花在 dominant normal patterns [32,33,37]."),
    ]),
    ("+ Mixture-of-Experts", PURPLE, WHITE, [
        ("核心思想", "多个 expert 各自专精不同 regime, gating 决定何时调谁 [38-40]."),
        ("为什么合身", "alarm threshold 附近的 dynamics 与 normal regime 显著不同 →  正适合 MoE 的『专家分工』."),
        ("产生模型", "ASProNet : MoE-Transformer with per-horizon gating."),
    ]),
    ("+ Tail-Aware Training", ACCENT, NAVY, [
        ("再补一刀", "样本距 alarm 越近, 给越大的训练权重."),
        ("效果", "把模型容量 '拉' 向 abnormal 边界."),
        ("最终目标", "improve prediction where prognosis is most operationally valuable: the early stage of an abnormal situation."),
    ]),
]

for k, (head, bg, txt_col, items) in enumerate(step_specs):
    x = sox + k * (sw + sgx)
    add_round(s, x, step_y, sw, step_h, WHITE, line=NAVY)
    add_rect(s, x, step_y, sw, Inches(0.5), bg)
    add_text(s, x, step_y + Inches(0.07), sw, Inches(0.4),
             head, size=15, bold=True, color=txt_col, align=PP_ALIGN.CENTER)
    blocks = []
    for h2, b2 in items:
        blocks.append([("● ", 13, ACCENT, True), (h2, 12, NAVY, True)])
        blocks.append([("    " + b2, 11, DGREY, False)])
        blocks.append([("", 4, NAVY, False)])
    rich(s, x + Inches(0.20), step_y + Inches(0.65), sw - Inches(0.4),
         step_h - Inches(0.8), blocks, line_spacing=1.2)
    if k < len(step_specs) - 1:
        arrow(s, x + sw, step_y + step_h / 2,
              x + sw + sgx, step_y + step_h / 2,
              color=ACCENT, w=2.5)

# Bottom: takeaway
add_text(s, Inches(0.5), Inches(6.55), Inches(12.4), Inches(0.4),
         "ASProNet = MoE-Transformer · 双分辨率 patch · 渐进 transient/steady 分解 · per-horizon gating · tail-aware loss.",
         size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# =============================================================================
# Slide 16 · Summary — the full logic on one slide
# =============================================================================
s = prs.slides.add_slide(BLANK)
page(s, 16, TOTAL, "总结  —  Introduction 一图记忆",
     "5 个 PART · 9 步推理 · 终点 = ASProNet 的设计动机",
     section="SUMMARY")

# Horizontal flow of 5 blocks
ox, oy = Inches(0.6), Inches(1.85)
bw, bh = Inches(2.34), Inches(4.4)
gx = Inches(0.13)

parts_info = [
    ("PART I",  "安全后果",      NAVY,   [
        "① 事故 ≈ 异常工况",
        "② 4088 死, $47B/yr",
        "结论: ASM 必须做."
    ]),
    ("PART II", "现状极限",      DEEP,   [
        "③ 现行 = FDD",
        "④ FDD = reactive +",
        "    alarm flood",
        "结论: 决策点太晚."
    ]),
    ("PART III","范式跃迁",      GREEN,  [
        "⑤ Prognosis = ",
        "    FDD 的前向延伸",
        "    在 alarm 之前预测",
        "结论: 范式可行."
    ]),
    ("PART IV", "障碍 (失配)",   RED,    [
        "⑥ 数据严重失衡",
        "⑦ Global MSE 偏正常",
        "    structural mismatch",
        "结论: 直接套不行."
    ]),
    ("PART V",  "方法论",        ACCENT, [
        "⑧ 时序 + 工况敏感",
        "⑨ MoE Transformer +",
        "    tail-aware",
        "→ ASProNet"
    ]),
]

for k, (lab, head, col, lines) in enumerate(parts_info):
    x = ox + k * (bw + gx)
    add_round(s, x, oy, bw, bh, WHITE, line=NAVY)
    add_rect(s, x, oy, bw, Inches(0.5), col)
    add_text(s, x, oy + Inches(0.07), bw, Inches(0.4),
             lab, size=13, bold=True,
             color=NAVY if col is ACCENT else WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, oy + Inches(0.6), bw, Inches(0.45),
             head, size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
    block = []
    for ln in lines:
        block.append([(ln, 12, DGREY, False)])
    rich(s, x + Inches(0.15), oy + Inches(1.25), bw - Inches(0.3),
         bh - Inches(1.4), block, line_spacing=1.5)
    if k < len(parts_info) - 1:
        arrow(s, x + bw, oy + bh / 2, x + bw + gx, oy + bh / 2,
              color=ACCENT, w=2.5)

# Bottom strip — final takeaway
add_round(s, Inches(0.5), Inches(6.45), Inches(12.4), Inches(0.55), NAVY)
add_text(s, Inches(0.5), Inches(6.55), Inches(12.4), Inches(0.4),
         "一句话: 因为现状 ‘出事再补救’ 代价太高, 又因为直接套 forecasting 会偏向正常态, 所以需要 ASProNet 这种 ‘对异常敏感’ 的预警模型.",
         size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# -----------------------------------------------------------------------------
out_path = "/home/aicode/sherwin/TSFM/paper/ASProNet_Introduction_逻辑推导.pptx"
prs.save(out_path)
print(f"saved {out_path}  ({len(prs.slides)} slides)")
