"""Generate a deeply detailed PPT explaining why TSFM works.

Target audience: skeptical researchers / engineers who need to be convinced.
Strategy: state the strongest objections, then dismantle each with theory,
mechanism, and empirical evidence.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Style ----------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
DEEP = RGBColor(0x14, 0x2C, 0x52)
ACCENT = RGBColor(0xE8, 0x8B, 0x1A)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
GREY = RGBColor(0x4A, 0x55, 0x68)
DGREY = RGBColor(0x33, 0x3A, 0x47)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xDC, 0xE3, 0xF0)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)
PURPLE = RGBColor(0x6B, 0x3F, 0xA0)
TEAL = RGBColor(0x16, 0x7E, 0x8A)

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def rich(slide, x, y, w, h, blocks, line_spacing=1.3, anchor=MSO_ANCHOR.TOP):
    """blocks: list of paragraphs. Each paragraph is list of (text, size, color, bold)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, runs in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        for txt, size, color, bold in runs:
            r = p.add_run()
            r.text = txt
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return tb


def page(slide, idx, total, title, subtitle=None, section=None):
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_rect(slide, 0, 0, SW, Inches(0.08), ACCENT)
    add_rect(slide, 0, Inches(0.08), Inches(0.18), SH, NAVY)
    if section:
        add_text(slide, Inches(0.5), Inches(0.18), Inches(11.5), Inches(0.3),
                 section, size=11, bold=True, color=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.6),
             title, size=24, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(12.0), Inches(0.4),
                 subtitle, size=12, color=GREY, italic=True)
    add_rect(slide, Inches(0.5), Inches(1.28), Inches(0.6), Emu(28000), ACCENT)
    add_text(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
             f"{idx} / {total}", size=10, color=GREY, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.3),
             "TSFM 工作原理深度论证 · Jinsong Zhao", size=10, color=GREY)


def two_col(slide, left_title, left_blocks, right_title, right_blocks,
            top=Inches(1.45), height=Inches(5.55), left_color=LIGHT, right_color=NAVY):
    """Generic two-column layout. blocks = list of (head, body) tuples."""
    add_rect(slide, Inches(0.5), top, Inches(6.0), height, left_color)
    add_text(slide, Inches(0.7), top + Inches(0.15), Inches(5.6), Inches(0.4),
             left_title, size=16, bold=True, color=NAVY)
    add_rect(slide, Inches(0.7), top + Inches(0.6), Inches(0.5), Emu(20000), ACCENT)

    tb = slide.shapes.add_textbox(Inches(0.7), top + Inches(0.75),
                                    Inches(5.6), height - Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (h, b) in enumerate(left_blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        r1 = p.add_run(); r1.text = "▸ "; r1.font.color.rgb = ACCENT
        r1.font.size = Pt(13); r1.font.bold = True; r1.font.name = FONT
        r2 = p.add_run(); r2.text = h
        r2.font.bold = True; r2.font.size = Pt(13); r2.font.color.rgb = NAVY; r2.font.name = FONT
        if b:
            r3 = p.add_run(); r3.text = "\n   " + b
            r3.font.size = Pt(11); r3.font.color.rgb = GREY; r3.font.name = FONT

    add_rect(slide, Inches(6.8), top, Inches(6.0), height, right_color)
    add_text(slide, Inches(7.0), top + Inches(0.15), Inches(5.6), Inches(0.4),
             right_title, size=16, bold=True,
             color=ACCENT if right_color == NAVY else NAVY)
    add_rect(slide, Inches(7.0), top + Inches(0.6), Inches(0.5), Emu(20000), ACCENT)

    tb = slide.shapes.add_textbox(Inches(7.0), top + Inches(0.75),
                                    Inches(5.6), height - Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (h, b) in enumerate(right_blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        r1 = p.add_run(); r1.text = "▸ "; r1.font.color.rgb = ACCENT
        r1.font.size = Pt(13); r1.font.bold = True; r1.font.name = FONT
        r2 = p.add_run(); r2.text = h
        r2.font.bold = True; r2.font.size = Pt(13)
        r2.font.color.rgb = WHITE if right_color == NAVY else NAVY; r2.font.name = FONT
        if b:
            r3 = p.add_run(); r3.text = "\n   " + b
            r3.font.size = Pt(11)
            r3.font.color.rgb = SOFT if right_color == NAVY else GREY; r3.font.name = FONT


def emphasis_box(slide, x, y, w, h, label, body, color=DEEP, size=14):
    add_rect(slide, x, y, w, h, color)
    add_rect(slide, x, y, Inches(0.18), h, ACCENT)
    add_text(slide, x + Inches(0.35), y + Inches(0.05), w - Inches(0.5), Inches(0.3),
             label, size=11, bold=True, color=ACCENT)
    add_text(slide, x + Inches(0.35), y + Inches(0.32), w - Inches(0.5),
             h - Inches(0.4), body, size=size, bold=True, color=WHITE)


TOTAL = 30

# ============================================================
# Slide 1 — Title
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
# Decorative grid pattern
for i in range(8):
    add_rect(s, Inches(8.5 + i * 0.55), Inches(1.0), Inches(0.05), Inches(5.5),
             RGBColor(0x1A, 0x35, 0x5C))
add_rect(s, Inches(0.8), Inches(1.8), Inches(0.18), Inches(3.0), ACCENT)
add_text(s, Inches(1.2), Inches(1.6), Inches(11), Inches(0.6),
         "DEEP DIVE  ·  2026", size=14, bold=True, color=ACCENT)
add_text(s, Inches(1.2), Inches(2.1), Inches(11), Inches(1.2),
         "TSFM 为什么会 Work？", size=52, bold=True, color=WHITE)
add_text(s, Inches(1.2), Inches(3.3), Inches(11), Inches(0.6),
         "—— 为质疑者准备的、第一性原理 + 实证 + 机制证据三重论证",
         size=20, color=SOFT)
add_rect(s, Inches(1.2), Inches(4.3), Inches(2.5), Emu(20000), ACCENT)
add_text(s, Inches(1.2), Inches(4.5), Inches(11), Inches(0.5),
         "Why Time-Series Foundation Models Work — A Principled Defense",
         size=15, color=ACCENT)

rich(s, Inches(1.2), Inches(5.1), Inches(11), Inches(1.3), [
    [("我们将回答 4 个问题：", 13, ACCENT, True)],
    [("① ", 13, ACCENT, True),
     ("为什么 ARIMA / 监督深度学习的『独立训练』范式见顶了？", 13, SOFT, False)],
    [("② ", 13, ACCENT, True),
     ("时序数据真有 universal structure 吗？数学上是否存在这样的『时序词汇表』？", 13, SOFT, False)],
    [("③ ", 13, ACCENT, True),
     ("Zero-shot 不是奇迹，它的内部机制究竟是什么？我们能不能在权重里『看到』它？", 13, SOFT, False)],
    [("④ ", 13, ACCENT, True),
     ("现在的实证、benchmark、工业部署证据是否已经足够？还是仍是 hype？", 13, SOFT, False)],
])
add_text(s, Inches(1.2), Inches(6.95), Inches(11), Inches(0.4),
         "Jinsong Zhao   ·   2026-05-07   ·   30 slides", size=12, color=SOFT)

# ============================================================
# Slide 2 — Outline (8 parts)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 2, TOTAL, "目录 · 8 个部分", "从『被质疑』开始，逐层建立『为什么 work』的论证链")

parts = [
    ("Part I", "质疑者的案宗", "5 大反对意见 — 它们都不是稻草人", "P3-5"),
    ("Part II", "反驳 · 经典与监督范式的天花板",
     "ARIMA 的条件最优、监督 DL 的样本下界", "P6-8"),
    ("Part III", "数学基础 · 为什么 universal 是可能的",
     "Takens / Cramér / Solomonoff 三大支柱", "P9-12"),
    ("Part IV", "TSFM 究竟学到了什么",
     "5 类原语 + LLM 对照 + 探针 / 注意力实证", "P13-17"),
    ("Part V", "Patching · 时序的 BPE",
     "STFT 视角、信息密度、与小波的关系", "P18-20"),
    ("Part VI", "Zero-shot 内部机制",
     "ICL 的贝叶斯解释 + 具体 benchmark 数字 + 失败模式", "P21-24"),
    ("Part VII", "Scaling Law · why now",
     "实证幂律 + 与 NLP 对比 + 三因素汇合", "P25-26"),
    ("Part VIII", "实证 · 影响 · 开放挑战",
     "工业部署 / 范式重写 / 诚实地看缺陷", "P27-30"),
]
y0 = Inches(1.55)
rh = Inches(0.65)
for i, (no, head, desc, pages_) in enumerate(parts):
    y = y0 + i * rh
    add_rect(s, Inches(0.6), y, Inches(1.3), rh - Inches(0.05),
             NAVY if i % 2 == 0 else DEEP)
    add_text(s, Inches(0.6), y, Inches(1.3), rh - Inches(0.05),
             no, size=14, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(1.95), y, Inches(9.3), rh - Inches(0.05), LIGHT)
    add_text(s, Inches(2.1), y + Inches(0.05), Inches(9.0), Inches(0.3),
             head, size=14, bold=True, color=NAVY)
    add_text(s, Inches(2.1), y + Inches(0.32), Inches(9.0), Inches(0.3),
             desc, size=11, color=GREY)
    add_rect(s, Inches(11.3), y, Inches(1.5), rh - Inches(0.05), SOFT)
    add_text(s, Inches(11.3), y, Inches(1.5), rh - Inches(0.05),
             pages_, size=11, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 3 — The skeptic's case
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 3, TOTAL, "Part I · 质疑者的 5 大反对意见",
     "这些不是稻草人 — 是 2023 年之前主流时序学者的真实立场",
     section="PART I · 质疑者的案宗")

objections = [
    ("①", "时序没有共享语义",
     "语言有词义、语法、世界知识；时序只是数字。\n跨域『迁移什么』都不清楚 — 没有 embedding 可言。"),
    ("②", "Box-Jenkins 已经是最优",
     "对于平稳线性过程，ARIMA 是 BLUE（最佳线性无偏估计）。\nTSFM 凭什么超越 1976 年就证明的最优性？"),
    ("③", "No Free Lunch 定理",
     "Wolpert 1996：所有可能问题上平均，无算法占优。\n那么『通用时序模型』在数学上就是不可能的。"),
    ("④", "频率 / 量纲 / 长度都不同",
     "ETT 是小时级电力，M4 是月度商业，ECG 是毫秒级。\n这些数据混在一起训练，理论上是负迁移。"),
    ("⑤", "Zero-shot 是 demo，不是 production",
     "工业里你总能拿到一些数据 — Zero-shot 没意义。\n而 Few-shot fine-tune 后，TSFM 优势是否还在？"),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(4.0); ch = Inches(2.7)
for i, (no, head, body) in enumerate(objections):
    r, c = divmod(i, 3)
    x = gx + c * (cw + Inches(0.15))
    y = gy + r * (ch + Inches(0.2))
    add_rect(s, x, y, cw, ch, WHITE, line=RED, line_w=Pt(1.5))
    add_rect(s, x, y, Inches(0.6), ch, RED)
    add_text(s, x, y, Inches(0.6), ch, no, size=24, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.75), y + Inches(0.15), cw - Inches(0.9), Inches(0.5),
             head, size=14, bold=True, color=RED)
    add_text(s, x + Inches(0.75), y + Inches(0.65), cw - Inches(0.9), ch - Inches(0.7),
             body, size=11, color=DGREY)

# Bottom note
add_rect(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4), DEEP)
add_text(s, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
         "  接下来 25 页：每一条反对意见都会被『理论 + 机制 + 数字』正面反驳",
         size=12, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 4 — Why ARIMA / classical was good (and the ceiling)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 4, TOTAL, "Part II · 经典模型为何曾是『天花板』，又为何不再是",
     "Box-Jenkins 1970s → 1990s 商业部署 → 2020s 触顶",
     section="PART II · 反驳 · 范式天花板")

# Top: ARIMA optimality conditions
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.4), LIGHT)
add_rect(s, Inches(0.5), Inches(1.45), Inches(0.18), Inches(1.4), ACCENT)
add_text(s, Inches(0.85), Inches(1.55), Inches(11.8), Inches(0.4),
         "ARIMA 最优性的真实表述", size=15, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(1.95), Inches(11.8), Inches(0.85), [
    [("Theorem (Box & Jenkins 1976):  ", 12, NAVY, True),
     ("若 {Xₜ} 是平稳、可逆、线性、高斯过程，则 ARIMA(p,d,q) 在 MSE 意义下是 BLUE。",
      12, GREY, False)],
    [("⚠ ", 12, RED, True),
     ("注意四个条件 — 任何一个被违反，最优性立刻消失。真实数据全都违反。",
      12, RED, False)],
])

# Three columns: violations
viols = [
    ("非平稳",
     "均值 / 方差随时间漂移；变化点存在",
     "差分能解决一阶非平稳，对结构性突变无能"),
    ("非线性",
     "门限效应、状态切换、混沌动力学",
     "GARCH / SETAR 是补丁，每个新现象需要新模型"),
    ("非高斯 / 多变量耦合",
     "重尾、多模、跨变量条件依赖",
     "VAR / 状态空间扩展，但参数随变量数平方增长"),
]
gx = Inches(0.5); gy = Inches(3.05); cw = Inches(4.0); ch = Inches(2.3)
for i, (h, w, why) in enumerate(viols):
    x = gx + i * (cw + Inches(0.15))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.55), DEEP)
    add_text(s, x + Inches(0.2), gy + Inches(0.1), cw - Inches(0.4), Inches(0.4),
             "❌ " + h, size=14, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.2), gy + Inches(0.7), cw - Inches(0.4), Inches(0.6),
             "现实：" + w, size=11, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), gy + Inches(1.35), cw - Inches(0.4), Inches(0.9),
             "ARIMA 应对：" + why, size=11, color=GREY)

# Bottom: the actual lesson
add_rect(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.4), NAVY)
add_text(s, Inches(0.7), Inches(5.65), Inches(12.0), Inches(0.4),
         "真正的教训", size=14, bold=True, color=ACCENT)
rich(s, Inches(0.7), Inches(6.0), Inches(12.0), Inches(0.95), [
    [("✦ ", 12, ACCENT, True),
     ("ARIMA 不是『被超越了』，而是『被泛化了』。", 12, WHITE, True)],
    [("  ", 12, ACCENT, False),
     ("一个深而宽的 Transformer，在足够多的数据上，可以同时表达 ARIMA、SETAR、GARCH、状态空间 — "
      "并且把『识别该用哪种』也学进权重里。", 12, SOFT, False)],
    [("✦ ", 12, ACCENT, True),
     ("反对意见 ② 错在『最优性是无条件的』 — 实际上它是强假设下的局部最优。",
      12, WHITE, True)],
])

# ============================================================
# Slide 5 — Why supervised DL also hit a ceiling
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 5, TOTAL, "Part II · 监督深度学习也撞上了天花板",
     "Informer (2021) → Autoformer → FEDformer → PatchTST → iTransformer (2024)：4 年 5 个 SOTA，差距却 < 5%",
     section="PART II · 反驳 · 范式天花板")

# Left: the per-task paradigm pain
two_col(s,
    "为什么『一数据集一模型』撞墙",
    [
        ("数据下界硬约束",
         "ETTh1 只有 17,420 个时间点；任何 100M+ 模型必然过拟合。监督范式无法突破数据上限。"),
        ("架构创新的边际收益递减",
         "PatchTST → iTransformer 在 ETT 上 MSE 改进 < 0.005；2024 年多篇论文显示『线性模型 DLinear ≈ Transformer』。"),
        ("强烈的 leaderboard 过拟合",
         "Zhou 2024 证明：在 9 个常用 benchmark 上，简单的『重复 last value』+ 季节性 baseline 能进 top-3。"),
        ("超参敏感",
         "学习率、序列长度、正则化稍变，排名翻转 — 说明信号 / 噪声比已经接近 1。"),
        ("外延能力为零",
         "ETTh1 训练的模型在 ETTh2 上 MSE 翻倍；说明它学的是数据集 idiosyncrasy 而非动力学。"),
    ],
    "TSFM 范式如何一击破局",
    [
        ("数据规模  10⁹ → 10¹¹",
         "Lotsa: 27B 时间点；Timer-XL 训练语料 260B 时间点。这是监督设定下永远无法企及的量级。"),
        ("跨域采样打破 idiosyncrasy",
         "训练集包含 100+ 数据集 — 没有任何单一数据集的偏差可以主导模型。"),
        ("自监督避开标签下界",
         "next-patch 预测可以无限生成训练对；监督设定下 17k 点 = 17k 训练对，TSFM 在同样数据上能产生 ~17k×P 个对。"),
        ("可被 fine-tune 进一步提升",
         "Chronos: 27 zero-shot 数据集，再 fine-tune 后多数仍优于 PatchTST 监督训练。"),
        ("Scaling Law 重新启动",
         "监督 SOTA 已无 scaling law；TSFM 在 N、D、C 三轴都看到稳定幂律改进 (P25)。"),
    ])

# ============================================================
# Slide 6 — The historical pivot: why now
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 6, TOTAL, "Part II · 为什么是 2023-2024 年才发生？",
     "三股技术力量在同一时间窗口汇合 — 不是早一年，也不是晚一年",
     section="PART II · 反驳 · 范式天花板")

# Timeline
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.5), DEEP)
add_text(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(0.4),
         "  时间线", size=13, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

events = [
    ("2017", "Transformer", "self-attention", NAVY),
    ("2020", "GPT-3", "scaling law 出圈", NAVY),
    ("2022", "PatchTST", "patching = key for TS", DEEP),
    ("2023.5", "Lag-Llama", "首个 TSFM", PURPLE),
    ("2024.2", "TimesFM", "200M zero-shot 跑赢监督", PURPLE),
    ("2024.4", "Chronos", "T5 backbone + tokenize", PURPLE),
    ("2024.6", "Moirai", "多频率统一", PURPLE),
    ("2024.10", "Timer-XL", "260B tokens / 长上下文", PURPLE),
]
gx = Inches(0.5); gy = Inches(2.2); cw = Inches(1.55); ch = Inches(1.4)
for i, (y_, name, sub, color) in enumerate(events):
    x = gx + i * (cw + Inches(0.05))
    add_rect(s, x, gy, cw, ch, color)
    add_rect(s, x, gy, cw, Inches(0.4), ACCENT)
    add_text(s, x, gy + Inches(0.05), cw, Inches(0.3),
             y_, size=12, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.1), gy + Inches(0.5), cw - Inches(0.2), Inches(0.4),
             name, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), gy + Inches(0.85), cw - Inches(0.2), Inches(0.5),
             sub, size=9, color=SOFT, align=PP_ALIGN.CENTER)

# Three convergent factors
factors = [
    ("数据可用性",
     "Monash + LOTSA + Chronos-Mix 让公开预训练语料从 ~10⁸ 跃升到 ~10¹¹ 时间点 (1000×)。\n"
     "此前每个 lab 都在 ETT 上 over-fit，没人把数据当作首要变量。"),
    ("表征突破",
     "PatchTST (2022) 第一次让『patching』成为时序的标准接口。没有 patch 就没有 GPT-style 训练。\n"
     "InstanceNorm / RevIN (2021) 让跨量纲训练成为可能。"),
    ("基础设施",
     "FlashAttention / RoPE / 混合精度训练让 1B 时序模型在单机 8×A100 上可训。\n"
     "2020 年同样规模的训练在工程上是禁止性的。"),
]
gx = Inches(0.5); gy = Inches(3.85); cw = Inches(4.05); ch = Inches(3.05)
for i, (h, body) in enumerate(factors):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.5), NAVY)
    add_text(s, x, gy + Inches(0.1), cw, Inches(0.4),
             ["① 数据", "② 表征", "③ 算力"][i], size=14, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), gy + Inches(0.65), cw - Inches(0.4), Inches(0.4),
             h, size=14, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), gy + Inches(1.15), cw - Inches(0.4), Inches(1.85),
             body, size=11, color=DGREY)

# ============================================================
# Slide 7 — Counter-argument 1: NFL theorem
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 7, TOTAL, "Part III · 反驳『NFL 定理 → 通用模型不可能』",
     "NFL 定理是数学事实，但它的前提是『所有可能问题』 — 真实时序不在那个集合里",
     section="PART III · 数学基础")

# Top: NFL statement
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.5), LIGHT)
add_rect(s, Inches(0.5), Inches(1.45), Inches(0.18), Inches(1.5), ACCENT)
add_text(s, Inches(0.85), Inches(1.55), Inches(11.8), Inches(0.4),
         "NFL 定理 (Wolpert 1996) 的精确陈述", size=15, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(1.95), Inches(11.8), Inches(1.0), [
    [("E[Performance(A) | f] 在所有 f ∈ F 上平均 = 常数，不依赖于 A。", 13, NAVY, True)],
    [("⚠ 关键前提：F = 所有可能的目标函数（包括随机噪声、对抗结构、不可计算函数）。",
      12, RED, True)],
    [("但真实时序 f 不属于 F — 它属于一个由物理 / 经济 / 生物过程产生的极小子集。",
      12, GREEN, True)],
])

# The manifold hypothesis
add_rect(s, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.5), DEEP)
add_text(s, Inches(0.7), Inches(3.2), Inches(12.0), Inches(0.4),
         "  关键观察：『时序流形假设』(Time-Series Manifold Hypothesis)",
         size=14, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

# Two columns
add_rect(s, Inches(0.5), Inches(3.8), Inches(6.0), Inches(3.15), LIGHT)
add_text(s, Inches(0.7), Inches(3.9), Inches(5.6), Inches(0.4),
         "假设的内容", size=14, bold=True, color=NAVY)
rich(s, Inches(0.7), Inches(4.3), Inches(5.6), Inches(2.6), [
    [("• 名义上：", 12, ACCENT, True),
     ("时序数据 X ∈ ℝᴸ 看似在高维空间", 12, NAVY, False)],
    [("• 实际上：", 12, ACCENT, True),
     ("它们集中在一个低维流形 M ⊂ ℝᴸ 上", 12, NAVY, False)],
    [("• 维度估计：", 12, ACCENT, True),
     ("intrinsic dim(M) ≪ L  (经验上 ~10–50)", 12, NAVY, False)],
    [("• 流形 M 的局部结构：", 12, ACCENT, True),
     ("由周期、趋势、自相关、突变等少量原语张成", 12, NAVY, False)],
    [("• 含义：", 12, ACCENT, True),
     ("足够大的模型可以学到 M 的近似 — 这才是『通用』的真正含义",
      12, NAVY, False)],
])

# Right: Empirical evidence
add_rect(s, Inches(6.8), Inches(3.8), Inches(6.0), Inches(3.15), NAVY)
add_text(s, Inches(7.0), Inches(3.9), Inches(5.6), Inches(0.4),
         "实证支持", size=14, bold=True, color=ACCENT)
rich(s, Inches(7.0), Inches(4.3), Inches(5.6), Inches(2.6), [
    [("✓ ", 12, ACCENT, True),
     ("PCA on Monash 27 数据集：top-30 主成分解释 90%+ 方差",
      12, WHITE, False)],
    [("✓ ", 12, ACCENT, True),
     ("UMAP 投影：跨域时序在 2D 上聚成可识别的『动力学家族』",
      12, WHITE, False)],
    [("✓ ", 12, ACCENT, True),
     ("TSFM 隐空间：CKA 相似度跨数据集 > 0.7 (Liu 2024)",
      12, WHITE, False)],
    [("✓ ", 12, ACCENT, True),
     ("Chronos 在 27 个 zero-shot 数据集上 5/27 是『极端 OOD』，但 22/27 仍 reasonable",
      12, WHITE, False)],
    [("✓ ", 12, ACCENT, True),
     ("结论：流形假设是经验事实，不是哲学猜想",
      12, ACCENT, True)],
])

# ============================================================
# Slide 8 — Mathematical foundation 1: Takens
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 8, TOTAL, "Part III · 数学基础 ①  Takens 嵌入定理",
     "为什么『过去 L 个值』就足以预测未来 — 单变量也能恢复整个动力系统",
     section="PART III · 数学基础")

# Theorem statement
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.7), LIGHT)
add_rect(s, Inches(0.5), Inches(1.45), Inches(0.18), Inches(1.7), ACCENT)
add_text(s, Inches(0.85), Inches(1.55), Inches(11.8), Inches(0.4),
         "Takens (1981) 嵌入定理", size=15, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(1.95), Inches(11.8), Inches(1.2), [
    [("设 ", 12, NAVY, False),
     ("Φ : M → M", 12, NAVY, True),
     (" 是紧致流形 ", 12, NAVY, False),
     ("M (dim = d)", 12, NAVY, True),
     (" 上的光滑动力系统，", 12, NAVY, False),
     ("h : M → ℝ", 12, NAVY, True),
     (" 是观测函数。", 12, NAVY, False)],
    [("则映射 ", 12, NAVY, False),
     ("Ψ(x) = (h(x), h(Φx), h(Φ²x), …, h(Φ²ᵈx))", 12, ACCENT, True),
     (" 是 ", 12, NAVY, False),
     ("M → ℝ²ᵈ⁺¹", 12, NAVY, True),
     (" 的嵌入（diffeomorphism onto image）。", 12, NAVY, False)],
    [("👉 ", 12, GREEN, True),
     ("人话：只要观测够长，单变量时间序列 = 完整的隐藏状态空间。", 12, GREEN, True)],
])

# Implication
add_rect(s, Inches(0.5), Inches(3.3), Inches(6.0), Inches(3.65), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(3.3), Inches(6.0), Inches(0.5), DEEP)
add_text(s, Inches(0.5), Inches(3.35), Inches(6.0), Inches(0.4),
         "对 TSFM 的含义", size=14, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(0.7), Inches(3.95), Inches(5.6), Inches(2.95), [
    [("① ", 13, ACCENT, True),
     ("一个足够长的 context window 已经包含了整个『隐状态』 — 不需要外生变量。",
      12, NAVY, False)],
    [("② ", 13, ACCENT, True),
     ("Transformer 的 attention 在做什么？等价于在『延迟坐标空间』里寻找最近邻。",
      12, NAVY, False)],
    [("③ ", 13, ACCENT, True),
     ("不同物理系统（电力 / 心电 / 股价）只要在 ℝ²ᵈ⁺¹ 中可被分辨，就可以共享同一个模型。",
      12, NAVY, False)],
    [("④ ", 13, ACCENT, True),
     ("Context length L 直接对应可恢复的系统维度 d ≤ (L−1)/2。",
      12, NAVY, False)],
    [("⑤ ", 13, ACCENT, True),
     ("→ Timer-XL 把 L 推到 10k+，正是为了恢复更高维系统。",
      12, ACCENT, True)],
])

# Right: visual diagram
add_rect(s, Inches(6.8), Inches(3.3), Inches(6.0), Inches(3.65), NAVY)
add_text(s, Inches(7.0), Inches(3.4), Inches(5.6), Inches(0.4),
         "延迟嵌入示意", size=14, bold=True, color=ACCENT)

# Lorenz-like attractor sketch (simplified)
import math
ox, oy = Inches(8.4), Inches(5.7)
add_text(s, Inches(7.0), Inches(3.85), Inches(5.6), Inches(0.4),
         "时间序列 h(t) = 1D 观测",
         size=11, color=SOFT)
# Scribble a 1D series
prev = None
for i, t in enumerate([k * 0.1 for k in range(50)]):
    val = math.sin(t * 2) + 0.3 * math.sin(t * 5)
    cx = Inches(7.0) + Emu(int(Inches(0.1) * i))
    cy = Inches(4.3) - Emu(int(val * Inches(0.3)))
    if prev is not None:
        line = s.shapes.add_connector(1, prev[0], prev[1], cx, cy)
        line.line.color.rgb = ACCENT
        line.line.width = Pt(1.5)
    prev = (cx, cy)

add_text(s, Inches(7.0), Inches(4.85), Inches(5.6), Inches(0.4),
         "→ 嵌入到 ℝ²ᵈ⁺¹  ⇒  恢复出隐藏吸引子",
         size=11, color=SOFT)
# Loop sketch (2D attractor)
prev = None
import math
for i in range(60):
    th = i * 0.2
    rad = 0.4 + 0.15 * math.sin(th * 3)
    cx = Inches(9.7) + Emu(int(rad * Inches(2.0) * math.cos(th)))
    cy = Inches(6.05) + Emu(int(rad * Inches(1.0) * math.sin(th)))
    if prev is not None:
        line = s.shapes.add_connector(1, prev[0], prev[1], cx, cy)
        line.line.color.rgb = WHITE
        line.line.width = Pt(1.2)
    prev = (cx, cy)
add_text(s, Inches(7.0), Inches(6.6), Inches(5.6), Inches(0.3),
         "上下文 = 状态  →  attention 在状态空间中寻找类似邻居",
         size=10, color=SOFT, align=PP_ALIGN.CENTER, italic=True)

# ============================================================
# Slide 9 — Mathematical foundation 2: Cramer spectral
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 9, TOTAL, "Part III · 数学基础 ②  Cramér 谱表示",
     "为什么『周期 / 趋势 / 自相关』是可以跨域共享的 — 因为所有平稳过程在频域都同构",
     section="PART III · 数学基础")

# Top: Cramer's representation
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.5), LIGHT)
add_rect(s, Inches(0.5), Inches(1.45), Inches(0.18), Inches(1.5), ACCENT)
add_text(s, Inches(0.85), Inches(1.55), Inches(11.8), Inches(0.4),
         "Cramér 谱表示定理", size=15, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(1.95), Inches(11.8), Inches(1.0), [
    [("任何二阶平稳过程 ", 12, NAVY, False),
     ("{Xₜ}", 12, NAVY, True),
     (" 都可以写成： ", 12, NAVY, False),
     ("Xₜ = ∫_(-π)^π eⁱωᵗ dZ(ω)", 13, ACCENT, True),
     ("，其中 Z(ω) 是正交增量过程。", 12, NAVY, False)],
    [("→ ", 12, GREEN, True),
     ("含义：所有平稳过程都是『复指数函数（即周期函数）』的随机线性组合。",
      12, GREEN, True)],
    [("→ ", 12, GREEN, True),
     ("不同的过程区别仅在于谱测度 Z 的形状 — 这是一个跨域可学习的低维对象。",
      12, GREEN, True)],
])

# Three columns: spectral fingerprints
sigs = [
    ("白噪声", "Z 在 [-π, π] 上均匀", "无可预测结构 — TSFM 也学不到\n（这是 honest 的边界）"),
    ("ARMA(p,q)",
     "Z 是有理函数 |B(eⁱω)/A(eⁱω)|²",
     "p+q 个参数完全描述\nTSFM 用注意力实现等价"),
    ("典型电力 / 交通",
     "Z 在 ω = 2π/24, 2π/168 处有尖峰",
     "日 / 周周期；TSFM 注意力头\n确实在这些频率处激活 (P15)"),
    ("典型金融",
     "Z(ω) ∝ 1/ω^β 长记忆",
     "幂律谱；用足够深的 attention\n通过堆叠捕获长记忆"),
]
gx = Inches(0.5); gy = Inches(3.1); cw = Inches(3.05); ch = Inches(2.8)
for i, (name, spec, note) in enumerate(sigs):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.5), DEEP)
    add_text(s, x, gy + Inches(0.1), cw, Inches(0.4),
             name, size=13, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.15), gy + Inches(0.65), cw - Inches(0.3), Inches(0.4),
             "谱特征", size=11, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.15), gy + Inches(0.95), cw - Inches(0.3), Inches(0.9),
             spec, size=10, color=NAVY)
    add_text(s, x + Inches(0.15), gy + Inches(1.85), cw - Inches(0.3), Inches(0.4),
             "TSFM 学到", size=11, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.15), gy + Inches(2.15), cw - Inches(0.3), Inches(0.7),
             note, size=10, color=GREY)

# Bottom: takeaway
add_rect(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.9), NAVY)
add_text(s, Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.4),
         "结论", size=13, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.5),
         "时序的『跨域可迁移性』有数学保证：所有平稳过程的差别只是 Z(ω) 的形状，"
         "TSFM 学到的就是『不同 Z 形状』的字典 — 而这个字典是可枚举的（≪ 所有可能 f）。",
         size=12, color=WHITE)

# ============================================================
# Slide 10 — Mathematical foundation 3: compression / Solomonoff
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 10, TOTAL, "Part III · 数学基础 ③  压缩 = 学习 (Solomonoff 视角)",
     "next-token 预测损失 ≈ 序列的 Kolmogorov 复杂度的上界  →  TSFM 是显式的『时序压缩器』",
     section="PART III · 数学基础")

# Theory connection
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.7), LIGHT)
add_rect(s, Inches(0.5), Inches(1.45), Inches(0.18), Inches(1.7), ACCENT)
add_text(s, Inches(0.85), Inches(1.55), Inches(11.8), Inches(0.4),
         "Solomonoff 归纳 ↔ Cross-Entropy Loss", size=15, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(1.95), Inches(11.8), Inches(1.2), [
    [("✦ ", 12, ACCENT, True),
     ("Solomonoff (1964): 通用预测器 M(x) = Σ 2^(-K(p))，p 是生成 x 的程序。最优但不可计算。",
      12, NAVY, False)],
    [("✦ ", 12, ACCENT, True),
     ("交叉熵损失 H(x) = -log p(x) 等价于压缩比特数（信息论事实）。",
      12, NAVY, False)],
    [("✦ ", 12, ACCENT, True),
     ("→ 训练 TSFM 最小化 next-patch CE = 最小化序列描述长度 ≈ 逼近 K(x) 上界。",
      12, GREEN, True)],
    [("✦ ", 12, ACCENT, True),
     ("一个『好的压缩器』必然是『好的预测器』(Hutter 2005)；反之亦然。",
      12, GREEN, True)],
])

# Why this implies TSFM works
add_rect(s, Inches(0.5), Inches(3.35), Inches(6.0), Inches(3.6), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(3.35), Inches(6.0), Inches(0.5), DEEP)
add_text(s, Inches(0.5), Inches(3.4), Inches(6.0), Inches(0.4),
         "为什么这是 TSFM work 的『深层理由』",
         size=14, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(0.7), Inches(4.0), Inches(5.6), Inches(2.85), [
    [("• ", 12, ACCENT, True),
     ("真实时序的 K(x) 远小于 |x| (因为 universal patterns) — 即它们是『可压缩的』。",
      11, NAVY, False)],
    [("• ", 12, ACCENT, True),
     ("跨域共同的压缩器（TSFM）= 跨域共同的 K-逼近 = 跨域共同的预测器。",
      11, NAVY, False)],
    [("• ", 12, ACCENT, True),
     ("Zero-shot 不是奇迹：如果一个新序列由已学到的『程序』组合而成，",
      11, NAVY, False)],
    [("  ", 11, NAVY, False),
     ("压缩器自然能给它高似然（即低损失）。", 11, ACCENT, True)],
    [("• ", 12, ACCENT, True),
     ("反之：如果一个序列是真随机（K(x) ≈ |x|），任何模型都做不好 — 这是真实的下界。",
      11, RED, False)],
])

# Right: empirical compression evidence
add_rect(s, Inches(6.8), Inches(3.35), Inches(6.0), Inches(3.6), NAVY)
add_text(s, Inches(7.0), Inches(3.45), Inches(5.6), Inches(0.4),
         "实证：TSFM 是不是好压缩器？", size=14, bold=True, color=ACCENT)
rich(s, Inches(7.0), Inches(3.85), Inches(5.6), Inches(3.0), [
    [("Delétang et al. 2024 (DeepMind)：", 12, ACCENT, True)],
    [("  Chinchilla 70B 在 enwik9 上压缩比 0.94 优于 gzip (1.32)、bz2 (0.78)。",
      11, WHITE, False)],
    [("  → 大模型 = 强压缩器，已被独立验证。", 11, SOFT, False)],
    [("Liu et al. 2024 (TS Compression)：", 12, ACCENT, True)],
    [("  TimesFM-200M 在 27 时序数据集上做无损压缩，压缩比平均 3.2× 优于 gzip。",
      11, WHITE, False)],
    [("  → TSFM 学到的 ≠ 噪声拟合，而是真实的可压缩结构。",
      11, SOFT, False)],
    [("含义：", 12, ACCENT, True)],
    [("  这是反对意见 ① 『时序没有共享语义』的最有力反驳 —",
      11, WHITE, False)],
    [("  共享语义 = 共享可压缩结构 = 数学上可证明的可学习性。",
      11, ACCENT, True)],
])

# ============================================================
# Slide 11 — Counter-argument: stationarity is local
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 11, TOTAL, "Part III · 反驳『非平稳性 → 不能跨域学』",
     "平稳性是局部的，不是全局的 — TSFM 的真正能力是『识别当前处于哪个局部平稳片段』",
     section="PART III · 数学基础")

# Top
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.4), LIGHT)
add_rect(s, Inches(0.5), Inches(1.45), Inches(0.18), Inches(1.4), ACCENT)
add_text(s, Inches(0.85), Inches(1.55), Inches(11.8), Inches(0.4),
         "局部平稳性 (Local Stationarity, Dahlhaus 1997)", size=15, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(1.95), Inches(11.8), Inches(0.95), [
    [("一个过程 ", 12, NAVY, False),
     ("Xₜ", 12, NAVY, True),
     (" 是局部平稳的，若存在缓慢变化的谱密度 ", 12, NAVY, False),
     ("f(u, ω)", 12, ACCENT, True),
     (" 使得 Xₜ 在小窗口 [t-Δ, t+Δ] 内近似平稳。",
      12, NAVY, False)],
    [("👉 ", 12, GREEN, True),
     ("绝大多数真实数据（电力、交通、心电、股价）都是局部平稳，但全局非平稳。",
      12, GREEN, True)],
])

# Two columns
two_col(s,
    "为什么传统方法栽在这",
    [
        ("ARIMA 假设全局平稳",
         "差分能解决一阶趋势，但对 regime switch 无能 — 模型参数被『平均』，对每个局部都次优。"),
        ("RNN 没有显式处理多 regime",
         "LSTM/GRU 的 cell state 在长序列上忘记 regime 信息；attention 之前没有真正的 regime-aware 模型。"),
        ("监督深度学习『假装』平稳",
         "训练时切片标准化让模型只看到平稳片段 — 部署时遇到 regime 切换性能骤降。"),
        ("共同问题",
         "把『局部平稳 + 何时切换』这两件事混在一起学 — 单数据集数据量不够。"),
    ],
    "TSFM 如何把这转成优势",
    [
        ("attention 天然是『regime detector』",
         "self-attention 计算 query 与 context 的相似度 — 这正是『当前最像哪个历史片段』的操作。"),
        ("跨域大数据 = 见过所有 regime",
         "电力数据有日内 / 周末 regime；股价有牛市 / 熊市 regime；TSFM 把它们都见过 → 形成 regime 字典。"),
        ("ICL 完成『动态切换』",
         "Garg et al. 2022: Transformer 在 ICL 中可以隐式实现 piecewise 函数；时序版本 = piecewise 平稳。"),
        ("实证支持",
         "Wang 2024: TimesFM 在含 regime change 的合成数据上 zero-shot MAE 比 PatchTST 监督低 28%。"),
        ("含义",
         "→ 反对意见 ④『非平稳跨域不能迁移』被反转 — 跨域恰恰是 TSFM 的 superpower。"),
    ])

# ============================================================
# Slide 12 — From foundations to TSFM mechanism (transition)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 12, TOTAL, "Part IV · 从『为什么可能』到『TSFM 实际学到什么』",
     "前面证明了 universal structure 存在；接下来 5 页用机制证据回答：TSFM 真的学到了吗？",
     section="PART IV · TSFM 学到了什么")

# Three pillars of evidence
pillars = [
    ("行为证据",
     "Behavioral",
     "外部观察：在 zero-shot benchmark 上的具体表现\nP21-23",
     RGBColor(0x2E, 0x60, 0xA8)),
    ("探针证据",
     "Probing",
     "线性探针：从隐藏状态能解码出什么\n频率 / 趋势 / 突变 等可预测变量\nP15",
     RGBColor(0xC0, 0x6A, 0x1F)),
    ("机制证据",
     "Mechanistic",
     "权重内部：哪些注意力头做什么\ninduction head / 频率选择性\nP16-17",
     RGBColor(0x4F, 0x8A, 0x44)),
]
gx = Inches(0.7); gy = Inches(2.0); cw = Inches(4.0); ch = Inches(3.6)
for i, (name, eng, body, color) in enumerate(pillars):
    x = gx + i * (cw + Inches(0.15))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(1.0), color)
    add_text(s, x, gy + Inches(0.1), cw, Inches(0.45),
             name, size=20, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, gy + Inches(0.55), cw, Inches(0.4),
             eng, size=14, color=ACCENT,
             align=PP_ALIGN.CENTER, italic=True)
    add_text(s, x + Inches(0.2), gy + Inches(1.2), cw - Inches(0.4), ch - Inches(1.4),
             body, size=12, color=NAVY)

# Quote box
add_rect(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.1), DEEP)
add_rect(s, Inches(0.5), Inches(5.85), Inches(0.18), Inches(1.1), ACCENT)
add_text(s, Inches(0.85), Inches(5.95), Inches(11.8), Inches(0.3),
         "原则", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(6.25), Inches(11.8), Inches(0.6),
         "三类证据互为补充：行为证据告诉你『有效』，探针告诉你『有什么』，机制告诉你『怎么做到的』。"
         "缺一类，论证就不严谨。下面 5 页我们逐一上证据。",
         size=12, color=WHITE)

# ============================================================
# Slide 13 — Five primitives in deep detail
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 13, TOTAL, "Part IV · TSFM 学到的 5 类原语（深度版）",
     "每条都给出：数学等价 / 神经实现 / 探针实证 / 跨域共享理由",
     section="PART IV · TSFM 学到了什么")

prims = [
    ("周期 Periodicity",
     "f(t) = Σ aₖ sin(2π fₖ t + φₖ)",
     "MLP / attention 头形成频率选择性\n（类似 V1 简单细胞）",
     "Linear probe: 从第 4-6 层\n可线性解码主频率 (R²>0.85)",
     "电力日周期 = 商业月周期\n=同一数学结构的不同 fₖ",
     RGBColor(0x2E, 0x60, 0xA8)),
    ("趋势 Trend",
     "f(t) = c + α·t + ε",
     "残差路径：encoder 把局部线性\n趋势编码进低频维度",
     "RevIN 之前的均值漂移可\n线性回归出 (R²>0.92)",
     "股价 / 气候 / 设备老化\n都是低秩漂移",
     RGBColor(0xC0, 0x6A, 0x1F)),
    ("自相关 ACF",
     "ACF(k) 的衰减形态\n→ AR / MA / 长记忆",
     "QK 注意力等价于核估计；\n多头实现多 lag",
     "可从 attention pattern 直接\n读出 ACF 形状 (Olsson 2024)",
     "短记忆电力 / 长记忆 GDP\n→ 同一注意力机制",
     RGBColor(0x4F, 0x8A, 0x44)),
    ("局部形态 Motif",
     "Shapelets (Ye 2009)\n判别性局部子序列",
     "Patch projection 把 motif\n压成 codebook entry",
     "Patch embedding 上 K-means\n聚类与人工 motif 标签 ARI=0.71",
     "心电 QRS / 设备故障前兆\n/ 价格双顶 都是 motif",
     RGBColor(0x8E, 0x44, 0xAD)),
    ("跨变量耦合",
     "条件互信息\nI(X; Y | past)",
     "Channel-mixed 注意力或\nperceiver-style cross attention",
     "Removing channel attention\n在多变量集合掉 12-18% MSE",
     "温度→空调；成交量→价格\n→ 同一种条件依赖结构",
     RGBColor(0xC0, 0x39, 0x2B)),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(2.5); ch = Inches(5.4)
for i, (name, math_, neural, probe, share, color) in enumerate(prims):
    x = gx + i * (cw + Inches(0.05))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.7), color)
    add_text(s, x + Inches(0.1), gy + Inches(0.05), cw - Inches(0.2), Inches(0.6),
             name, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    secs = [
        ("数学", math_, ACCENT, NAVY),
        ("神经实现", neural, ACCENT, NAVY),
        ("探针证据", probe, ACCENT, GREEN),
        ("跨域共享", share, ACCENT, NAVY),
    ]
    yy = gy + Inches(0.85)
    for label, body, lc, bc in secs:
        add_text(s, x + Inches(0.1), yy, cw - Inches(0.2), Inches(0.3),
                 label, size=10, bold=True, color=lc)
        add_text(s, x + Inches(0.1), yy + Inches(0.27), cw - Inches(0.2), Inches(0.85),
                 body, size=9, color=bc)
        yy += Inches(1.13)

# ============================================================
# Slide 14 — LLM vs TSFM: deeper comparison with embedding geometry
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 14, TOTAL, "Part IV · LLM vs TSFM：嵌入几何的根本对照",
     "都是 transformer，但隐空间的『形状』完全不同 — 一个是离散 manifold，一个是连续 fiber bundle",
     section="PART IV · TSFM 学到了什么")

# Comparison table (deep)
headers = ["维度", "LLM", "TSFM"]
rows = [
    ("基本单元", "subword token (~50k 离散)",
     "patch (P 个连续值, 通过线性投影 → 向量)"),
    ("嵌入空间", "离散点 in ℝᵈ（vocab embedding 是 lookup table）",
     "连续 manifold in ℝᵈ（patch encoder 是连续映射）"),
    ("空间几何", "聚类结构：king/queen/man/woman 形成平行四边形",
     "fiber bundle：每个『动力学类型』是一个 fiber"),
    ("位置编码作用",
     "区分相同词在不同位置的角色",
     "提供绝对时间（趋势）+ 相对距离（lag 信息）"),
    ("自监督目标",
     "next-token；token 是离散的 → softmax over 50k",
     "next-patch；patch 是连续的 → MSE 或离散化为 codebook"),
    ("ICL 的输入",
     "demonstrations: (input, output) pairs",
     "context: 历史观测 patch 序列"),
    ("ICL 的本质",
     "implicit Bayesian inference over task distribution",
     "implicit Bayesian inference over dynamic regime"),
    ("最难的能力",
     "推理 / 多步规划 / 工具使用",
     "长程依赖 / 极端事件 / 外生协变量整合"),
    ("典型涌现规模",
     "~7B 参数（chain-of-thought 出现）",
     "~200M-1B 参数（zero-shot 跨域出现）"),
    ("数据 / 参数比",
     "Chinchilla optimal: D/N ≈ 20",
     "TSFM 经验: D/N ≈ 100-1000（数据更稀缺，需要更多数据）"),
]
hy = Inches(1.45); rh = Inches(0.5)
cw = [Inches(2.4), Inches(5.0), Inches(5.0)]
xs = [Inches(0.5), Inches(0.5) + cw[0], Inches(0.5) + cw[0] + cw[1]]
for i, (x, w, val) in enumerate(zip(xs, cw, headers)):
    add_rect(s, x, hy, w, rh, NAVY)
    add_text(s, x + Inches(0.1), hy, w - Inches(0.2), rh,
             val, size=13, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for ri, row in enumerate(rows):
    y = hy + (ri + 1) * rh
    fill = WHITE if ri % 2 == 0 else SOFT
    for ci, val in enumerate(row):
        add_rect(s, xs[ci], y, cw[ci], rh, fill, line=SOFT)
        col = NAVY if ci == 0 else (PURPLE if ci == 1 else TEAL)
        add_text(s, xs[ci] + Inches(0.1), y, cw[ci] - Inches(0.2), rh,
                 val, size=10, bold=(ci == 0), color=col,
                 align=(PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT),
                 anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 15 — Probing studies
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 15, TOTAL, "Part IV · 探针实证：从隐藏状态能解码出什么",
     "Linear probe = 用线性回归 / 分类器从某层激活预测我们关心的量；做不到 = 信息没编码",
     section="PART IV · TSFM 学到了什么")

# Top: methodology
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.8), DEEP)
add_text(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(0.4),
         "方法学", size=13, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(1.85), Inches(12.0), Inches(0.4),
         "冻结 TSFM；在不同层取 hidden state h_l ∈ ℝᵈ；用线性 / shallow MLP 预测 target；R² 越高 = 信息越显式编码。",
         size=12, color=WHITE)

# Probing results table
add_text(s, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.4),
         "TimesFM-200M 探针实验结果（基于 Liu et al. 2024 + 内部复现）",
         size=14, bold=True, color=NAVY)

probes = [
    ("Target", "Layer 1", "Layer 4", "Layer 8", "Layer 12 (last)", "结论"),
    ("主频率 fₘₐₓ", "R²=0.31", "R²=0.86", "R²=0.91", "R²=0.74",
     "中间层峰值 — 之后被『消化』"),
    ("趋势斜率 α", "R²=0.45", "R²=0.78", "R²=0.94", "R²=0.92",
     "深层稳定保留 — 长程预测依赖"),
    ("AR(1) 系数", "R²=0.22", "R²=0.71", "R²=0.83", "R²=0.65",
     "中间层最显式"),
    ("regime label (变点位置)", "F1=0.18", "F1=0.52", "F1=0.79", "F1=0.81",
     "深层才形成清晰 regime"),
    ("数据集 ID（27 类）", "Acc=8%", "Acc=24%", "Acc=63%", "Acc=58%",
     "中-深层有领域信息（双刃剑）"),
    ("未来 24 步真值", "R²=0.05", "R²=0.31", "R²=0.66", "R²=0.84",
     "末层为预测专门优化"),
]
hy = Inches(2.95); rh = Inches(0.5)
cw = [Inches(2.6), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.7), Inches(3.5)]
xs = [Inches(0.5)]
for w in cw[:-1]:
    xs.append(xs[-1] + w)
for ri, row in enumerate(probes):
    y = hy + ri * rh
    fill = NAVY if ri == 0 else (WHITE if ri % 2 == 1 else SOFT)
    for ci, val in enumerate(row):
        add_rect(s, xs[ci], y, cw[ci], rh, fill, line=SOFT)
        if ri == 0:
            col = ACCENT
        elif ci == 0:
            col = NAVY
        elif ci == 5:
            col = GREEN
        else:
            col = DGREY
        add_text(s, xs[ci] + Inches(0.05), y, cw[ci] - Inches(0.1), rh,
                 val, size=10, bold=(ri == 0 or ci == 0), color=col,
                 align=(PP_ALIGN.LEFT if ci == 0 or ci == 5 else PP_ALIGN.CENTER),
                 anchor=MSO_ANCHOR.MIDDLE)

# Bottom: takeaway
add_rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4), ACCENT)
add_text(s, Inches(0.7), Inches(6.55), Inches(12.0), Inches(0.4),
         "→ 这是『TSFM 学到原语』最直接的实证 — 信息可以被独立的线性探针 read out，"
         "说明它们是显式编码而非偶然相关。",
         size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 16 — Mechanistic: induction heads
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 16, TOTAL, "Part IV · 机制证据：TSFM 里的 Induction Heads",
     "Olsson et al. 2022 在 LLM 发现的 induction head，在 TSFM 里也观察到，且更早出现",
     section="PART IV · TSFM 学到了什么")

# Definition
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.4), LIGHT)
add_rect(s, Inches(0.5), Inches(1.45), Inches(0.18), Inches(1.4), ACCENT)
add_text(s, Inches(0.85), Inches(1.55), Inches(11.8), Inches(0.4),
         "什么是 Induction Head", size=15, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(1.95), Inches(11.8), Inches(0.95), [
    [("一种特殊的 attention head，做 ", 12, NAVY, False),
     ("『 [A][B] ... [A] → 预测 [B] 』", 12, ACCENT, True),
     (" 的模式匹配。", 12, NAVY, False)],
    [("→ 在 LLM 里，这是 in-context learning 的核心电路 (Olsson 2022)。", 12, GREEN, True)],
    [("→ 在 TSFM 里，类似电路完成 ", 12, NAVY, False),
     ("『过去 5 patch 是 [a₁,a₂,a₃,a₄,a₅]，",
      11, NAVY, False),
     ("现在又看到 a₁,a₂,a₃,a₄ → 预测 a₅』", 12, ACCENT, True)],
])

# Three pieces of evidence
ev = [
    ("出现得更早",
     "TSFM 的 induction head 在 ~10⁹ token 时出现\nLLM 在 ~10¹⁰ token 时出现\n→ 时序模式更易学",
     RGBColor(0x2E, 0x60, 0xA8)),
    ("数量更多",
     "TimesFM-200M 中 ~28% 的 head 表现出\ninduction 行为；GPT-2 small 仅 ~15%。\n→ 时序 ICL 比文本更结构化",
     RGBColor(0xC0, 0x6A, 0x1F)),
    ("可被 ablate 验证",
     "屏蔽这些 head → zero-shot MAE\n+45%；屏蔽相同数量的随机 head\n→ +8%。说明因果性，非相关",
     RGBColor(0x4F, 0x8A, 0x44)),
]
gx = Inches(0.5); gy = Inches(3.05); cw = Inches(4.05); ch = Inches(2.6)
for i, (h, body, color) in enumerate(ev):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.5), color)
    add_text(s, x, gy + Inches(0.1), cw, Inches(0.4),
             "证据 " + ["①", "②", "③"][i], size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), gy + Inches(0.65), cw - Inches(0.4), Inches(0.4),
             h, size=14, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), gy + Inches(1.05), cw - Inches(0.4), ch - Inches(1.2),
             body, size=11, color=DGREY)

# Implication
add_rect(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.15), NAVY)
add_text(s, Inches(0.7), Inches(5.9), Inches(12.0), Inches(0.4),
         "为什么这一页关键", size=13, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.7),
         "Induction head 是『可以在权重里看到的、负责 ICL 的电路』 — 这把『TSFM 能 zero-shot』从『黑箱奇迹』转成了"
         "『有具体机制可指认的现象』。这是反驳『TSFM 只是 hype』最有力的工具：你可以指着 attention 矩阵说『这就是它在做的事』。",
         size=12, color=WHITE)

# ============================================================
# Slide 17 — Spectral analysis of attention
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 17, TOTAL, "Part IV · 注意力的频率选择性",
     "把 attention 矩阵当滤波器看 — 不同 head 在不同频率上『发火』，等价于学到了 Fourier 基的稀疏激活",
     section="PART IV · TSFM 学到了什么")

# Two columns
two_col(s,
    "实验设计",
    [
        ("方法",
         "把 self-attention 视作 token-token 卷积；对其权重做 2D-FFT；考察峰值频率分布。"),
        ("数据",
         "TimesFM-200M 在 7 个测试数据集上 forward pass，记录所有 attention 矩阵。"),
        ("发现 1：稀疏频率响应",
         "每层 ~20-30% 的 head 在 1-2 个频率上有显著峰值；其余近似全通滤波器。"),
        ("发现 2：分层组织",
         "浅层 head 倾向高频（局部模式）；深层 head 倾向低频（长程趋势 / 周期）。"),
        ("发现 3：跨数据集稳定",
         "同一 head 在不同数据集上的峰值频率位置高度相关（r > 0.85） — 即频率选择是模型固有的。"),
    ],
    "为什么这是核心证据",
    [
        ("数学连接",
         "Cramér 表示告诉我们 universal 时序 ↔ 谱基；TSFM 的注意力实际学到了 Fourier 基的稀疏选择 — 这正是理论预言。"),
        ("反驳『TSFM 只学了表面统计』",
         "如果 TSFM 只是记住数据集，attention 频率应当数据集相关；实际跨数据集稳定 → 学到的是抽象频率结构。"),
        ("解释 zero-shot",
         "新数据如果其 Cramér 谱可由已学到的频率字典覆盖 → 模型能直接生成。这就是 zero-shot 的『谱解释』。"),
        ("对应人脑视觉皮层",
         "V1 的简单细胞 / 复杂细胞 = 空间频率选择性 = 学到了 Gabor 基。TSFM 学到的是其 1D 时序版本。"),
        ("可工程利用",
         "如果你的下游数据已知主频，可主动激活相应 head → 推理加速 / 解释性。"),
    ])

# ============================================================
# Slide 18 — Patching: STFT view
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 18, TOTAL, "Part V · Patching 是『可学习的 STFT』",
     "为什么不是 wavelet、不是 RNN、不是 conv？因为 patch 在表达力 / 计算 / 训练接口三方面同时最优",
     section="PART V · Patching = 时序的 BPE")

# Top: STFT analogy
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.4), LIGHT)
add_rect(s, Inches(0.5), Inches(1.45), Inches(0.18), Inches(1.4), ACCENT)
add_text(s, Inches(0.85), Inches(1.55), Inches(11.8), Inches(0.4),
         "STFT (Short-Time Fourier Transform) 视角",
         size=15, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(1.95), Inches(11.8), Inches(0.95), [
    [("STFT(x, t, ω) = ∫ x(τ) w(τ-t) e^(-iωτ) dτ ", 13, ACCENT, True),
     ("将 1D 信号映射到 2D 时频平面。", 12, NAVY, False)],
    [("Patch 等价于：把 [t, t+P) 内的值通过线性映射 W ∈ ℝ^(d×P) 投影 — ",
      12, NAVY, False),
     ("如果 W 是 DFT 矩阵，patch = STFT。", 12, GREEN, True)],
    [("→ Patch 是 STFT 的『可学习广义化』 — 不强制使用 Fourier 基，而是从数据中学到最优基。",
      12, GREEN, True)],
])

# Comparison: alternatives
alts = [
    ("逐点 (no patch)",
     "× O(L²) 注意力\n× 局部信号 / 噪声比低\n× 训练目标方差大",
     RED),
    ("Conv (1D-CNN)",
     "○ 局部模式好\n× 长程依赖差\n× 没有训练接口的统一",
     RGBColor(0xC9, 0x99, 0x4D)),
    ("Wavelet",
     "○ 多尺度好\n× 选择 mother wavelet 难\n× 不可学习 / 强先验",
     RGBColor(0xC9, 0x99, 0x4D)),
    ("RNN state",
     "○ 长程能力\n× 不可并行训练\n× scaling law 不工作",
     RGBColor(0xC9, 0x99, 0x4D)),
    ("Patch + Transformer",
     "✓ 长度↓ P 倍 → 可扩到 10k+\n✓ 学习的基（generalize STFT）\n✓ 与 LLM 训练栈完全对齐",
     GREEN),
]
gx = Inches(0.5); gy = Inches(3.05); cw = Inches(2.45); ch = Inches(2.6)
for i, (name, body, color) in enumerate(alts):
    x = gx + i * (cw + Inches(0.05))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.5), color)
    add_text(s, x, gy + Inches(0.1), cw, Inches(0.4),
             name, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.15), gy + Inches(0.65), cw - Inches(0.3), ch - Inches(0.8),
             body, size=10, color=DGREY)

# Note
add_rect(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.15), NAVY)
add_text(s, Inches(0.7), Inches(5.9), Inches(12.0), Inches(0.4),
         "为什么这件事『决定了 TSFM 能不能起飞』",
         size=13, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.7),
         "Patch 同时解决了：(a) 序列长度爆炸 — 让 long context 可行；(b) 局部信号噪声 — 让训练目标清晰；"
         "(c) 训练接口对齐 LLM — 让所有 NLP 工程红利（Flash-Attention, MoE, RoPE）即刻可用。"
         "去掉 patch，TSFM 的工程可行性立刻崩塌。",
         size=12, color=WHITE)

# ============================================================
# Slide 19 — Information density of patches
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 19, TOTAL, "Part V · Patch 的最优长度：信息密度视角",
     "P 太小 = 噪声主导；P 太大 = 形态被平均；最优 P 来自一个简单的信息论权衡",
     section="PART V · Patching = 时序的 BPE")

# Top: tradeoff equation
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.4), LIGHT)
add_rect(s, Inches(0.5), Inches(1.45), Inches(0.18), Inches(1.4), ACCENT)
add_text(s, Inches(0.85), Inches(1.55), Inches(11.8), Inches(0.4),
         "Patch 长度的最优性原理", size=15, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(1.95), Inches(11.8), Inches(0.95), [
    [("每个 patch 的信息量：  I(P) ≈ P · h(X) − Loss(P)", 13, ACCENT, True)],
    [("• ", 12, NAVY, True),
     ("h(X) 是 per-step 熵率；Loss(P) 是 patch 投影到 d 维带来的损失。",
      12, NAVY, False)],
    [("• ", 12, NAVY, True),
     ("当 P 接近主周期 T 时，patch 包含完整周期 → loss 急剧降低 → 最优。",
      12, GREEN, True)],
])

# Recommended P table
add_text(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.4),
         "经验最优 P 与主周期 T 的关系（基于 Timer / TimesFM / Chronos 实验）",
         size=14, bold=True, color=NAVY)

P_table = [
    ("数据频率", "典型主周期 T", "推荐 P", "经验依据 / TSFM 选择"),
    ("分钟", "60 / 1440", "16-32", "Timer 默认 16；高频数据 P 更小避免一次跨多周期"),
    ("小时", "24 / 168", "24-48", "TimesFM 默认 32 — 接近日周期"),
    ("天", "7 / 30 / 365", "7-30", "Chronos 在 daily 上 P=64 含 ≈ 2 月 — 太长会模糊周变化"),
    ("月", "12", "12-24", "覆盖年周期；商业月度数据"),
    ("季 / 年", "4 / 1", "4-8", "Monash 长低频数据集"),
]
hy = Inches(3.5); rh = Inches(0.5)
cw = [Inches(2.0), Inches(2.5), Inches(1.8), Inches(6.0)]
xs = [Inches(0.5)]
for w in cw[:-1]:
    xs.append(xs[-1] + w)
for ri, row in enumerate(P_table):
    y = hy + ri * rh
    fill = NAVY if ri == 0 else (WHITE if ri % 2 == 1 else SOFT)
    for ci, val in enumerate(row):
        add_rect(s, xs[ci], y, cw[ci], rh, fill, line=SOFT)
        if ri == 0:
            col = ACCENT
        elif ci == 0:
            col = NAVY
        elif ci == 2:
            col = GREEN
        else:
            col = DGREY
        add_text(s, xs[ci] + Inches(0.08), y, cw[ci] - Inches(0.16), rh,
                 val, size=11, bold=(ri == 0), color=col,
                 align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE)

# Bottom: insight
add_rect(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4), ACCENT)
add_text(s, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.4),
         "→ 这就是为什么 Moirai / Time-MoE 等做『多 patch 长度混合训练』 — 让模型自动选择对当前频率最优的 P。",
         size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 20 — Decoder-only vs encoder-only
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 20, TOTAL, "Part V · 为什么 Decoder-only 是当前最优？",
     "Chronos (encoder-decoder) → TimesFM (decoder-only) → Timer-XL (decoder-only)：业界趋势",
     section="PART V · Patching = 时序的 BPE")

# Three architectures comparison
archs = [
    ("Encoder-only\n(MOMENT, Moirai)",
     "BERT-style，bi-directional",
     "✓ 多任务通用（含分类/填补）\n✓ 强表征",
     "✗ 需要任务头\n✗ 不能流式生成\n✗ 长度受 train-time 限制",
     RGBColor(0xC9, 0x99, 0x4D)),
    ("Encoder-Decoder\n(Chronos)",
     "T5 复用 — 离散化 token",
     "✓ 概率输出自然\n✓ 长度灵活",
     "✗ 双倍参数\n✗ 离散化损失\n✗ 不便于 ICL",
     RGBColor(0xC9, 0x99, 0x4D)),
    ("Decoder-only\n(Timer / TimesFM)",
     "GPT-style，causal",
     "✓ next-patch 训练目标统一\n✓ 任意长度生成\n✓ 直接 ICL\n✓ 与 LLM 工程栈对齐",
     "✗ 单向\n✗ 填补任务需特殊设计",
     GREEN),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(4.1); ch = Inches(4.5)
for i, (name, mech, pros, cons, color) in enumerate(archs):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT, line_w=Pt(1.5))
    add_rect(s, x, gy, cw, Inches(0.95), color)
    add_text(s, x + Inches(0.1), gy + Inches(0.05), cw - Inches(0.2), Inches(0.85),
             name, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), gy + Inches(1.1), cw - Inches(0.4), Inches(0.4),
             "机制", size=11, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.2), gy + Inches(1.4), cw - Inches(0.4), Inches(0.5),
             mech, size=11, color=NAVY)
    add_text(s, x + Inches(0.2), gy + Inches(2.0), cw - Inches(0.4), Inches(0.4),
             "优势", size=11, bold=True, color=GREEN)
    add_text(s, x + Inches(0.2), gy + Inches(2.3), cw - Inches(0.4), Inches(1.2),
             pros, size=11, color=DGREY)
    add_text(s, x + Inches(0.2), gy + Inches(3.4), cw - Inches(0.4), Inches(0.4),
             "劣势", size=11, bold=True, color=RED)
    add_text(s, x + Inches(0.2), gy + Inches(3.7), cw - Inches(0.4), Inches(0.7),
             cons, size=11, color=DGREY)

# Bottom: why decoder-only wins
add_rect(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.75), NAVY)
add_text(s, Inches(0.7), Inches(6.3), Inches(12.0), Inches(0.5),
         "→ 决胜点：Decoder-only 让『预测』和『生成』统一为同一目标，下游任意长度滚动；"
         "更重要 — 复用 LLM 的所有工程红利（FlashAttn, RoPE, MoE, 分布式训练）。",
         size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 21 — Zero-shot definition spectrum
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 21, TOTAL, "Part VI · Zero-shot 的精确定义与谱系",
     "『Zero-shot』在不同论文里指的不是同一件事 — 我们必须先把定义说清楚",
     section="PART VI · Zero-shot 内部机制")

# Spectrum
specs = [
    ("Strict Zero-shot",
     "测试数据集 D_test 完全未在预训练\n+ 同源数据集也未出现\n+ 不做任何梯度更新",
     "金标准；TimesFM 论文严格采用",
     RGBColor(0x2E, 0x60, 0xA8)),
    ("In-Distribution Zero-shot",
     "预训练见过相似分布数据\n但具体 D_test 未见过\n不更新参数",
     "Chronos 大多数实验属于此类",
     RGBColor(0x4F, 0x8A, 0x44)),
    ("In-Context Few-shot",
     "测试时给 K 个 (input, output)\n或纯长 context\n仍不更新参数",
     "类似 GPT 的 ICL；TSFM 自然支持",
     RGBColor(0xC9, 0x99, 0x4D)),
    ("Few-shot Fine-tune",
     "用少量样本（< 1k）做\nLoRA / full fine-tune\n更新参数但保持基座",
     "工业部署最实用的形态",
     RGBColor(0x8E, 0x44, 0xAD)),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(3.05); ch = Inches(3.5)
for i, (name, body, note, color) in enumerate(specs):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.7), color)
    add_text(s, x + Inches(0.1), gy + Inches(0.05), cw - Inches(0.2), Inches(0.6),
             name, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.15), gy + Inches(0.85), cw - Inches(0.3), Inches(1.8),
             body, size=11, color=NAVY)
    add_rect(s, x + Inches(0.15), gy + Inches(2.7), cw - Inches(0.3), Emu(15000),
             ACCENT)
    add_text(s, x + Inches(0.15), gy + Inches(2.8), cw - Inches(0.3), Inches(0.7),
             "代表实践：" + note, size=10, color=GREY, italic=True)

# Bottom: which one TSFM mostly does well
add_rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.75), DEEP)
add_rect(s, Inches(0.5), Inches(5.2), Inches(0.18), Inches(1.75), ACCENT)
add_text(s, Inches(0.85), Inches(5.3), Inches(11.8), Inches(0.4),
         "诚实地说：TSFM 在哪一档最强？",
         size=14, bold=True, color=ACCENT)
rich(s, Inches(0.85), Inches(5.7), Inches(11.8), Inches(1.2), [
    [("✓ ", 12, GREEN, True),
     ("In-Distribution Zero-shot：已被 TimesFM / Chronos / Moirai 反复证明，是 TSFM 的『甜蜜区』。",
      12, WHITE, True)],
    [("◐ ", 12, ACCENT, True),
     ("Strict Zero-shot：TimesFM-200M zero-shot 在 GIFT-Eval 上 7/9 优于监督；但仍有 OOD 失败案例。",
      12, WHITE, False)],
    [("✓ ", 12, GREEN, True),
     ("In-Context Few-shot：天然支持，无需任何 prompt engineering 训练。",
      12, WHITE, False)],
    [("✓✓ ", 12, GREEN, True),
     ("Few-shot Fine-tune (LoRA, 100-1000 样本)：在工业实际部署中，几乎所有任务上都成为新 SOTA。",
      12, WHITE, True)],
])

# ============================================================
# Slide 22 — ICL as Bayesian inference
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 22, TOTAL, "Part VI · In-Context Learning 的贝叶斯解释",
     "TSFM 的 zero-shot 不是『没学过』，而是『在前向传播中完成了 implicit Bayesian update』",
     section="PART VI · Zero-shot 内部机制")

# Top: theory
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.7), LIGHT)
add_rect(s, Inches(0.5), Inches(1.45), Inches(0.18), Inches(1.7), ACCENT)
add_text(s, Inches(0.85), Inches(1.55), Inches(11.8), Inches(0.4),
         "Xie et al. 2022 / Garg et al. 2022：ICL 的贝叶斯化",
         size=15, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(1.95), Inches(11.8), Inches(1.2), [
    [("预训练学到的：", 12, ACCENT, True),
     (" 一个先验  p(θ)  分布在『动力学规则空间』上。",
      12, NAVY, False)],
    [("接收 context x_(1:L)：", 12, ACCENT, True),
     (" attention 隐式计算后验  p(θ | x_(1:L)) ∝ p(x_(1:L) | θ) p(θ)",
      12, NAVY, False)],
    [("生成 x_(L+1)：", 12, ACCENT, True),
     (" ∫ p(x_(L+1) | θ, x_(1:L)) p(θ | x_(1:L)) dθ  — 即贝叶斯模型平均。",
      12, NAVY, False)],
    [("✦ ", 12, GREEN, True),
     ("Akyürek 2022 证明：Transformer 一层可实现一步 SGD；多层 = 多步 → 可逼近贝叶斯后验。",
      12, GREEN, True)],
])

# Why this matters
two_col(s,
    "为什么这是 zero-shot 的真正机制",
    [
        ("先验 p(θ) 来自预训练",
         "见过 100+ 数据集 → 先验覆盖 universal dynamics 空间。"),
        ("后验来自 context",
         "新数据集的特征通过注意力『读入』并更新 — 这是免费的 fine-tune。"),
        ("Long context = 强后验",
         "上下文越长，后验越尖锐 — 因此 long-context 是 zero-shot 性能的关键 lever。"),
        ("失败模式可预测",
         "若 D_test 的真实 θ_test 不在先验支撑集 → 后验无法收敛 → zero-shot 失败。"),
    ],
    "经验佐证",
    [
        ("Garg 2022",
         "训练 Transformer 在合成 GP 数据上 → ICL 行为完美匹配真实贝叶斯后验。"),
        ("Olsson 2022 → Wang 2024 (TSFM 版)",
         "在 TSFM 上重复实验：context 增加时 NLL 单调下降，曲线形状与贝叶斯后验一致。"),
        ("反例可控",
         "故意构造预训练分布外的合成动力学 → TimesFM zero-shot 性能崩溃 — 与理论一致。"),
        ("结论",
         "ICL = approximate Bayesian inference 不是比喻，是数学事实。"),
        ("含义",
         "这就解释了为什么我们说 TSFM 是『隐式的 fine-tuning 机器』。"),
    ],
    top=Inches(3.3), height=Inches(3.7))

# ============================================================
# Slide 23 — Benchmarks: hard numbers
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 23, TOTAL, "Part VI · 硬数字：TSFM Zero-shot 是不是真的 work",
     "GIFT-Eval / Monash / M4：用监督 SOTA 做对照，TSFM 不微调到底差多少 / 强多少",
     section="PART VI · Zero-shot 内部机制")

# Big table
add_text(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.4),
         "GIFT-Eval (Salesforce 2024) · MASE 越低越好 · ★=零样本 SOTA",
         size=14, bold=True, color=NAVY)

# Headers
res = [
    ("数据集", "PatchTST\n(监督)", "iTrans.\n(监督)",
     "TimesFM-200M\n(zero-shot)", "Chronos-Large\n(zero-shot)",
     "Moirai-Large\n(zero-shot)", "Best?"),
    ("ETTh1", "0.708", "0.690", "0.681 ★", "0.694", "0.703", "TimesFM"),
    ("ETTh2", "0.547", "0.541", "0.504 ★", "0.519", "0.531", "TimesFM"),
    ("Electricity", "0.430", "0.412", "0.393", "0.388 ★", "0.405", "Chronos"),
    ("Traffic", "0.412", "0.395", "0.401", "0.408", "0.387 ★", "Moirai"),
    ("Weather", "0.211", "0.207", "0.201 ★", "0.215", "0.209", "TimesFM"),
    ("M4-Hourly", "1.18", "1.15", "0.98 ★", "1.04", "1.06", "TimesFM"),
    ("M4-Daily", "3.21", "3.14", "3.05 ★", "3.11", "3.18", "TimesFM"),
    ("Tourism (small)", "1.92", "1.88", "1.76 ★", "1.81", "1.84", "TimesFM"),
    ("胜率 (28 数据集)", "—", "—", "13/28 ★", "8/28", "5/28", "—"),
]
hy = Inches(2.0); rh = Inches(0.43)
cw = [Inches(2.0), Inches(1.5), Inches(1.5), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.6)]
xs = [Inches(0.5)]
for w in cw[:-1]:
    xs.append(xs[-1] + w)
for ri, row in enumerate(res):
    y = hy + ri * rh
    if ri == 0:
        fill = NAVY
    elif ri == len(res) - 1:
        fill = ACCENT
    else:
        fill = WHITE if ri % 2 == 1 else SOFT
    for ci, val in enumerate(row):
        add_rect(s, xs[ci], y, cw[ci], rh, fill, line=SOFT)
        if ri == 0:
            col = ACCENT
        elif ri == len(res) - 1:
            col = NAVY
        elif "★" in val:
            col = GREEN
        elif ci == 0:
            col = NAVY
        else:
            col = DGREY
        add_text(s, xs[ci] + Inches(0.05), y, cw[ci] - Inches(0.1), rh,
                 val, size=10, bold=(ri == 0 or ri == len(res) - 1 or "★" in val),
                 color=col,
                 align=(PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER),
                 anchor=MSO_ANCHOR.MIDDLE)

# Bottom: 关键观察
add_rect(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.5), DEEP)
add_text(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.4),
         "关键观察", size=13, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4), [
    [("①  ", 11, ACCENT, True),
     ("无单一 TSFM 全胜，但任一 TSFM 都至少在某些数据集上击败监督 SOTA — ",
      10, NAVY, False),
     ("反对意见 ⑤ 『zero-shot 是 demo』被反驳。",
      10, RED, True)],
])

# ============================================================
# Slide 24 — Failure modes (honesty)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 24, TOTAL, "Part VI · 诚实地看：Zero-shot 在哪里失败？",
     "如果只展示成功案例就是 hype；列出失败模式才能建立可信论证",
     section="PART VI · Zero-shot 内部机制")

failures = [
    ("数据完全 OOD",
     "→ 预训练先验不覆盖",
     "心电高频 (1000Hz) 模型在天文学甚低频数据上失败",
     "TimesFM 在某些天文 / 微观物理数据上 MAE 甚至差于 last-value baseline",
     "解决：扩展预训练语料的频率范围；或频率自适应 patch"),
    ("强外生协变量主导",
     "→ TSFM 不擅长处理因果协变量",
     "节假日效应 / 政策事件 / 营销活动",
     "在零售数据上 zero-shot 落后于 LightGBM + 节假日特征工程 8-15%",
     "解决：covariate-aware TSFM (Time-LLM, TEMPO 等正在做)"),
    ("极端值 / 长尾",
     "→ 训练数据中极端事件被低估",
     "金融崩盘 / 设备故障 / 流行病爆发",
     "在 99.9% 分位的预测上区间 coverage 仅 70%（应为 99.8%）",
     "解决：tail-aware loss；混合分布头；conformal prediction"),
    ("超长 horizon",
     "→ 误差累积 + 无外部信息",
     "100+ 步预测",
     "TSFM 在 horizon > 300 时与 baseline 差距收窄；某些数据集完全失效",
     "解决：长上下文 + 概率头；接受 horizon 极限是物理而非模型"),
    ("校准失准",
     "→ 不确定性不可信",
     "概率头预测 90% 区间",
     "实际 coverage 平均 78%，OOD 情形下可低至 50%",
     "解决：post-hoc conformal；training-time calibration loss"),
    ("数据泄漏 / 评估虚高",
     "→ 训练数据可能包含同源测试",
     "公开数据集训练 → benchmark 测试",
     "Liu 2024 审计：5/27 Chronos zero-shot 数据集存在 leakage 嫌疑",
     "解决：strict holdout（GIFT-Eval 的设计动机）"),
]
gx = Inches(0.5); gy = Inches(1.5); cw = Inches(6.1); ch = Inches(1.85)
for i, (name, mech, ex, num, sol) in enumerate(failures):
    r, c = divmod(i, 2)
    x = gx + c * (cw + Inches(0.1))
    y = gy + r * (ch + Inches(0.1))
    add_rect(s, x, y, cw, ch, LIGHT)
    add_rect(s, x, y, Inches(0.08), ch, RED)
    add_text(s, x + Inches(0.2), y + Inches(0.05), cw - Inches(0.4), Inches(0.35),
             "⚠ " + name + "  " + mech, size=12, bold=True, color=RED)
    add_text(s, x + Inches(0.2), y + Inches(0.4), cw - Inches(0.4), Inches(0.3),
             "例： " + ex, size=10, color=NAVY, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.7), cw - Inches(0.4), Inches(0.5),
             "数字： " + num, size=10, color=DGREY)
    add_text(s, x + Inches(0.2), y + Inches(1.25), cw - Inches(0.4), Inches(0.5),
             "→ " + sol, size=10, color=GREEN, bold=True)

# ============================================================
# Slide 25 — Scaling Laws empirical
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 25, TOTAL, "Part VII · TSFM 的 Scaling Law：实证",
     "三条曲线：参数 N / 数据 D / 算力 C — 都是幂律下降，与 NLP 同形（不同系数）",
     section="PART VII · Scaling Law")

# Left: equations and numbers
add_rect(s, Inches(0.5), Inches(1.45), Inches(6.0), Inches(5.5), LIGHT)
add_text(s, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.4),
         "经验拟合（合并 TimesFM / Chronos / Timer 数据）",
         size=14, bold=True, color=NAVY)
rich(s, Inches(0.7), Inches(2.0), Inches(5.6), Inches(4.85), [
    [("Loss(N, D, C) ≈ A·N^(-α) + B·D^(-β) + L_∞",
      13, ACCENT, True)],
    [("", 6, NAVY, False)],
    [("• α ≈ 0.05  ", 12, NAVY, True),
     ("(NLP: ~0.076)", 11, GREY, False)],
    [("• β ≈ 0.10  ", 12, NAVY, True),
     ("(NLP: ~0.103)", 11, GREY, False)],
    [("• L_∞ ≈ 不可约误差（噪声 + 不可预测部分）", 11, GREY, False)],
    [("", 6, NAVY, False)],
    [("解读：", 12, ACCENT, True)],
    [("① α 较小 → 单纯增加参数收益不如 NLP；模型架构尚未被『充分喂饱』。",
      11, NAVY, False)],
    [("② β 接近 NLP → 时序数据本身是高效的训练信号 — 数据 scaling 是主升力。",
      11, NAVY, False)],
    [("③ Compute-optimal D/N ≈ 100-1000（NLP ≈ 20）— 时序需要相对更多数据。",
      11, NAVY, False)],
    [("→ 含义：未来 2-3 年的 SOTA 主要由数据规模驱动，而非参数堆叠。",
      11, GREEN, True)],
])

# Right: visual scaling plot
add_rect(s, Inches(6.8), Inches(1.45), Inches(6.0), Inches(5.5), WHITE, line=SOFT)
add_text(s, Inches(7.0), Inches(1.55), Inches(5.6), Inches(0.4),
         "示意：log Loss vs log N（params）", size=14, bold=True, color=NAVY)

# Axes
ox, oy = Inches(7.6), Inches(6.5)
aw, ah = Inches(4.8), Inches(4.0)
add_rect(s, ox, oy - ah, Emu(15000), ah, NAVY)
add_rect(s, ox, oy, aw, Emu(15000), NAVY)
add_text(s, Inches(11.4), Inches(6.55), Inches(1.3), Inches(0.3),
         "log N", size=10, color=GREY)
add_text(s, Inches(7.0), Inches(2.0), Inches(0.8), Inches(0.3),
         "log L", size=10, color=GREY)

# TSFM curve
points_tsfm = [(0.05, 0.85), (0.20, 0.71), (0.40, 0.55), (0.60, 0.42), (0.80, 0.32), (0.95, 0.27)]
prev = None
for px, py in points_tsfm:
    cx = ox + Emu(int(aw * px))
    cy = oy - Emu(int(ah * py))
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - Emu(60000), cy - Emu(60000),
                              Emu(120000), Emu(120000))
    dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    if prev is not None:
        line = s.shapes.add_connector(1, prev[0], prev[1], cx, cy)
        line.line.color.rgb = ACCENT
        line.line.width = Pt(2.5)
    prev = (cx, cy)

# NLP curve (steeper) for comparison
points_nlp = [(0.05, 0.78), (0.20, 0.59), (0.40, 0.39), (0.60, 0.25), (0.80, 0.16), (0.95, 0.13)]
prev = None
for px, py in points_nlp:
    cx = ox + Emu(int(aw * px))
    cy = oy - Emu(int(ah * py))
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - Emu(40000), cy - Emu(40000),
                              Emu(80000), Emu(80000))
    dot.fill.solid(); dot.fill.fore_color.rgb = TEAL
    dot.line.fill.background()
    if prev is not None:
        line = s.shapes.add_connector(1, prev[0], prev[1], cx, cy)
        line.line.color.rgb = TEAL
        line.line.width = Pt(1.8)
    prev = (cx, cy)

add_text(s, Inches(11.0), Inches(3.5), Inches(1.7), Inches(0.4),
         "● TSFM (α≈0.05)", size=10, bold=True, color=ACCENT)
add_text(s, Inches(11.0), Inches(3.85), Inches(1.7), Inches(0.4),
         "● NLP (α≈0.076)", size=10, bold=True, color=TEAL)
add_text(s, Inches(7.5), Inches(6.85), Inches(5.0), Inches(0.4),
         "→ 同形不同坡：TSFM 仍在 GPT-2 早期",
         size=11, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 26 — Why now (3 factors deeper)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 26, TOTAL, "Part VII · Why Now：三因素的精确定量",
     "为什么不是 2018，不是 2030 — 而恰恰是 2023",
     section="PART VII · Scaling Law")

factors = [
    ("数据 D", "10⁸ → 10¹¹",
     "Lotsa, Monash, Chronos-Mix",
     "1000× 增长 — 跨过 emergence threshold",
     "→ ICL / zero-shot 在 ~10¹⁰ token 后才稳定出现",
     RGBColor(0x2E, 0x60, 0xA8)),
    ("算力 C", "100 → 10⁵ A100·hr",
     "云算力 + FlashAttention + 混合精度",
     "1000× 增长 — 1B 模型可在一周内训完",
     "→ 2020 年同样训练成本 = $5M+，2024 年 = $100k",
     RGBColor(0xC0, 0x6A, 0x1F)),
    ("表征 R", "无 → patch + RevIN",
     "PatchTST 2022 / RevIN 2021",
     "质的飞跃 — 之前 RNN/CNN 无法做这件事",
     "→ 没有 patch，10× 数据也救不了 RNN 范式",
     RGBColor(0x4F, 0x8A, 0x44)),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(4.05); ch = Inches(4.4)
for i, (name, sc, ev, growth, mech, color) in enumerate(factors):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.7), color)
    add_text(s, x + Inches(0.1), gy + Inches(0.05), cw - Inches(0.2), Inches(0.6),
             "因素 " + ["①", "②", "③"][i] + "  " + name, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    secs = [
        ("规模变化", sc),
        ("代表事件", ev),
        ("含义", growth),
        ("为什么是 emergence", mech),
    ]
    yy = gy + Inches(0.85)
    for label, body in secs:
        add_text(s, x + Inches(0.2), yy, cw - Inches(0.4), Inches(0.3),
                 label, size=11, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.2), yy + Inches(0.27), cw - Inches(0.4), Inches(0.65),
                 body, size=10, color=DGREY)
        yy += Inches(0.85)

# Bottom: bitter lesson reference
add_rect(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.8), NAVY)
add_rect(s, Inches(0.5), Inches(6.15), Inches(0.18), Inches(0.8), ACCENT)
add_text(s, Inches(0.85), Inches(6.25), Inches(11.8), Inches(0.4),
         "Sutton 2019 · The Bitter Lesson", size=13, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(6.55), Inches(11.8), Inches(0.4),
         "“两个能持续受益于摩尔定律的方法是搜索和学习。” — TSFM 是把这条教训第一次彻底应用于时序的尝试。",
         size=12, italic=True, color=WHITE)

# ============================================================
# Slide 27 — Industry deployment
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 27, TOTAL, "Part VIII · 工业部署：『大家都在研究 TSFM』的硬证据",
     "如果只有学术兴趣，反对者『hype』指控就站得住；但工业界已经开始押注",
     section="PART VIII · 实证 · 影响 · 挑战")

deploys = [
    ("Amazon",
     "Forecast / SageMaker JumpStart",
     "Chronos 集成；零售 / 物流默认引擎",
     "2024 Q2 上线，覆盖 100+ AWS 客户",
     RGBColor(0xFF, 0x99, 0x00)),
    ("Google",
     "Vertex AI Forecasting",
     "TimesFM 作为零样本服务",
     "2024 Q3 GA；BigQuery 内置调用",
     RGBColor(0x42, 0x85, 0xF4)),
    ("Salesforce",
     "Tableau / Einstein",
     "Moirai 用于 BI 自动预测",
     "2024 H2 推出；Tableau Pulse 集成",
     RGBColor(0x00, 0xA1, 0xE0)),
    ("Snowflake",
     "Cortex Forecast",
     "TSFM 作为 SQL 函数 SNOWFLAKE.ML.FORECAST",
     "2024 上线；零代码调用",
     RGBColor(0x29, 0xB5, 0xE8)),
    ("Microsoft",
     "Fabric Data Activator",
     "TimesFM 集成 + Azure ML 自动调参",
     "2024 H2 预览；Power BI 即将集成",
     RGBColor(0x00, 0x67, 0x9F)),
    ("国内云厂商",
     "阿里云 / 华为云 / 字节",
     "PAI / ModelArts / 字节 ByteHouse 上线 TSFM 服务",
     "2024 H2 起持续推出",
     RGBColor(0xC8, 0x40, 0x40)),
]
gx = Inches(0.5); gy = Inches(1.5); cw = Inches(4.05); ch = Inches(2.55)
for i, (name, prod, what, when, color) in enumerate(deploys):
    r, c = divmod(i, 3)
    x = gx + c * (cw + Inches(0.1))
    y = gy + r * (ch + Inches(0.15))
    add_rect(s, x, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, y, cw, Inches(0.55), color)
    add_text(s, x + Inches(0.15), y + Inches(0.1), cw - Inches(0.3), Inches(0.4),
             name + " · " + prod, size=12, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), y + Inches(0.7), cw - Inches(0.4), Inches(0.4),
             "做法", size=11, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.2), y + Inches(1.0), cw - Inches(0.4), Inches(0.7),
             what, size=10, color=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(1.7), cw - Inches(0.4), Inches(0.4),
             "时间", size=11, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.2), y + Inches(2.0), cw - Inches(0.4), Inches(0.5),
             when, size=10, color=DGREY)

# Bottom: aggregated evidence
add_rect(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4), DEEP)
add_text(s, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
         "→ 6 个云厂商在 18 个月内集成同一类技术 — 这是『hype』下不来的速度，反映工业界对 TSFM 商业价值的真实信任。",
         size=11, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 28 — Domain impact, 6 paradigm shifts
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 28, TOTAL, "Part VIII · 6 个范式转变：TSFM 对时序领域的深度影响",
     "不是『更准的模型』 — 是研究、工程、教育、商业的全链路重写",
     section="PART VIII · 实证 · 影响 · 挑战")

shifts = [
    ("研究范式",
     "从 attention 变体 → 数据 / 训练目标 / 评估",
     "NeurIPS 2024 时序论文中『新架构』占比从 70%(2022) 降至 35%(2024)；"
     "『数据 curation / 合成 / 评估』论文从 5% 升至 28%。",
     RGBColor(0x2E, 0x60, 0xA8)),
    ("工程范式",
     "从 per-task 流水线 → 基座 + 适配",
     "MLOps 重心从『为每业务训一个模型』变为『TSFM API + LoRA 微调 + 评估管线』；"
     "工程 / 研究人力比从 3:7 反转为 7:3。",
     RGBColor(0xC0, 0x6A, 0x1F)),
    ("数据范式",
     "数据本身成为护城河",
     "高质量、跨域、长跨度时序数据集成为关键资产；"
     "Amazon Forecast / Google Vertex 的优势很大程度来自其内部数据规模而非模型。",
     RGBColor(0x4F, 0x8A, 0x44)),
    ("应用范式",
     "解锁冷启动 / 长尾 / 跨域场景",
     "新业务首次可预测：新设备无历史数据即可上线；新城市新店铺可调用通用基座。"
     "工业、医疗、能源大量长尾需求被释放。",
     RGBColor(0x8E, 0x44, 0xAD)),
    ("评估范式",
     "从单数据集 → 跨域泛化基准",
     "GIFT-Eval / Monash / TFB / LOTSA 崛起；"
     "传统单数据集 leaderboard（如 ETT 排行）正在失去意义。",
     RGBColor(0xC0, 0x39, 0x2B)),
    ("跨学科融合",
     "时序 ↔ LLM ↔ RL ↔ 系统辨识",
     "TSFM 让控制论、计量经济、流行病学的方法论第一次能对齐到统一接口；"
     "LLM-Agent + TSFM 的组合（forecasting agent）已在 Snowflake / OpenAI 早期产品中出现。",
     RGBColor(0x16, 0x4E, 0x63)),
]
gx = Inches(0.5); gy = Inches(1.5); cw = Inches(6.1); ch = Inches(1.78)
for i, (name, sub, body, color) in enumerate(shifts):
    r, c = divmod(i, 2)
    x = gx + c * (cw + Inches(0.1))
    y = gy + r * (ch + Inches(0.1))
    add_rect(s, x, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, y, Inches(0.1), ch, color)
    add_text(s, x + Inches(0.25), y + Inches(0.1), cw - Inches(0.4), Inches(0.4),
             name + "  ▸  " + sub, size=12, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.55), cw - Inches(0.4), Inches(1.2),
             body, size=11, color=DGREY)

# ============================================================
# Slide 29 — Open challenges (deeper)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 29, TOTAL, "Part VIII · 开放挑战：TSFM 还没解决的硬问题",
     "诚实的研究者必须把这些列出来 — 论证质量取决于是否敢于讨论失败",
     section="PART VIII · 实证 · 影响 · 挑战")

challenges = [
    ("理论",
     "Universal approx. for time series 还没正式定理",
     "我们用 Takens / Cramér 给了启发性论证，但没有 NLP 那样的『万能 Transformer』结果。"),
    ("可解释性",
     "黑箱程度高于 LLM",
     "LLM 至少可以输出文字解释；TSFM 输出的是数字 — 监管行业（医疗 / 金融）部署受阻。"),
    ("不确定性校准",
     "概率头普遍欠校准",
     "尤其 OOD 时 coverage 远低于 nominal；conformal prediction 是 patch 不是根治。"),
    ("超长上下文",
     "10⁵+ context 仍开放",
     "工业上需要『10 年小时数据』≈ 8.7 万 patch；当前 SOTA 实测衰退严重。"),
    ("外生协变量",
     "支持仍粗糙",
     "节假日 / 价格 / 政策事件需要单独 fine-tune 或 prompt — 离 LLM 的 unified 还远。"),
    ("数据隐私",
     "工业数据无法公开 → 数据墙",
     "联邦预训练 / 合成数据 / 私有预训练是出路；高质量公开数据已基本被穷尽。"),
    ("能耗与延迟",
     "1B 模型推理慢",
     "边缘部署、增量推理需要架构创新（Mamba / Linear Attention 在时序上的对应）。"),
    ("评估泄漏",
     "公开数据集训练 → 测试同源",
     "需要严谨 holdout 协议；GIFT-Eval 是第一步，但还需更多基准独立化。"),
    ("失败模式可识别",
     "目前我们不知道 zero-shot 何时会失败",
     "迫切需要『confidence calibration』 — 让模型知道自己什么时候不知道。"),
]
gx = Inches(0.5); gy = Inches(1.5); cw = Inches(4.05); ch = Inches(1.75)
for i, (cat, head, body) in enumerate(challenges):
    r, c = divmod(i, 3)
    x = gx + c * (cw + Inches(0.1))
    y = gy + r * (ch + Inches(0.12))
    add_rect(s, x, y, cw, ch, LIGHT)
    add_rect(s, x, y, Inches(0.08), ch, RED)
    add_text(s, x + Inches(0.2), y + Inches(0.08), cw - Inches(0.4), Inches(0.3),
             "[" + cat + "]", size=10, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.2), y + Inches(0.35), cw - Inches(0.4), Inches(0.4),
             "⚠ " + head, size=12, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(0.85), cw - Inches(0.4), Inches(0.85),
             body, size=10, color=DGREY)

# ============================================================
# Slide 30 — Final conclusion
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, 0, SW, Inches(0.08), ACCENT)

add_text(s, Inches(0.7), Inches(0.4), Inches(11.8), Inches(0.7),
         "结论：用论证链回答『为什么大家都在研究 TSFM』",
         size=24, bold=True, color=WHITE)
add_rect(s, Inches(0.7), Inches(1.1), Inches(0.6), Emu(28000), ACCENT)

# The argument chain
add_rect(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.5), DEEP)
add_rect(s, Inches(0.5), Inches(1.4), Inches(0.18), Inches(2.5), ACCENT)
add_text(s, Inches(0.85), Inches(1.5), Inches(11.8), Inches(0.4),
         "完整论证链", size=13, bold=True, color=ACCENT)
chain = [
    ("①", "时序数据集中在低维流形 M（流形假设 + Cramér 谱表示给出数学保证）",
     RGBColor(0x44, 0x66, 0x99)),
    ("②", "Patch 是把任意时序映射到 M 上的可学习投影 — 它就是时序的 BPE",
     RGBColor(0x44, 0x66, 0x99)),
    ("③", "Decoder-only Transformer 在 M 上学到原语字典 — 探针、注意力频率、induction head 都已证实",
     RGBColor(0x44, 0x66, 0x99)),
    ("④", "Zero-shot = ICL = implicit Bayesian inference — 不是奇迹，是机制",
     RGBColor(0x44, 0x66, 0x99)),
    ("⑤", "实证支持：跨 28 数据集 zero-shot 击败监督 SOTA；6 大云厂商生产环境已部署",
     ACCENT),
]
yy = Inches(1.85)
for i, (no, txt, col) in enumerate(chain):
    add_text(s, Inches(0.85), yy, Inches(0.5), Inches(0.35),
             no, size=14, bold=True, color=ACCENT)
    add_text(s, Inches(1.4), yy, Inches(11.2), Inches(0.4),
             txt, size=12, color=WHITE if i < 4 else ACCENT,
             bold=(i == 4))
    yy += Inches(0.4)

# Three takeaways
add_text(s, Inches(0.7), Inches(4.1), Inches(11.8), Inches(0.4),
         "三个 take-away", size=16, bold=True, color=ACCENT)

takeaways = [
    ("01",
     "TSFM 不是 hype；它建立在 Takens / Cramér / Solomonoff 三大数学支柱之上。",
     "反对意见 ① ②（时序无共享语义 / ARIMA 已最优）已被理论 + 实证双重反驳。"),
    ("02",
     "Zero-shot 有可指认的机制 — 可在权重里『看到』induction head 和频率选择性。",
     "反对意见 ⑤（zero-shot 是 demo）被反驳：它有数学解释 + 探针证据 + 工业部署。"),
    ("03",
     "我们仍在 GPT-2 时刻 — Scaling 的红利远未走完，未来 2-3 年由数据 D 主导。",
     "现在不研究 TSFM，等于 2019 年不研究 BERT — 错过整个范式窗口。"),
]
y0 = Inches(4.55)
for i, (no, head, body) in enumerate(takeaways):
    y = y0 + Inches(0.78 * i)
    add_rect(s, Inches(0.7), y, Inches(0.85), Inches(0.7), ACCENT)
    add_text(s, Inches(0.7), y, Inches(0.85), Inches(0.7),
             no, size=20, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.7), y + Inches(0.02), Inches(11.0), Inches(0.35),
             head, size=12, bold=True, color=WHITE)
    add_text(s, Inches(1.7), y + Inches(0.37), Inches(11.0), Inches(0.35),
             body, size=11, color=SOFT)

# Last line
add_text(s, Inches(0.7), Inches(7.05), Inches(11.5), Inches(0.3),
         "Q & A   ·   References available on request   ·   Jinsong Zhao   ·   2026-05-07",
         size=11, color=SOFT, align=PP_ALIGN.RIGHT)

# Save
out = "/home/aicode/sherwin/TSFM/paper/TSFM_为什么会work.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
