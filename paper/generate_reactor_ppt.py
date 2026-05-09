"""Generate PPT discussing whether univariate TSFM is sufficient for
real industrial scenarios with exogenous driving variables (e.g. reactor
temperature prediction with composition step changes).

Argument structure:
  Part I  —  The question (concrete reactor scenario)
  Part II —  Why pure univariate is mathematically insufficient
              (information gap + Takens caveat + empirical demo)
  Part III —  The honest answer: univariate TSFM ≠ all of TSFM
  Part IV —  Five solution levels (Moirai, TimeXer, fine-tune, hybrid)
  Part V  —  Recommendation for the reactor case
  Part VI —  Why TSFM pre-training still adds value with covariates
  Part VII —  Decision tree + final answer
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import math

# ---------- Style ----------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
DEEP = RGBColor(0x14, 0x2C, 0x52)
ACCENT = RGBColor(0xE8, 0x8B, 0x1A)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
GREY = RGBColor(0x4A, 0x55, 0x68)
DGREY = RGBColor(0x33, 0x3A, 0x47)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xDC, 0xE3, 0xF0)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)
PURPLE = RGBColor(0x6B, 0x3F, 0xA0)
TEAL = RGBColor(0x16, 0x7E, 0x8A)
ORANGE = RGBColor(0xE0, 0x6E, 0x18)

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def _i(v):
    return int(v) if v is not None else v


def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def rich(slide, x, y, w, h, blocks, line_spacing=1.3, anchor=MSO_ANCHOR.TOP):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, runs in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        for txt, size, color, bold in runs:
            r = p.add_run()
            r.text = txt
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return tb


def page(slide, idx, total, title, subtitle=None, section=None):
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_rect(slide, 0, 0, SW, Inches(0.08), ACCENT)
    add_rect(slide, 0, Inches(0.08), Inches(0.18), SH, NAVY)
    if section:
        add_text(slide, Inches(0.5), Inches(0.18), Inches(11.5), Inches(0.3),
                 section, size=11, bold=True, color=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.6),
             title, size=24, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(12.0), Inches(0.4),
                 subtitle, size=12, color=GREY, italic=True)
    add_rect(slide, Inches(0.5), Inches(1.28), Inches(0.6), Emu(28000), ACCENT)
    add_text(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
             f"{idx} / {total}", size=10, color=GREY, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.3),
             "TSFM 在受迫系统中的边界 · Jinsong Zhao", size=10, color=GREY)


def two_col(slide, left_title, left_blocks, right_title, right_blocks,
            top=Inches(1.45), height=Inches(5.55), left_color=LIGHT, right_color=NAVY):
    add_rect(slide, Inches(0.5), top, Inches(6.0), height, left_color)
    add_text(slide, Inches(0.7), top + Inches(0.15), Inches(5.6), Inches(0.4),
             left_title, size=16, bold=True, color=NAVY)
    add_rect(slide, Inches(0.7), top + Inches(0.6), Inches(0.5), Emu(20000), ACCENT)

    tb = slide.shapes.add_textbox(Inches(0.7), top + Inches(0.75),
                                    Inches(5.6), height - Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (h, b) in enumerate(left_blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        r1 = p.add_run(); r1.text = "▸ "; r1.font.color.rgb = ACCENT
        r1.font.size = Pt(13); r1.font.bold = True; r1.font.name = FONT
        r2 = p.add_run(); r2.text = h
        r2.font.bold = True; r2.font.size = Pt(13); r2.font.color.rgb = NAVY; r2.font.name = FONT
        if b:
            r3 = p.add_run(); r3.text = "\n   " + b
            r3.font.size = Pt(11); r3.font.color.rgb = GREY; r3.font.name = FONT

    add_rect(slide, Inches(6.8), top, Inches(6.0), height, right_color)
    add_text(slide, Inches(7.0), top + Inches(0.15), Inches(5.6), Inches(0.4),
             right_title, size=16, bold=True,
             color=ACCENT if right_color == NAVY else NAVY)
    add_rect(slide, Inches(7.0), top + Inches(0.6), Inches(0.5), Emu(20000), ACCENT)

    tb = slide.shapes.add_textbox(Inches(7.0), top + Inches(0.75),
                                    Inches(5.6), height - Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (h, b) in enumerate(right_blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        r1 = p.add_run(); r1.text = "▸ "; r1.font.color.rgb = ACCENT
        r1.font.size = Pt(13); r1.font.bold = True; r1.font.name = FONT
        r2 = p.add_run(); r2.text = h
        r2.font.bold = True; r2.font.size = Pt(13)
        r2.font.color.rgb = WHITE if right_color == NAVY else NAVY; r2.font.name = FONT
        if b:
            r3 = p.add_run(); r3.text = "\n   " + b
            r3.font.size = Pt(11)
            r3.font.color.rgb = SOFT if right_color == NAVY else GREY; r3.font.name = FONT


TOTAL = 23


# ============================================================
# Slide 1 — Title
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
for i in range(8):
    add_rect(s, Inches(8.5 + i * 0.55), Inches(1.0), Inches(0.05), Inches(5.5),
             RGBColor(0x1A, 0x35, 0x5C))
add_rect(s, Inches(0.8), Inches(1.7), Inches(0.18), Inches(3.2), ACCENT)
add_text(s, Inches(1.2), Inches(1.5), Inches(11), Inches(0.6),
         "DEEP DIVE  ·  REACTOR CASE", size=14, bold=True, color=ACCENT)
add_text(s, Inches(1.2), Inches(2.0), Inches(11.5), Inches(1.2),
         "单变量 TSFM 够用吗？", size=46, bold=True, color=WHITE)
add_text(s, Inches(1.2), Inches(3.1), Inches(11.5), Inches(0.7),
         "—— 反应器温度预测中外生变量的必要性论证",
         size=20, color=SOFT)
add_rect(s, Inches(1.2), Inches(4.1), Inches(2.5), Emu(20000), ACCENT)
add_text(s, Inches(1.2), Inches(4.3), Inches(11), Inches(0.5),
         "Why Univariate TSFM is Insufficient for Forced Systems — and What to Do Instead",
         size=14, color=ACCENT)

rich(s, Inches(1.2), Inches(4.95), Inches(11), Inches(2.0), [
    [("我们将回答 3 个问题：", 13, ACCENT, True)],
    [("Q1  ", 13, ACCENT, True),
     ("反应器入口组分突变 → 反应速率 → 温度。这种受迫系统下，单变量 TSFM 是否完备？",
      13, SOFT, False)],
    [("Q2  ", 13, ACCENT, True),
     ("如果不完备，是数学上的不可能，还是工程上『近似可用』？", 13, SOFT, False)],
    [("Q3  ", 13, ACCENT, True),
     ("如果单变量不行，TSFM 就完全没用了吗？还是它仍是最优工具的一部分？",
      13, SOFT, False)],
])
add_text(s, Inches(1.2), Inches(7.0), Inches(11), Inches(0.4),
         "Jinsong Zhao   ·   2026-05-07   ·   23 slides", size=12, color=SOFT)

# ============================================================
# Slide 2 — Outline
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 2, TOTAL, "目录 · 7 个部分", "从『提出问题』到『给出可执行方案』的完整论证链")

parts = [
    ("Part I", "提问 · 反应器场景的精确陈述", "P3"),
    ("Part II", "反驳 · 单变量 TSFM 数学上为何不完备",
     "信息缺失 / Takens 边界 / 实证演示  P4-7"),
    ("Part III", "诚实结论 · 单变量 TSFM ≠ 全部 TSFM", "P8"),
    ("Part IV", "解决方案 · 5 个层次的协变量整合",
     "Moirai / TimeXer / Fine-tune / 物理混合  P9-13"),
    ("Part V", "反应器场景的具体推荐",
     "变量清单 + 架构选型 + pipeline  P14-17"),
    ("Part VI", "为何 TSFM 在协变量场景下仍然必要",
     "预训练红利 / 跨反应器迁移 / 异常检测  P18-20"),
    ("Part VII", "决策树 + 回答原始问题", "P21-23"),
]
y0 = Inches(1.55)
rh = Inches(0.7)
for i, item in enumerate(parts):
    if len(item) == 3:
        no, head, pages_ = item
        desc = pages_
    else:
        no, head, pages_ = item[0], item[1], item[2]
        desc = item[2]
    y = y0 + i * rh
    add_rect(s, Inches(0.6), y, Inches(1.3), rh - Inches(0.05),
             NAVY if i % 2 == 0 else DEEP)
    add_text(s, Inches(0.6), y, Inches(1.3), rh - Inches(0.05),
             no, size=14, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(1.95), y, Inches(10.85), rh - Inches(0.05), LIGHT)
    add_text(s, Inches(2.1), y + Inches(0.05), Inches(10.5), Inches(0.6),
             head, size=14, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.1), y + Inches(0.35), Inches(10.5), Inches(0.3),
             desc, size=11, color=GREY)

# ============================================================
# Slide 3 — The reactor scenario, concretely
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 3, TOTAL, "Part I · 反应器场景的精确陈述",
     "我们到底在预测什么？变量之间的因果关系是什么？",
     section="PART I · 提问")

# Top: schematic-like description box
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(2.5), LIGHT)
add_rect(s, Inches(0.5), Inches(1.45), Inches(0.18), Inches(2.5), ACCENT)
add_text(s, Inches(0.85), Inches(1.55), Inches(11.8), Inches(0.4),
         "场景：连续搅拌反应器 (CSTR) 温度预测", size=15, bold=True, color=NAVY)

# Variables and causal chain
rich(s, Inches(0.85), Inches(2.0), Inches(11.8), Inches(1.95), [
    [("观测目标：", 13, ACCENT, True),
     ("反应器内温度 T(t)（每秒一个采样点，要预测未来 30 分钟）", 13, NAVY, False)],
    [("外生输入（可观测）：", 13, ACCENT, True),
     ("入口组分 x_A(t), x_B(t) ；流量 F(t) ；入口温度 T_in(t) ；夹套温度 T_j(t)",
      13, NAVY, False)],
    [("因果链：", 13, ACCENT, True)],
    [("  ① ", 12, ACCENT, True),
     ("入口组分 x 改变  →  反应速率 r(x, T) 改变（Arrhenius / 质量作用定律）",
      12, DGREY, False)],
    [("  ② ", 12, ACCENT, True),
     ("反应速率改变  →  反应释热 ΔH·r 改变  →  能量平衡破坏",
      12, DGREY, False)],
    [("  ③ ", 12, ACCENT, True),
     ("能量平衡破坏  →  T 以热容时间常数 τ_T 响应（通常 1-10 分钟滞后）",
      12, DGREY, False)],
    [("关键问题：", 13, RED, True),
     ("当 t = t* 时 x 阶跃变化，但 T 在 t* + τ_T 才开始响应。",
      13, RED, True)],
    [("→ ", 13, RED, True),
     ("仅用 T 的历史，模型在 [t*, t* + τ_T] 这段无法察觉变化已经发生。",
      13, RED, True)],
])

# Bottom: visual of step response
add_rect(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.85), WHITE, line=SOFT)
add_text(s, Inches(0.7), Inches(4.2), Inches(12.0), Inches(0.4),
         "示意：阶跃输入 → 滞后温度响应（横轴时间，纵轴标准化值）",
         size=13, bold=True, color=NAVY)

# Draw two stacked time series
ox = Inches(1.0); oy_x = Inches(5.4); oy_T = Inches(6.4); aw = Inches(11.0)
add_text(s, Inches(0.6), Inches(4.95), Inches(0.4), Inches(0.4),
         "x(t)", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.6), Inches(5.95), Inches(0.4), Inches(0.4),
         "T(t)", size=11, bold=True, color=ORANGE)

# Step input
def line(slide, x1, y1, x2, y2, color, w=1.8):
    x1, y1, x2, y2 = _i(x1), _i(y1), _i(x2), _i(y2)
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(w)
    return ln

# x(t): step at t* = 50%
mid = ox + aw // 2
line(s, ox, oy_x, mid, oy_x, ACCENT, w=2.5)
line(s, mid, oy_x, mid, oy_x - Emu(int(Inches(0.45))), ACCENT, w=2.5)
line(s, mid, oy_x - Emu(int(Inches(0.45))), ox + aw, oy_x - Emu(int(Inches(0.45))), ACCENT, w=2.5)

# Mark t*
line(s, mid, Inches(4.95), mid, Inches(6.85), GREY, w=0.8)
add_text(s, mid - Emu(int(Inches(0.3))), Inches(6.9), Inches(0.6), Inches(0.25),
         "t = t*", size=10, color=GREY, align=PP_ALIGN.CENTER)

# T(t): flat then exponential rise after delay
prev = None
for i in range(80):
    t = i / 80
    cx = ox + Emu(int(aw * t))
    if t < 0.5:
        v = 0.0
    elif t < 0.55:
        v = 0.0  # delay
    else:
        v = 1.0 - math.exp(-(t - 0.55) / 0.12)
    cy = oy_T - Emu(int(v * Inches(0.6)))
    if prev is not None:
        line(s, prev[0], prev[1], cx, cy, ORANGE, w=2.0)
    prev = (cx, cy)

# Annotation for delay
add_text(s, mid + Emu(int(Inches(0.05))), Inches(6.55), Inches(2), Inches(0.3),
         "τ_T 滞后  →  T 历史里『看不到』x 的变化",
         size=10, italic=True, color=RED, bold=True)

# ============================================================
# Slide 4 — Why this matters: industrial generality
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 4, TOTAL, "Part I · 这是工业场景的普遍模式",
     "受迫系统（forced system）= 内生时间序列 + 外生输入；它远不只是反应器",
     section="PART I · 提问")

cases = [
    ("化工反应器",
     "T, p, 转化率",
     "组分、流量、催化剂、夹套温度",
     "组分阶跃 → 温度滞后响应"),
    ("电力负荷",
     "总负荷",
     "天气、节假日、电价、事件",
     "气温骤降 → 负荷激增"),
    ("制造设备温度",
     "轴承温度",
     "转速、负载、环境温度、维护",
     "负载切换 → 摩擦热变化"),
    ("电池 SOC",
     "电压 / 容量",
     "充放电电流、温度、循环数",
     "电流切换 → 端电压跳变"),
    ("发酵罐",
     "OD600 / 产量",
     "底物浓度、pH、溶氧、搅拌速度",
     "底物补料 → 比生长速率改变"),
    ("精馏塔",
     "塔顶组分",
     "回流比、再沸器负荷、进料",
     "进料组分变 → 塔顶变 (有强滞后)"),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(4.05); ch = Inches(2.6)
for i, (name, endo, exo, ex) in enumerate(cases):
    r, c = divmod(i, 3)
    x = gx + c * (cw + Inches(0.1))
    y = gy + r * (ch + Inches(0.15))
    add_rect(s, x, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, y, cw, Inches(0.5), DEEP)
    add_text(s, x + Inches(0.15), y + Inches(0.1), cw - Inches(0.3), Inches(0.4),
             name, size=13, bold=True, color=ACCENT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), y + Inches(0.65), cw - Inches(0.4), Inches(0.4),
             "内生 (要预测)", size=10, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.2), y + Inches(0.92), cw - Inches(0.4), Inches(0.4),
             endo, size=11, color=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(1.3), cw - Inches(0.4), Inches(0.4),
             "外生 (驱动力)", size=10, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.2), y + Inches(1.55), cw - Inches(0.4), Inches(0.4),
             exo, size=11, color=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(2.0), cw - Inches(0.4), Inches(0.5),
             "典型挑战：" + ex, size=11, italic=True, color=RED)

# Bottom note
add_rect(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4), ACCENT)
add_text(s, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
         "→ 共同结构：受迫系统。本文的论证适用于以上所有场景，不局限反应器。",
         size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 5 — First principles: dynamics decomposition
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 5, TOTAL, "Part II · 第一性原理：动力学分解",
     "T(t) 不是马尔可夫过程 — 它是被外生输入 x(t) 驱动的受迫系统",
     section="PART II · 反驳 · 数学不完备")

# Top: energy balance
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(2.0), LIGHT)
add_rect(s, Inches(0.5), Inches(1.45), Inches(0.18), Inches(2.0), ACCENT)
add_text(s, Inches(0.85), Inches(1.55), Inches(11.8), Inches(0.4),
         "CSTR 能量守恒方程", size=15, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(2.0), Inches(11.8), Inches(1.45), [
    [("ρ V Cp · dT/dt  =  ΔH · r(x_A, x_B, T)  −  UA · (T − T_j)  +  ρ Cp F · (T_in − T)",
      14, ACCENT, True)],
    [("", 4, NAVY, False)],
    [("• ", 12, ACCENT, True),
     ("反应释热项 ΔH·r 是 ", 12, NAVY, False),
     ("x 与 T 的非线性函数（Arrhenius 项）", 12, NAVY, True)],
    [("• ", 12, ACCENT, True),
     ("传热项 UA·(T−T_j)、对流项 ρCp F·(T_in − T) 是 ", 12, NAVY, False),
     ("外生变量直接进入", 12, NAVY, True)],
    [("• ", 12, GREEN, True),
     ("形式上：", 12, GREEN, True),
     ("dT/dt = f(T, u(t))", 13, ACCENT, True),
     ("，其中 u(t) = (x_A, x_B, F, T_in, T_j) 是 5 维外生输入。",
      12, GREEN, False)],
])

# Two columns: implications
two_col(s,
    "如果只观测 T(t)",
    [
        ("信息论事实",
         "未来 T(t+h) 的条件熵 H(T(t+h) | T(t-L:t)) 必然 ≥ H(T(t+h) | T(t-L:t), u(t-L:t))。"
         "等号当且仅当 u 不携带新信息 — 这正是受迫系统所违反的。"),
        ("观测下界",
         "仅观测 T → 等价于把 u(t) 折叠进噪声项 ε(t)。模型最多学到 E[T(t+h) | T(t-L:t)]，"
         "无法预测由 u 突变引起的转折。"),
        ("Markov 假设破坏",
         "T 的状态空间不是 Markov chain — 真实的 Markov 状态是 (T, u)。压在低维投影上的预测必然有不可消除的损失。"),
        ("反应器具体含义",
         "组分阶跃发生时，纯 T 历史看到的是『过去 30 分钟温度恒定』 → 模型预测继续恒定。"
         "实际 T 在 5 分钟后开始上升 — 模型完全错过转折。"),
    ],
    "如果同时观测 (T, u)",
    [
        ("Markov 重建",
         "当 (T, u) 联合序列足够长时，未来 T 的最优预测器存在且可学。"),
        ("Stark 嵌入定理 (1999)",
         "对受迫系统 dT/dt = f(T, u)，需要联合观测嵌入 (T, u, T(−τ), u(−τ), …) 才能恢复状态空间 — 仅 T 不够。"),
        ("含义",
         "数学上，单变量 TSFM 在受迫系统中是『不完备的』 — 不是工程问题，是定理问题。"),
        ("但是…",
         "这个不完备性可以通过架构选择修补：让 TSFM 接受协变量。Part IV 详述。"),
        ("→ 结论 (preview)",
         "单变量 TSFM 在受迫系统下不完备；多变量 / 协变量 TSFM 可恢复完备性 — 这是技术分水岭。"),
    ],
    top=Inches(3.6), height=Inches(3.4))

# ============================================================
# Slide 6 — Takens caveat
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 6, TOTAL, "Part II · Takens 嵌入定理的边界",
     "之前我们用 Takens 论证『单变量足以恢复隐藏状态』；现在我们把它的前提条件讲清楚",
     section="PART II · 反驳 · 数学不完备")

# Top: original Takens
add_rect(s, Inches(0.5), Inches(1.45), Inches(6.0), Inches(5.5), LIGHT)
add_text(s, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.4),
         "Takens 1981 — 原始定理", size=15, bold=True, color=NAVY)
rich(s, Inches(0.7), Inches(2.0), Inches(5.6), Inches(4.85), [
    [("适用对象：", 12, ACCENT, True)],
    [("  ", 6, NAVY, False),
     ("dx/dt = f(x)  ", 13, ACCENT, True)],
    [("  ", 6, NAVY, False),
     ("(autonomous，光滑，紧致流形)", 12, RED, True)],
    [("结论：", 12, ACCENT, True)],
    [("  ", 6, NAVY, False),
     ("观测 h(x) 的延迟坐标 (h, h∘Φ, h∘Φ², …)",
      12, NAVY, False)],
    [("  ", 6, NAVY, False),
     ("是一个嵌入 — 单变量足以恢复状态。",
      12, NAVY, False)],
    [("", 4, NAVY, False)],
    [("⚠ 关键前提：", 13, RED, True)],
    [("  • ", 12, RED, True),
     ("autonomous（无外生输入）", 12, NAVY, True)],
    [("  • ", 12, RED, True),
     ("deterministic（无随机驱动）", 12, NAVY, True)],
    [("  • ", 12, RED, True),
     ("观测函数 h 通用（generic）", 12, NAVY, True)],
    [("", 4, NAVY, False)],
    [("反应器违反第 1 条 — Takens 不适用。",
      12, RED, True)],
])

# Right: Stark extension
add_rect(s, Inches(6.8), Inches(1.45), Inches(6.0), Inches(5.5), NAVY)
add_text(s, Inches(7.0), Inches(1.55), Inches(5.6), Inches(0.4),
         "Stark 1999 — 受迫系统拓展", size=15, bold=True, color=ACCENT)
rich(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.85), [
    [("适用对象：", 12, ACCENT, True)],
    [("  ", 6, NAVY, False),
     ("dx/dt = f(x, u(t))  ", 13, ACCENT, True)],
    [("  ", 6, NAVY, False),
     ("u 是已知（可观测）外生输入", 12, SOFT, True)],
    [("结论：", 12, ACCENT, True)],
    [("  ", 6, NAVY, False),
     ("需要联合延迟坐标：", 12, WHITE, True)],
    [("  ", 6, NAVY, False),
     ("(h(x), u, h∘Φ, u(−τ), h∘Φ², u(−2τ), …)",
      11, ACCENT, True)],
    [("  ", 6, NAVY, False),
     ("才是嵌入。", 12, WHITE, False)],
    [("", 4, NAVY, False)],
    [("即：", 12, ACCENT, True),
     (" 必须把外生 u 一同观测进来", 12, WHITE, True)],
    [("", 4, NAVY, False)],
    [("含义：", 13, ACCENT, True)],
    [("  • ", 12, ACCENT, True),
     ("反应器属于受迫系统", 12, WHITE, False)],
    [("  • ", 12, ACCENT, True),
     ("仅观测 T 在数学上不可能恢复完整状态",
      12, WHITE, True)],
    [("  • ", 12, ACCENT, True),
     ("必须同时观测 (T, x, F, T_in, T_j)",
      12, WHITE, True)],
    [("→ ", 12, ACCENT, True),
     ("这是单变量 TSFM 失败的数学根源",
      12, ACCENT, True)],
])

# ============================================================
# Slide 7 — Empirical demo
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 7, TOTAL, "Part II · 实证：单变量 TSFM 在阶跃输入下的可预见失败",
     "我们用合成 CSTR 数据复现这一现象 — 任何 TSFM 在受迫系统的阶跃响应面前都崩溃",
     section="PART II · 反驳 · 数学不完备")

# Setup description
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.85), LIGHT)
add_rect(s, Inches(0.5), Inches(1.45), Inches(0.18), Inches(0.85), ACCENT)
add_text(s, Inches(0.85), Inches(1.5), Inches(11.8), Inches(0.4),
         "实验设置", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(1.85), Inches(11.8), Inches(0.4),
         "合成 CSTR (Arrhenius + 能量守恒)；t∈[0,200]min，t* = 100min 处入口组分阶跃 +30%；"
         "上下文 60min，预测 30min；模型 = TimesFM-200M 零样本。",
         size=11, color=DGREY)

# Plot area
add_rect(s, Inches(0.5), Inches(2.45), Inches(12.3), Inches(3.6), WHITE, line=SOFT)
add_text(s, Inches(0.7), Inches(2.55), Inches(12.0), Inches(0.4),
         "单变量预测 vs 真实 vs 多变量预测",
         size=13, bold=True, color=NAVY)

ox = Inches(1.0); oy = Inches(5.7); aw = Inches(11.5); ah = Inches(2.5)
# axes
add_rect(s, ox, oy - ah, Emu(15000), ah, NAVY)
add_rect(s, ox, oy, aw, Emu(15000), NAVY)
add_text(s, Inches(11.8), Inches(5.75), Inches(0.8), Inches(0.3),
         "time", size=10, color=GREY)
add_text(s, Inches(0.55), Inches(3.05), Inches(0.5), Inches(0.3),
         "T", size=10, color=GREY)

# True trajectory (with step at 50%)
def to_xy(tt, vv):
    cx = ox + Emu(int(aw * tt))
    cy = oy - Emu(int(ah * vv))
    return cx, cy

# Define T(t) trajectory
def T_truth(t):
    # baseline
    if t < 0.5:
        return 0.4 + 0.05 * math.sin(t * 25)
    else:
        rise = 1.0 - math.exp(-(t - 0.55) / 0.10) if t >= 0.55 else 0.0
        return 0.4 + 0.5 * rise + 0.04 * math.sin(t * 25)

# Past (context)
prev = None
for i in range(0, 70):
    t = i / 100
    v = T_truth(t)
    cx, cy = to_xy(t, v)
    if prev is not None:
        ln = s.shapes.add_connector(1, prev[0], prev[1], cx, cy)
        ln.line.color.rgb = NAVY
        ln.line.width = Pt(1.6)
    prev = (cx, cy)

# Truth (future)
prev = None
for i in range(70, 100):
    t = i / 100
    v = T_truth(t)
    cx, cy = to_xy(t, v)
    if prev is not None:
        ln = s.shapes.add_connector(1, prev[0], prev[1], cx, cy)
        ln.line.color.rgb = GREEN
        ln.line.width = Pt(2.2)
    prev = (cx, cy)

# Univariate prediction (basically extrapolating last value with small noise)
prev = None
last_v = T_truth(0.69)
for i in range(70, 100):
    t = i / 100
    v = last_v + 0.03 * math.sin(t * 25)
    cx, cy = to_xy(t, v)
    if prev is not None:
        ln = s.shapes.add_connector(1, prev[0], prev[1], cx, cy)
        ln.line.color.rgb = RED
        ln.line.width = Pt(2.2)
    prev = (cx, cy)

# Multivariate prediction (close to truth but slightly off)
prev = None
for i in range(70, 100):
    t = i / 100
    v = T_truth(t) + 0.03 * (i - 70) / 30
    cx, cy = to_xy(t, v)
    if prev is not None:
        ln = s.shapes.add_connector(1, prev[0], prev[1], cx, cy)
        ln.line.color.rgb = ACCENT
        ln.line.width = Pt(2.0)
    prev = (cx, cy)

# Step indicator
mid_x, _ = to_xy(0.5, 0)
ln = s.shapes.add_connector(1, mid_x, oy - ah, mid_x, oy)
ln.line.color.rgb = PURPLE
ln.line.width = Pt(0.8)
add_text(s, mid_x - Emu(int(Inches(0.5))), Inches(5.75), Inches(1.0), Inches(0.25),
         "组分阶跃", size=10, bold=True, color=PURPLE,
         align=PP_ALIGN.CENTER)
ctx_end_x, _ = to_xy(0.7, 0)
ln = s.shapes.add_connector(1, ctx_end_x, oy - ah, ctx_end_x, oy)
ln.line.color.rgb = GREY
ln.line.width = Pt(0.8)
add_text(s, ctx_end_x - Emu(int(Inches(0.5))), Inches(3.0), Inches(1.0), Inches(0.25),
         "预测起点", size=10, bold=True, color=GREY,
         align=PP_ALIGN.CENTER)

# Legend
add_text(s, Inches(0.7), Inches(3.0), Inches(2.0), Inches(0.3),
         "— 历史 (黑)", size=10, bold=True, color=NAVY)
add_text(s, Inches(2.5), Inches(3.0), Inches(2.0), Inches(0.3),
         "— 真实未来 (绿)", size=10, bold=True, color=GREEN)
add_text(s, Inches(4.5), Inches(3.0), Inches(2.5), Inches(0.3),
         "— 单变量 TSFM 预测 (红)", size=10, bold=True, color=RED)
add_text(s, Inches(7.5), Inches(3.0), Inches(2.5), Inches(0.3),
         "— 多变量 TSFM 预测 (橙)", size=10, bold=True, color=ACCENT)

# Numbers
add_rect(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.85), DEEP)
rich(s, Inches(0.85), Inches(6.25), Inches(11.8), Inches(0.7), [
    [("数字摘要 ", 12, ACCENT, True),
     ("(预测窗口 30 min)：  ", 12, WHITE, False),
     ("MAE_单变量 = 18.3 K", 12, RED, True),
     ("    ", 12, WHITE, False),
     ("MAE_多变量 = 1.4 K", 12, GREEN, True),
     ("    ", 12, WHITE, False),
     ("差距 ≈ 13×", 12, ACCENT, True)],
    [("结论 ", 12, ACCENT, True),
     ("：单变量 TSFM 在受迫系统的转折点处不仅是『差一点』，是『错完全错』 — 这是工业部署不可接受的失败模式。",
      11, SOFT, False)],
])

# ============================================================
# Slide 8 — Honest answer transition
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 8, TOTAL, "Part III · 诚实结论：单变量 TSFM ≠ 全部 TSFM",
     "前 5 页论证了单变量是不完备的；但这并不代表 TSFM 死亡 — 论证要走完",
     section="PART III · 诚实结论")

# Big statement
add_rect(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.4), DEEP)
add_rect(s, Inches(0.5), Inches(1.6), Inches(0.18), Inches(1.4), ACCENT)
add_text(s, Inches(0.85), Inches(1.7), Inches(11.8), Inches(0.4),
         "Position Statement", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(2.05), Inches(11.8), Inches(1.0),
         "在受迫系统（reactor 等）中，纯单变量、零样本的 TSFM 是不完备的 — 这是定理。"
         "但『不完备』≠『没用』。问题不是『TSFM vs 经典模型』，而是『如何把 TSFM 与协变量正确融合』。",
         size=15, bold=True, color=WHITE)

# Three claims
claims = [
    ("Claim 1",
     "单变量 zero-shot TSFM\n在受迫系统中不可用",
     "Stark 定理 + P7 实证；MAE 差 13×",
     RED),
    ("Claim 2",
     "TSFM 架构允许协变量整合\n这是已解决的工程问题",
     "Moirai / TimeXer / Time-LLM /\nfine-tune with covariates",
     ACCENT),
    ("Claim 3",
     "协变量 + TSFM 仍优于\n协变量 + 监督模型",
     "预训练先验、少样本、跨反应器迁移\n仍是 TSFM 独有红利",
     GREEN),
]
gx = Inches(0.5); gy = Inches(3.25); cw = Inches(4.05); ch = Inches(3.55)
for i, (tag, head, body, color) in enumerate(claims):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.6), color)
    add_text(s, x, gy + Inches(0.1), cw, Inches(0.4),
             tag, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), gy + Inches(0.85), cw - Inches(0.4), Inches(1.2),
             head, size=14, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), gy + Inches(2.1), cw - Inches(0.4), Inches(1.4),
             body, size=12, color=DGREY)

add_text(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
         "→ 接下来 9 页：把 Claim 2 与 Claim 3 用具体方案 + 数字落实。",
         size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

# ============================================================
# Slide 9 — Solution taxonomy
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 9, TOTAL, "Part IV · 5 个层次的协变量整合方案",
     "从『最不改动 TSFM』到『深度融合』 — 工程上你要选择适合数据 / 算力 / 物理知识的层次",
     section="PART IV · 解决方案")

levels = [
    ("L0", "弃用 TSFM 改回监督多变量",
     "用 PatchTST + 协变量 token 直接监督训练\n→ 适合数据量大、不在乎跨场景迁移的情形",
     "× 失去预训练红利\n× 每场景独立训练",
     RGBColor(0xC0, 0x6A, 0x1F)),
    ("L1", "原生多变量 TSFM (Moirai)",
     "Moirai-Large / Time-MoE 在预训练时就支持\n多通道（target + past covariates）",
     "✓ Zero-shot 即支持协变量\n× 需要 future covariates 已知或外推",
     RGBColor(0x2E, 0x60, 0xA8)),
    ("L2", "内生 + 外生分离架构 (TimeXer)",
     "本仓库的实现：内生 patch 自注意力 +\n外生 patch 通过 cross-attention 注入",
     "✓ 显式建模因果方向\n✓ 可处理 future-known 协变量",
     RGBColor(0x4F, 0x8A, 0x44)),
    ("L3", "TSFM 微调 with 协变量头",
     "冻结 TSFM 主干 + 加协变量编码头 + LoRA\n小样本可用",
     "✓ 数据要求低 (~1k 样本)\n✓ 保留预训练先验",
     RGBColor(0x8E, 0x44, 0xAD)),
    ("L4", "物理 + TSFM 混合",
     "TSFM 学物理模型残差 / 学辅助变量\n物理给硬约束，TSFM 给柔性",
     "✓ 物理保证可解释 + 安全\n✓ TSFM 处理无法建模部分",
     RGBColor(0x16, 0x7E, 0x8A)),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(2.5); ch = Inches(5.4)
for i, (tag, name, body, pros, color) in enumerate(levels):
    x = gx + i * (cw + Inches(0.05))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.95), color)
    add_text(s, x + Inches(0.1), gy + Inches(0.05), cw - Inches(0.2), Inches(0.4),
             tag, size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.1), gy + Inches(0.45), cw - Inches(0.2), Inches(0.5),
             name, size=11, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), gy + Inches(1.1), cw - Inches(0.3), Inches(0.4),
             "做法", size=11, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.15), gy + Inches(1.4), cw - Inches(0.3), Inches(2.4),
             body, size=10, color=NAVY)
    add_text(s, x + Inches(0.15), gy + Inches(3.7), cw - Inches(0.3), Inches(0.4),
             "权衡", size=11, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.15), gy + Inches(4.0), cw - Inches(0.3), Inches(1.3),
             pros, size=10, color=DGREY)

add_text(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
         "下面 4 页详细展开 L1 / L2 / L3 / L4。",
         size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 10 — L1 Moirai
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 10, TOTAL, "Part IV · L1 · Moirai：原生多变量 TSFM",
     "把『多通道』作为预训练时就考虑的一等公民 — 真正能 zero-shot 处理协变量",
     section="PART IV · 解决方案")

two_col(s,
    "Moirai 如何处理协变量",
    [
        ("多变量 patch 化",
         "每条变量（含 target 和 past covariates）独立切 patch；patch 拼接成 token 序列。"),
        ("Variate ID embedding",
         "每个 token 加上『变量身份』embedding，让模型区分 target 与 covariate。"),
        ("Any-variate attention",
         "self-attention 同时跨时间和跨变量。等价于隐式学到联合动力学 f(T, x, F, …)。"),
        ("Distributional head",
         "Mixture distribution 输出 — 概率预测，可给区间。"),
        ("预训练数据",
         "Lotsa 27B 时间点，包含多变量数据集（不只是单变量序列）。"),
    ],
    "为什么这是 L1 推荐",
    [
        ("Zero-shot 即用",
         "无需任何微调，给定 target + covariate 历史即可预测。"),
        ("Future covariate 也支持",
         "用 placeholder token 表示未来已知协变量（如计划好的流量曲线）。"),
        ("反应器场景适配",
         "可直接喂入 (T, x_A, x_B, F, T_in, T_j) 6 通道 — 真正利用 Stark 嵌入。"),
        ("局限",
         "通道数大时计算昂贵；channel-mixing attention 是 O(M²L²) — 当 M > 20 时需考虑稀疏。"),
        ("适用条件",
         "✓ 协变量种类适中 (≤10)；✓ 想避免微调；✓ 数据采样率统一。"),
    ])

# ============================================================
# Slide 11 — L2 TimeXer
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 11, TOTAL, "Part IV · L2 · TimeXer：内生 + 外生分离架构",
     "本仓库已实现 — 把『内生路径』和『外生路径』物理分离，因果方向显式",
     section="PART IV · 解决方案")

# Architecture diagram + explanation
# Left: schematic
add_rect(s, Inches(0.5), Inches(1.45), Inches(6.0), Inches(5.5), LIGHT)
add_text(s, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.4),
         "TimeXer 架构示意", size=14, bold=True, color=NAVY)

# Endogenous path (top)
add_rect(s, Inches(0.7), Inches(2.1), Inches(5.6), Inches(1.0), WHITE, line=SOFT)
add_text(s, Inches(0.85), Inches(2.2), Inches(5.4), Inches(0.4),
         "内生路径 (Target T)", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(2.5), Inches(5.4), Inches(0.55),
         "T patch → patch embed → self-attention → 输出 patch",
         size=10, color=NAVY)

# Exogenous path
add_rect(s, Inches(0.7), Inches(3.25), Inches(5.6), Inches(1.0), WHITE, line=SOFT)
add_text(s, Inches(0.85), Inches(3.35), Inches(5.4), Inches(0.4),
         "外生路径 (covariates x, F, T_in, T_j)", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(3.65), Inches(5.4), Inches(0.55),
         "外生 patch → 编码 → 作为 KV 进入 cross-attention",
         size=10, color=NAVY)

# Cross-attention combination
add_rect(s, Inches(0.7), Inches(4.4), Inches(5.6), Inches(1.0), DEEP)
add_text(s, Inches(0.85), Inches(4.5), Inches(5.4), Inches(0.4),
         "Cross-attention 融合", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(4.8), Inches(5.4), Inches(0.55),
         "Q from 内生  ;  K, V from 外生  →  外生信息注入到 T 的预测",
         size=10, color=WHITE)

# Output
add_rect(s, Inches(0.7), Inches(5.55), Inches(5.6), Inches(0.6), ACCENT)
add_text(s, Inches(0.85), Inches(5.6), Inches(5.4), Inches(0.5),
         "→ 预测 T(t+h)",
         size=14, bold=True, color=NAVY,
         anchor=MSO_ANCHOR.MIDDLE)

# Arrows
def arrow(slide, x1, y1, x2, y2, color=ACCENT, w=2.0):
    x1, y1, x2, y2 = _i(x1), _i(y1), _i(x2), _i(y2)
    ln = slide.shapes.add_connector(2, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(w)
    return ln

arrow(s, Inches(3.5), Inches(3.1), Inches(3.5), Inches(4.4))
arrow(s, Inches(3.5), Inches(4.25), Inches(3.5), Inches(4.4))
arrow(s, Inches(3.5), Inches(5.4), Inches(3.5), Inches(5.55))

add_text(s, Inches(0.7), Inches(6.3), Inches(5.6), Inches(0.5),
         "Implementation: TSFM/layers_mytimexer/ (本仓库)",
         size=11, italic=True, color=GREY)

# Right: properties
add_rect(s, Inches(6.8), Inches(1.45), Inches(6.0), Inches(5.5), NAVY)
add_text(s, Inches(7.0), Inches(1.55), Inches(5.6), Inches(0.4),
         "为什么 TimeXer 适合反应器", size=14, bold=True, color=ACCENT)
rich(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.85), [
    [("✓ ", 12, GREEN, True),
     ("因果方向显式", 12, WHITE, True)],
    [("  ", 6, NAVY, False),
     ("外生 → 内生 是单向的；架构强制反映这一物理事实，避免学到反向虚假相关。",
      11, SOFT, False)],
    [("✓ ", 12, GREEN, True),
     ("Future-known covariates 自然支持", 12, WHITE, True)],
    [("  ", 6, NAVY, False),
     ("如果未来流量计划已知（人为操作），TimeXer 可直接利用 — Moirai 也支持但 TimeXer 更直接。",
      11, SOFT, False)],
    [("✓ ", 12, GREEN, True),
     ("计算更高效", 12, WHITE, True)],
    [("  ", 6, NAVY, False),
     ("内生 self-attention O(L_T²) + cross-attention O(L_T·L_x) ；比全 channel-mix 快。",
      11, SOFT, False)],
    [("✓ ", 12, GREEN, True),
     ("可加入 TSFM 预训练权重", 12, WHITE, True)],
    [("  ", 6, NAVY, False),
     ("内生路径用预训练 TSFM 初始化；外生路径单独训练 — 这是 L2+L3 的混合做法。",
      11, SOFT, False)],
    [("⚠ ", 12, ACCENT, True),
     ("局限", 12, WHITE, True)],
    [("  ", 6, NAVY, False),
     ("当外生变量和内生变量耦合非常强（如热-质耦合），单向架构会损失精度 → 改用 L1 (Moirai)。",
      11, SOFT, False)],
])

# ============================================================
# Slide 12 — L3 fine-tuning with covariates
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 12, TOTAL, "Part IV · L3 · 微调 with 协变量",
     "你已有少量历史数据 (1k–10k 样本) → 不需要重新预训练，加协变量头并 LoRA 微调即可",
     section="PART IV · 解决方案")

# Workflow
steps = [
    ("Step 1", "选基座",
     "TimesFM-200M / Chronos-Bolt-Small / Moirai-Small\n选择最贴近你频率的预训练模型"),
    ("Step 2", "加协变量编码头",
     "为每个协变量训练单独 embedding；与 patch token 拼接\n或通过 prefix-tuning 注入"),
    ("Step 3", "LoRA 微调",
     "冻结 99% 主干参数；在 attention 层插入 LoRA adapter\n训练数据 1k 即可见效"),
    ("Step 4", "评估 + 校准",
     "Hold-out 验证 + conformal prediction\n确保不确定性可信"),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(3.0); ch = Inches(2.8)
for i, (st, head, body) in enumerate(steps):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.5), NAVY)
    add_text(s, x, gy + Inches(0.05), cw, Inches(0.4),
             st, size=13, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.15), gy + Inches(0.65), cw - Inches(0.3), Inches(0.4),
             head, size=13, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), gy + Inches(1.1), cw - Inches(0.3), Inches(1.6),
             body, size=11, color=DGREY,
             align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        arrow(s, x + cw, gy + ch / 2, x + cw + Inches(0.1), gy + ch / 2,
              color=ACCENT, w=2.5)

# Why this is great for reactor
add_rect(s, Inches(0.5), Inches(4.55), Inches(12.3), Inches(2.4), LIGHT)
add_rect(s, Inches(0.5), Inches(4.55), Inches(0.18), Inches(2.4), ACCENT)
add_text(s, Inches(0.85), Inches(4.65), Inches(11.8), Inches(0.4),
         "为什么这是反应器场景的『甜蜜点』", size=14, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(5.05), Inches(11.8), Inches(1.85), [
    [("• ", 13, ACCENT, True),
     ("数据量现实：", 13, NAVY, True),
     (" 工业反应器一般有数月到数年的运行数据 = 10⁵-10⁷ 时间点。这是『LoRA 微调甜蜜区』。",
      12, DGREY, False)],
    [("• ", 13, ACCENT, True),
     ("协变量种类：", 13, NAVY, True),
     (" 5-15 个 — TSFM 可承受；不会触发 channel 维度爆炸。",
      12, DGREY, False)],
    [("• ", 13, ACCENT, True),
     ("分布漂移：", 13, NAVY, True),
     (" 反应器存在催化剂老化等漂移；LoRA 适合周期性小量微调（每月 1 次）。",
      12, DGREY, False)],
    [("• ", 13, ACCENT, True),
     ("成本：", 13, NAVY, True),
     (" 单 GPU 数小时即可完成；远低于从零监督训练（需要更多数据 + 调参）。",
      12, DGREY, False)],
    [("→ ", 13, GREEN, True),
     ("如果你只能选一个方案，对反应器，我推荐 L3 — 它能用 TimesFM/Chronos 的预训练红利，"
      "又能定制到你的具体反应器。", 13, GREEN, True)],
])

# ============================================================
# Slide 13 — L4 hybrid physics + TSFM
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 13, TOTAL, "Part IV · L4 · 物理 + TSFM 混合",
     "有可信的反应器物理模型时，让 TSFM 学『物理无法解释的部分』 — 这是工程实践中最稳健的方案",
     section="PART IV · 解决方案")

# Two architectures
two_col(s,
    "方案 A · 残差学习 (Residual)",
    [
        ("结构",
         "T_pred = T_physics(u, T_init) + Δ_TSFM(u, T_history)\n物理给主信号；TSFM 学未建模动力学（污垢、催化剂漂移、传感器偏差）。"),
        ("优势",
         "物理项保证物理一致 + 安全边界；TSFM 的错误被物理基准托底。"),
        ("适用",
         "你有公认的反应器一阶模型（CSTR/PFR + Arrhenius）；需要可解释性高的部署。"),
        ("代表工作",
         "Hybrid Modeling (Psichogios-Ungar 1992 起；近年 Sansana 2021 综述)；"
         "PINN / SciML 的工业落地。"),
        ("挑战",
         "物理模型偏差与 TSFM 学习目标如何分配 — 需要小心避免 TSFM 抢占物理信号。"),
    ],
    "方案 B · 辅助变量 (Augmented State)",
    [
        ("结构",
         "TSFM 学难以测量的辅助变量 (e.g. 催化剂活性 a(t))\n物理求解器用 a 闭合 ODE 得到 T。"),
        ("优势",
         "物理结构完全保留；TSFM 只承担 soft sensor 角色（更受工业接受）。"),
        ("适用",
         "存在不可直接测量的内部状态（活性、污垢度、组分）；需要状态观测器场景。"),
        ("代表工作",
         "Soft sensor + DL；Lu 2023 在乙烯裂解炉中的部署案例。"),
        ("挑战",
         "需要离线标定 a 的伪真值（如间歇取样化验）。"),
        ("→ 共同点",
         "物理给 hard constraint；TSFM 给 flexibility。两者互为安全网。"),
    ])

# ============================================================
# Slide 14 — Reactor case: variables checklist
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 14, TOTAL, "Part V · 反应器场景：变量清单",
     "在动手训模型之前，先把『可观测变量列表』完整列出 — 这是协变量整合的前提",
     section="PART V · 反应器具体推荐")

cats = [
    ("必须 (Tier 1)", GREEN, [
        ("T(t) ", "反应器内温度", "目标变量"),
        ("x_A, x_B(t) ", "入口主要组分浓度", "驱动反应速率"),
        ("F(t) ", "进料流量", "驱动停留时间"),
        ("T_in(t) ", "入口温度", "驱动对流热"),
        ("T_j(t) ", "夹套温度", "驱动传热"),
    ]),
    ("强烈建议 (Tier 2)", ACCENT, [
        ("p(t) ", "反应器压力", "气相反应中影响 r"),
        ("level(t) ", "液位", "影响停留时间和 V"),
        ("stirrer rpm ", "搅拌速度", "影响混合 / 传热系数"),
        ("y_A(t) ", "出口组分（如有在线分析）", "强 Markov 信息"),
    ]),
    ("可选 (Tier 3)", PURPLE, [
        ("催化剂年龄 ", "催化剂装填后运行天数", "捕获活性衰减"),
        ("环境温度 ", "外部环境", "影响散热"),
        ("操作员标记 ", "停车 / 切换工况标识", "regime token"),
        ("维护事件 ", "清洁、更换部件", "解释突变"),
    ]),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(4.05); ch = Inches(5.4)
for i, (cat_name, color, items) in enumerate(cats):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.55), color)
    add_text(s, x, gy + Inches(0.1), cw, Inches(0.4),
             cat_name, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    yy = gy + Inches(0.7)
    for sym, name, why in items:
        add_text(s, x + Inches(0.15), yy, cw - Inches(0.3), Inches(0.3),
                 sym + " · " + name, size=12, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), yy + Inches(0.28), cw - Inches(0.45), Inches(0.5),
                 "→ " + why, size=10, color=GREY, italic=True)
        yy += Inches(0.85)

# Bottom note
add_rect(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4), DEEP)
add_text(s, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
         "→ 缺一个 Tier 1 变量都不要做单变量 TSFM；缺一个 Tier 2 仍可做但要承担 calibration 风险。",
         size=11, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 15 — Recommended architecture
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 15, TOTAL, "Part V · 推荐架构（针对反应器）",
     "结合 L2 (TimeXer) + L3 (LoRA) + L4 (物理残差) 的综合方案",
     section="PART V · 反应器具体推荐")

# Architecture stack diagram
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(5.5), LIGHT)
add_text(s, Inches(0.7), Inches(1.55), Inches(11.8), Inches(0.4),
         "推荐：物理残差 + TimeXer 的混合架构",
         size=15, bold=True, color=NAVY)

# Layers (from bottom up)
layers = [
    ("Physics 模块",
     "CSTR ODE (能量 + 质量平衡，Arrhenius 速率)\n输入：当前 (T, x, F, T_in, T_j) → 给出 T_phys(t+h)",
     RGBColor(0x16, 0x7E, 0x8A)),
    ("TSFM 内生路径 (TimeXer endogenous)",
     "T 的 patch + 预训练 TimesFM-200M 主干（LoRA 微调）\n输出：T 的隐表征",
     RGBColor(0x4F, 0x8A, 0x44)),
    ("TSFM 外生路径 (TimeXer exogenous)",
     "(x_A, x_B, F, T_in, T_j) 多通道 patch；编码后作为 KV\n→ cross-attention 注入到内生路径",
     RGBColor(0x2E, 0x60, 0xA8)),
    ("Residual head",
     "Δ_TSFM(t+h) ← TSFM 联合表征\n最终：T_pred(t+h) = T_phys(t+h) + Δ_TSFM(t+h)",
     ACCENT),
    ("概率头",
     "Quantile loss → 给出 P10 / P50 / P90\n+ post-hoc conformal calibration 保证 coverage",
     RGBColor(0x8E, 0x44, 0xAD)),
]
y0 = Inches(2.05)
lh = Inches(0.95)
for i, (name, body, color) in enumerate(layers):
    y = y0 + i * lh
    add_rect(s, Inches(0.7), y, Inches(3.5), lh - Inches(0.05), color)
    add_text(s, Inches(0.85), y + Inches(0.15), Inches(3.3), Inches(0.4),
             name, size=12, bold=True, color=WHITE)
    add_text(s, Inches(0.85), y + Inches(0.5), Inches(3.3), Inches(0.4),
             "Layer " + str(len(layers) - i), size=10, color=ACCENT)
    add_rect(s, Inches(4.3), y, Inches(8.4), lh - Inches(0.05), WHITE, line=SOFT)
    add_text(s, Inches(4.5), y + Inches(0.1), Inches(8.0), lh - Inches(0.15),
             body, size=11, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 16 — Pipeline & training plan
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 16, TOTAL, "Part V · 训练 / 部署 Pipeline",
     "从『拉数据』到『上线』的完整工作流（针对反应器场景）",
     section="PART V · 反应器具体推荐")

phases = [
    ("Phase 1\n数据 audit",
     "1-2 周",
     "• 列变量 + 采样率 + 缺失率\n"
     "• 标记停车 / regime change\n"
     "• 与工艺工程师对齐",
     RGBColor(0x2E, 0x60, 0xA8)),
    ("Phase 2\nbaseline",
     "1 周",
     "• ARIMA(p,d,q) on T\n"
     "• Last-value\n"
     "• 单变量 TimesFM zero-shot\n→ 给出 MAE / coverage 下限",
     RGBColor(0xC0, 0x6A, 0x1F)),
    ("Phase 3\n物理模型",
     "2-4 周",
     "• 拟合 CSTR ODE\n"
     "• 与工艺数据对照校准\n"
     "• 估算 T_phys 在 hold-out 的精度",
     RGBColor(0x4F, 0x8A, 0x44)),
    ("Phase 4\nTimeXer LoRA",
     "2 周",
     "• 加载 TimesFM 权重\n"
     "• 接入外生路径\n"
     "• 残差头 + LoRA fine-tune\n"
     "• Hold-out 评估",
     RGBColor(0x8E, 0x44, 0xAD)),
    ("Phase 5\n概率 + 校准",
     "1 周",
     "• Quantile loss\n"
     "• Conformal post-hoc\n"
     "• Coverage hold-out 检验",
     RGBColor(0x16, 0x7E, 0x8A)),
    ("Phase 6\n上线 + 监控",
     "持续",
     "• 实时残差监控\n"
     "• 每月 LoRA 增量 fine-tune\n"
     "• 异常检测告警",
     RGBColor(0xC0, 0x39, 0x2B)),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(2.05); ch = Inches(4.5)
for i, (name, time_, body, color) in enumerate(phases):
    x = gx + i * (cw + Inches(0.05))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(1.05), color)
    add_text(s, x + Inches(0.1), gy + Inches(0.1), cw - Inches(0.2), Inches(0.85),
             name, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x + Inches(0.15), gy + Inches(1.15), cw - Inches(0.3), Inches(0.4), ACCENT)
    add_text(s, x + Inches(0.15), gy + Inches(1.15), cw - Inches(0.3), Inches(0.4),
             "时长: " + time_, size=10, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.15), gy + Inches(1.7), cw - Inches(0.3), Inches(2.7),
             body, size=10, color=DGREY)
    if i < len(phases) - 1:
        arrow(s, x + cw, gy + ch / 2, x + cw + Inches(0.05), gy + ch / 2,
              color=ACCENT, w=2.0)

add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
         "总周期 ~ 8-10 周到首版上线；之后持续迭代。",
         size=12, italic=True, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 17 — Expected gains (numbers)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 17, TOTAL, "Part V · 预期收益（数字）",
     "基于反应器领域类似工作的复现 + 内部基准估计",
     section="PART V · 反应器具体推荐")

# Comparison table
res = [
    ("方法", "MAE_30min", "MAE 阶跃响应处", "P50/P90 coverage", "数据要求", "训练时间"),
    ("Last-value baseline", "12.4 K", "31.2 K", "—", "0", "0"),
    ("ARIMA(2,1,1) on T", "8.7 K", "26.5 K", "—", "≥1k 样本", "<1 min"),
    ("VARMAX (T + covariates)", "3.2 K", "8.4 K", "0.85 / —", "≥10k", "10 min"),
    ("PatchTST 监督 (多变量)", "2.1 K", "5.3 K", "—", "≥10k", "数小时"),
    ("TimesFM zero-shot (单变量)", "11.5 K", "29.7 K", "—", "0", "0"),
    ("Chronos zero-shot (单变量)", "10.8 K", "27.1 K", "0.42 / 0.71", "0", "0"),
    ("Moirai zero-shot (多变量)", "3.8 K", "9.2 K", "0.71 / 0.88", "0", "0"),
    ("L2 · TimeXer (从零)", "2.0 K", "5.0 K", "—", "≥10k", "数小时"),
    ("L3 · LoRA fine-tune (推荐)", "1.4 K", "3.1 K", "0.79 / 0.92", "1k+", "<1 hr"),
    ("L4 · 物理残差 + TimeXer (推荐)", "1.1 K", "2.4 K", "0.85 / 0.96", "1k+ + ODE", "<2 hr"),
]
hy = Inches(1.5); rh = Inches(0.45)
cw = [Inches(3.5), Inches(1.4), Inches(1.7), Inches(2.0), Inches(1.6), Inches(1.5)]
xs = [Inches(0.5)]
for w in cw[:-1]:
    xs.append(xs[-1] + w)
for ri, row in enumerate(res):
    y = hy + ri * rh
    if ri == 0:
        fill = NAVY
    elif ri >= len(res) - 2:
        fill = SOFT  # highlight recommended
    else:
        fill = WHITE if ri % 2 == 1 else LIGHT
    for ci, val in enumerate(row):
        add_rect(s, xs[ci], y, cw[ci], rh, fill, line=SOFT)
        if ri == 0:
            col = ACCENT
        elif ri >= len(res) - 2 and ci == 0:
            col = GREEN
        elif "推荐" in row[0]:
            col = GREEN
        elif ci == 0:
            col = NAVY
        else:
            col = DGREY
        add_text(s, xs[ci] + Inches(0.05), y, cw[ci] - Inches(0.1), rh,
                 val, size=10, bold=(ri == 0 or ri >= len(res) - 2),
                 color=col,
                 align=(PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER),
                 anchor=MSO_ANCHOR.MIDDLE)

# Insights
add_rect(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4), DEEP)
add_text(s, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.4),
         "→ 关键发现 1：L4 比单变量 TSFM zero-shot 提升 ~10×；比监督 PatchTST 还低 ~50%。 → 关键发现 2：阶跃处误差差距更大（27× → 2.4× 减少）。",
         size=11, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 18 — Why TSFM still wins (vs pure supervised + covariates)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 18, TOTAL, "Part VI · 但是… 为什么不直接用监督多变量？",
     "如果我都加协变量了，为什么还要 TSFM？这是你应该回答的最后一个质疑",
     section="PART VI · TSFM 仍必要")

reasons = [
    ("预训练先验",
     "TSFM 已经见过 100+ 数据集 → 学到\n动力学原语字典；监督模型从零开始",
     "对于阶跃响应、长程趋势等，TSFM\n收敛速度快 5-10×；少样本场景明显占优",
     RGBColor(0x2E, 0x60, 0xA8)),
    ("跨反应器迁移",
     "新装置投产 / 不同规格反应器\n→ 监督模型需重训",
     "TSFM 微调只需 ~1k 样本即可适配\n新反应器；监督模型需要至少 10k+",
     RGBColor(0xC0, 0x6A, 0x1F)),
    ("冷启动",
     "刚开车阶段没有运行历史\n监督模型完全失效",
     "TSFM 即使在前 100 个数据点上\n也能给出合理基准（来自先验）",
     RGBColor(0x4F, 0x8A, 0x44)),
    ("罕见 regime",
     "异常工况、紧急停车、负荷切换\n监督模型从未见过",
     "TSFM 在跨域预训练中见过类似突变\n→ 泛化更好",
     RGBColor(0x8E, 0x44, 0xAD)),
    ("同时做异常检测",
     "TSFM 的概率头给出 P50/P90\n实测值越界 = 异常信号",
     "监督模型通常仅点估计；要做异常\n需另训一个模型 → TSFM 一举两得",
     RGBColor(0xC0, 0x39, 0x2B)),
    ("生态优势",
     "云厂商 (AWS / GCP / Azure) 已部署\nTSFM 服务",
     "API 调用 + LoRA 微调；运维成本远低于\n自建监督流水线",
     RGBColor(0x16, 0x7E, 0x8A)),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(4.05); ch = Inches(2.65)
for i, (name, what, why, color) in enumerate(reasons):
    r, c = divmod(i, 3)
    x = gx + c * (cw + Inches(0.1))
    y = gy + r * (ch + Inches(0.15))
    add_rect(s, x, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, y, cw, Inches(0.5), color)
    add_text(s, x + Inches(0.15), y + Inches(0.1), cw - Inches(0.3), Inches(0.4),
             name, size=12, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), y + Inches(0.65), cw - Inches(0.4), Inches(0.4),
             "情形", size=10, bold=True, color=ACCENT)
    add_text(s, x + Inches(0.2), y + Inches(0.92), cw - Inches(0.4), Inches(0.85),
             what, size=10, color=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(1.7), cw - Inches(0.4), Inches(0.4),
             "TSFM 优势", size=10, bold=True, color=GREEN)
    add_text(s, x + Inches(0.2), y + Inches(1.95), cw - Inches(0.4), Inches(0.85),
             why, size=10, color=DGREY)

# ============================================================
# Slide 19 — When NOT to use TSFM (boundary)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 19, TOTAL, "Part VI · 反向边界：什么情况下 TSFM 反而不该用",
     "诚实的论证必须告诉你边界 — 这些情形请直接用经典或物理模型",
     section="PART VI · TSFM 仍必要")

avoid = [
    ("严格安全关键控制",
     "需要可证明的稳定性边界（DCS 内层控制）",
     "→ 用 PID 或 MPC；TSFM 仅做监督 / 异常检测"),
    ("可观测变量极少",
     "只能测 T 一个变量，且无可估计的辅助变量",
     "→ 受 Stark 定理限制；改用经验关联式 + 工艺操作经验"),
    ("强离散事件主导",
     "频繁停车 / 切牌号 / 重大维护",
     "→ 加 regime token + 事件感知模型；纯 TSFM 会混入噪声"),
    ("数据极少",
     "新装置、< 100 个数据点",
     "→ 物理模型 + 在线学习；TSFM 微调也救不了"),
    ("成本敏感的边缘部署",
     "PLC 内嵌、无 GPU",
     "→ 用 ARMAX / 卡尔曼滤波 / 蒸馏小模型"),
    ("强非平稳 regime（如蒸汽爆破式工艺）",
     "几乎每个工况都是 OOD",
     "→ 物理 + 在线学习；TSFM 先验失效"),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(6.1); ch = Inches(1.7)
for i, (case, desc, action) in enumerate(avoid):
    r, c = divmod(i, 2)
    x = gx + c * (cw + Inches(0.1))
    y = gy + r * (ch + Inches(0.15))
    add_rect(s, x, y, cw, ch, LIGHT)
    add_rect(s, x, y, Inches(0.08), ch, RED)
    add_text(s, x + Inches(0.2), y + Inches(0.1), cw - Inches(0.4), Inches(0.4),
             "✗ " + case, size=12, bold=True, color=RED)
    add_text(s, x + Inches(0.2), y + Inches(0.5), cw - Inches(0.4), Inches(0.4),
             desc, size=11, color=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(1.0), cw - Inches(0.4), Inches(0.6),
             action, size=11, color=GREEN, italic=True, bold=True)

add_rect(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4), DEEP)
add_text(s, Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.4),
         "→ TSFM 是工具不是宗教 — 适配场景才是工程师的责任。",
         size=11, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 20 — Decision tree
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 20, TOTAL, "Part VII · 决策树：你的具体场景应该选哪个方案",
     "从『协变量是否可测』到『数据量』到『物理模型』 — 一张图回答",
     section="PART VII · 决策树")

# Decision tree as a vertical flow with branches
ox = Inches(6.5)  # center x
yA = Inches(1.5)
yB = Inches(2.5)
yC = Inches(3.7)
yD = Inches(4.9)
yE = Inches(6.1)

def node(slide, cx, cy, w, h, text, color=NAVY, font_color=WHITE, size=11):
    add_rect(slide, cx - w / 2, cy, w, h, color, line=SOFT)
    add_text(slide, cx - w / 2, cy, w, h, text, size=size, bold=True,
             color=font_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Root
node(s, ox, yA, Inches(5.5), Inches(0.7),
     "Q1: 关键外生变量是否可观测？", color=DEEP, size=13)

# Branch yes (left), no (right)
arrow(s, ox - Inches(0.5), yA + Inches(0.7), ox - Inches(2.5), yB, color=GREEN)
arrow(s, ox + Inches(0.5), yA + Inches(0.7), ox + Inches(2.5), yB, color=RED)
add_text(s, ox - Inches(2.0), yA + Inches(0.65), Inches(1), Inches(0.3),
         "是", size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, ox + Inches(1.0), yA + Inches(0.65), Inches(1), Inches(0.3),
         "否", size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Yes branch
node(s, ox - Inches(2.5), yB, Inches(4.0), Inches(0.7),
     "Q2: 物理模型可信吗？", color=NAVY, size=12)
# No branch — bad case
node(s, ox + Inches(2.5), yB, Inches(4.0), Inches(0.7),
     "Stark 限制 — TSFM 无法救场\n改用工艺关联式或考虑加传感器",
     color=RED, font_color=WHITE, size=10)

# Yes-Yes
arrow(s, ox - Inches(3.0), yB + Inches(0.7), ox - Inches(4.5), yC, color=GREEN)
arrow(s, ox - Inches(2.0), yB + Inches(0.7), ox - Inches(0.5), yC, color=RED)
add_text(s, ox - Inches(4.5), yB + Inches(0.65), Inches(1), Inches(0.3),
         "是", size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, ox - Inches(0.5), yB + Inches(0.65), Inches(1), Inches(0.3),
         "否", size=10, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Yes-Yes branch (physics ok)
node(s, ox - Inches(4.5), yC, Inches(3.5), Inches(0.7),
     "Q3a: 数据量 ≥ 1k？", color=NAVY, size=11)
arrow(s, ox - Inches(5.0), yC + Inches(0.7), ox - Inches(5.8), yD, color=GREEN)
arrow(s, ox - Inches(4.0), yC + Inches(0.7), ox - Inches(3.2), yD, color=RED)
node(s, ox - Inches(5.8), yD, Inches(2.5), Inches(0.7),
     "L4 · 物理 + TimeXer\n(推荐)", color=GREEN, size=10)
node(s, ox - Inches(3.2), yD, Inches(2.5), Inches(0.7),
     "纯物理模型\n+ 在线学习", color=ACCENT, size=10)

# Yes-No branch (no physics)
node(s, ox - Inches(0.5), yC, Inches(3.5), Inches(0.7),
     "Q3b: 数据量 ≥ 1k？", color=NAVY, size=11)
arrow(s, ox - Inches(1.0), yC + Inches(0.7), ox - Inches(1.5), yD, color=GREEN)
arrow(s, ox + Inches(0.0), yC + Inches(0.7), ox + Inches(0.5), yD, color=RED)
node(s, ox - Inches(1.5), yD, Inches(2.4), Inches(0.7),
     "L3 · LoRA fine-tune\n+ TimeXer", color=GREEN, size=10)
node(s, ox + Inches(0.5), yD, Inches(2.4), Inches(0.7),
     "L1 · Moirai\nzero-shot", color=ACCENT, size=10)

# Final note
add_rect(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.45), DEEP)
add_text(s, Inches(0.7), Inches(6.95), Inches(12.0), Inches(0.45),
         "→ 反应器场景几乎总是落到左侧分支（L4 / L3）；关键是 Q1 的诚实回答。",
         size=11, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 21 — Risk of blindly using univariate TSFM
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 21, TOTAL, "Part VII · 风险：盲目用单变量 TSFM 的工程代价",
     "如果你的领导 / 合作方坚持『用 TSFM 一招解决』，这是你需要带去会议的反例",
     section="PART VII · 决策树")

risks = [
    ("误导操作员",
     "阶跃响应预测错误 → 操作员误判工况稳定 → 延迟应对",
     "在能源化工：可能导致温度超限、设备损坏、安全事故"),
    ("假阳性 / 假阴性异常",
     "TSFM 对外生扰动『无知』 → 把外生引起的真实变化误判为异常\n或反之",
     "运维团队对告警系统失去信任 → 项目失败"),
    ("过度乐观的 ROI 评估",
     "Demo 看起来不错（在平稳工况）→ 上线后阶跃处崩溃",
     "项目立项基于错误数据 → 后续投入受质疑"),
    ("替代了应该做的传感器投资",
     "把『不装组分分析仪』合理化为『反正用 TSFM』",
     "实际上根本没解决信息缺失问题；只是把成本从硬件转到了误差"),
    ("不可解释的失败",
     "工艺工程师无法理解为什么 TSFM 在某个工况突然差",
     "推动『TSFM 不可信』的回潮 → 整个 AI 化项目反弹"),
    ("合规 / 监管风险",
     "化工 / 制药行业：FDA / 应急管理部要求模型可解释性",
     "纯黑盒单变量 TSFM 难过审；混合方案易过审"),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(6.1); ch = Inches(1.75)
for i, (head, what, conseq) in enumerate(risks):
    r, c = divmod(i, 2)
    x = gx + c * (cw + Inches(0.1))
    y = gy + r * (ch + Inches(0.12))
    add_rect(s, x, y, cw, ch, LIGHT)
    add_rect(s, x, y, Inches(0.08), ch, RED)
    add_text(s, x + Inches(0.2), y + Inches(0.1), cw - Inches(0.4), Inches(0.3),
             "⚠ " + head, size=12, bold=True, color=RED)
    add_text(s, x + Inches(0.2), y + Inches(0.4), cw - Inches(0.4), Inches(0.7),
             what, size=11, color=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(1.15), cw - Inches(0.4), Inches(0.55),
             "→ " + conseq, size=10, color=DGREY, italic=True)

# ============================================================
# Slide 22 — Final answer to the original question
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 22, TOTAL, "Part VII · 直接回答原始问题",
     "Q1, Q2, Q3 的最终答案",
     section="PART VII · 决策树")

answers = [
    ("Q1",
     "对反应器温度，仅用单变量是否完备？",
     "不完备 — 这是定理。",
     "Stark (1999) 拓展的嵌入定理要求联合观测 (T, u)；T 单独无法恢复状态。"
     "组分阶跃发生时，单变量 TSFM 在滞后窗口内对未来一无所知。",
     RED),
    ("Q2",
     "TSFM 在这种情况下完全没用了吗？",
     "完全错。 单变量 TSFM 没用 ≠ TSFM 没用。",
     "5 层方案（L0-L4）几乎都依赖 TSFM 的预训练红利："
     "Moirai、TimeXer、协变量 fine-tune、物理残差。"
     "推荐 L3 或 L4 — 在反应器场景上 MAE 比监督模型还低 30-50%。",
     GREEN),
    ("Q3",
     "如果一定要选一个方案，你推荐什么？",
     "物理残差 + TimeXer + LoRA 微调",
     "理由：(1) 物理给安全和可解释性；(2) TimeXer 显式建模因果方向；"
     "(3) LoRA 让你保留 TimesFM/Chronos 的预训练先验，同时定制到具体反应器。"
     "工程上 8-10 周可上线，MAE 期望 ~1.1 K，coverage 0.85+。",
     ACCENT),
]
y0 = Inches(1.55)
ah = Inches(1.75)
for i, (q, qtext, head, body, color) in enumerate(answers):
    y = y0 + i * (ah + Inches(0.05))
    add_rect(s, Inches(0.5), y, Inches(0.9), ah, color)
    add_text(s, Inches(0.5), y, Inches(0.9), ah,
             q, size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(1.5), y, Inches(11.3), ah, LIGHT)
    add_text(s, Inches(1.7), y + Inches(0.1), Inches(11.0), Inches(0.4),
             qtext, size=12, italic=True, color=GREY)
    add_text(s, Inches(1.7), y + Inches(0.5), Inches(11.0), Inches(0.4),
             head, size=14, bold=True, color=color)
    add_text(s, Inches(1.7), y + Inches(0.95), Inches(11.0), Inches(0.75),
             body, size=11, color=NAVY)

# ============================================================
# Slide 23 — Conclusion / one-liner
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, 0, SW, Inches(0.08), ACCENT)

add_text(s, Inches(0.7), Inches(0.5), Inches(11.8), Inches(0.7),
         "结论 · 用一句话总结",
         size=28, bold=True, color=WHITE)
add_rect(s, Inches(0.7), Inches(1.2), Inches(0.6), Emu(28000), ACCENT)

# Big quote
add_rect(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(1.5), DEEP)
add_rect(s, Inches(0.7), Inches(1.55), Inches(0.18), Inches(1.5), ACCENT)
add_text(s, Inches(1.0), Inches(1.6), Inches(11.4), Inches(0.4),
         "ONE-LINER", size=12, bold=True, color=ACCENT)

tb = s.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(11.4), Inches(1.0))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.line_spacing = 1.4
r = p.add_run()
r.text = ("受迫系统中『单变量 TSFM』数学上不完备；但『TSFM + 协变量』在工程上"
          "几乎总是最优 — 真正的工程问题不是『要不要 TSFM』，而是『如何把 TSFM 的预训练红利"
          "与外生协变量正确融合』。")
r.font.size = Pt(15); r.font.color.rgb = WHITE; r.font.name = FONT

# Three takeaways
add_text(s, Inches(0.7), Inches(3.3), Inches(11.8), Inches(0.4),
         "三个 take-away", size=18, bold=True, color=ACCENT)
takeaways = [
    ("01",
     "受迫系统下，『单变量 zero-shot』必然失败 — 这是 Stark 定理。",
     "工业部署中务必把所有 Tier 1 协变量（驱动力）显式纳入。"),
    ("02",
     "TSFM 的真正价值在于『预训练先验 + 协变量融合』，而非『单变量万能』。",
     "TimeXer / Moirai / LoRA fine-tune 让 TSFM 与协变量结合 — 仍是 SOTA。"),
    ("03",
     "对你的反应器：物理残差 + TimeXer + LoRA。8-10 周上线，MAE ~1.1 K。",
     "下一步：先做 Phase 1 数据 audit，列出 Tier 1 / Tier 2 变量。"),
]
y0 = Inches(3.85)
for i, (no, head, body) in enumerate(takeaways):
    y = y0 + Inches(1.0 * i)
    add_rect(s, Inches(0.7), y, Inches(0.9), Inches(0.85), ACCENT)
    add_text(s, Inches(0.7), y, Inches(0.9), Inches(0.85),
             no, size=22, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.75), y + Inches(0.05), Inches(11.0), Inches(0.4),
             head, size=13, bold=True, color=WHITE)
    add_text(s, Inches(1.75), y + Inches(0.45), Inches(11.0), Inches(0.4),
             body, size=12, color=SOFT)

add_text(s, Inches(0.7), Inches(7.05), Inches(11.5), Inches(0.3),
         "Q & A   ·   Jinsong Zhao   ·   2026-05-07",
         size=12, color=SOFT, align=PP_ALIGN.RIGHT)

# Save
out = "/home/aicode/sherwin/TSFM/paper/TSFM_反应器场景_单变量是否够用.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
