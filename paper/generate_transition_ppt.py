"""TSFM for multi-condition transition forecasting & early warning.

Grounded in the ZJSH PI10102 OOD diagnostics:
  - target Δx KS rejection FDR=18.84% — near i.i.d. across 24 transitions
  - controller OPs KS 75-87% — strongly regime-dependent
  - ADF+KPSS unanimously stationary on Δx for all 24 transitions
  - Hurst 0.4-0.55, VR(60) 0.2-0.4 (mean-reverting at long horizons)
  - Permutation entropy ≈ 0.98+ (high ordinal randomness)

Core argument:
  - Three sub-tasks must be separated:
      A. Within-transition trajectory forecast
      B. Transition early warning (leading prediction)
      C. Forecasting in a never-seen-before regime
  - For A:  univariate TSFM is the sweet spot (marginal i.i.d. + ICL)
  - For B:  univariate is insufficient — controller OPs are leading indicators
  - For C:  TSFM helps but needs OOD detection layer
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import math

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
DEEP = RGBColor(0x14, 0x2C, 0x52)
ACCENT = RGBColor(0xE8, 0x8B, 0x1A)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
GREY = RGBColor(0x4A, 0x55, 0x68)
DGREY = RGBColor(0x33, 0x3A, 0x47)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xDC, 0xE3, 0xF0)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)
PURPLE = RGBColor(0x6B, 0x3F, 0xA0)
TEAL = RGBColor(0x16, 0x7E, 0x8A)
ORANGE = RGBColor(0xE0, 0x6E, 0x18)
GOLD = RGBColor(0xD4, 0xA0, 0x2A)

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _i(v):
    return int(v) if v is not None else v


def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_oval(slide, x, y, w, h, fill, line=None):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def rich(slide, x, y, w, h, blocks, line_spacing=1.3, anchor=MSO_ANCHOR.TOP):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, runs in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        for txt, size, color, bold in runs:
            r = p.add_run()
            r.text = txt
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return tb


def line(slide, x1, y1, x2, y2, color, w=1.8):
    x1, y1, x2, y2 = _i(x1), _i(y1), _i(x2), _i(y2)
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(w)
    return ln


def arrow(slide, x1, y1, x2, y2, color=ACCENT, w=2.0):
    x1, y1, x2, y2 = _i(x1), _i(y1), _i(x2), _i(y2)
    ln = slide.shapes.add_connector(2, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(w)
    return ln


def page(slide, idx, total, title, subtitle=None, section=None):
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_rect(slide, 0, 0, SW, Inches(0.08), ACCENT)
    add_rect(slide, 0, Inches(0.08), Inches(0.18), SH, NAVY)
    if section:
        add_text(slide, Inches(0.5), Inches(0.18), Inches(11.5), Inches(0.3),
                 section, size=11, bold=True, color=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.6),
             title, size=22, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(12.0), Inches(0.4),
                 subtitle, size=12, color=GREY, italic=True)
    add_rect(slide, Inches(0.5), Inches(1.28), Inches(0.6), Emu(28000), ACCENT)
    add_text(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
             f"{idx} / {total}", size=10, color=GREY, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.3),
             "TSFM 在多工况过渡过程中的潜力 · ZJSH 实证 · Jinsong Zhao",
             size=10, color=GREY)


TOTAL = 22

# ============================================================
# Slide 1 — Title
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, 0, SW, Inches(0.15), ACCENT)
for i in range(8):
    add_rect(s, Inches(0.5 + i * 1.6), Inches(0.4), Inches(0.6), Emu(15000), ACCENT)

add_text(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.7),
         "TSFM 在多工况过渡过程预测预警中的潜力",
         size=36, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.15), Inches(11.5), Inches(0.5),
         "Univariate or Multivariate? — A ZJSH-grounded analysis",
         size=18, color=ACCENT, italic=True)

add_rect(s, Inches(0.8), Inches(2.95), Inches(0.7), Emu(30000), ACCENT)

rich(s, Inches(0.8), Inches(3.15), Inches(11.5), Inches(3.4), [
    [("基于 ", 16, SOFT, False),
     ("ZJSH 数据集 PI10102", 16, ACCENT, True),
     (" — 16 变量, 524k 样本, 24 个过渡 — 的 OOD 实证分析",
      16, SOFT, False)],
    [("", 8, SOFT, False)],
    [("核心结论 (一句话):", 16, ACCENT, True)],
    [("过渡过程的预测预警 ", 14, SOFT, False),
     ("不是单一任务", 14, ACCENT, True),
     (", 必须分解为 ", 14, SOFT, False),
     ("3 个子任务", 14, ACCENT, True),
     (" — TSFM 在每个子任务上的合理范式 ", 14, SOFT, False),
     ("不一样", 14, ACCENT, True), (".", 14, SOFT, False)],
    [("", 6, SOFT, False)],
    [("• ", 14, ACCENT, True),
     ("过渡内部轨迹预测", 14, WHITE, True),
     (": ", 14, SOFT, False),
     ("✓ 单变量 TSFM 可能最优 (Δ 边际近似 i.i.d.)",
      14, GREEN, True)],
    [("• ", 14, ACCENT, True),
     ("过渡早期预警 (leading indicator)", 14, WHITE, True),
     (": ", 14, SOFT, False),
     ("✗ 必须 multivariate (controller OP)",
      14, RED, True)],
    [("• ", 14, ACCENT, True),
     ("新工况首次出现 (extreme OOD)", 14, WHITE, True),
     (": ", 14, SOFT, False),
     ("△ TSFM 有韧性, 但需 OOD 检测层", 14, ORANGE, True)],
])

add_rect(s, Inches(0.8), Inches(6.7), Inches(11.5), Emu(15000), ACCENT)
add_text(s, Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.4),
         "Jinsong Zhao    ·    清华大学    ·    2026-05-08",
         size=14, color=ACCENT)

# ============================================================
# Slide 2 — Question reformulation
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 2, TOTAL, "把问题精确化",
     "『TSFM 在多工况过渡过程预测预警中是否有潜力 — 单变量是否够?』",
     section="PART I · 问题")

add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.6), LIGHT)
add_rect(s, Inches(0.5), Inches(1.55), Inches(0.18), Inches(1.6), ACCENT)
add_text(s, Inches(0.85), Inches(1.65), Inches(11.8), Inches(0.4),
         "用户问题的隐含三元组", size=14, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(2.05), Inches(11.8), Inches(1.05), [
    [("Q.1   多工况过渡过程下, TSFM 是否有潜力?",
      13, NAVY, True)],
    [("Q.2   单变量 TSFM (univariate-to-univariate) 这种『TSFM 主流模式』是否更合适?",
      13, NAVY, True)],
    [("Q.3   什么时候够、什么时候不够?",
      13, NAVY, True)],
])

# Hidden ambiguity: "prediction" is not one task
add_rect(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.55), DEEP)
add_text(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.55),
         "关键观察: 『过渡过程预测预警』不是一个任务, 而是 至少三个 不同任务",
         size=14, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

tasks = [
    ("Task A · 内部轨迹",
     "已经进入了过渡, 预测下一段轨迹",
     "输入: 过渡早期 60-300 步\n输出: 之后 30-120 步轨迹\n问: TSFM 单变量能否抓住?",
     ACCENT),
    ("Task B · 早期预警",
     "稳态时, 预警过渡 即将 发生",
     "输入: 稳态运行历史\n输出: 概率: 未来 X 步内会发生过渡?\n问: 谁带来 leading 信号?",
     RED),
    ("Task C · 新工况",
     "进入历史从未见过的工况",
     "输入: 训练数据中没出现过的目标态\n输出: 该工况下的轨迹预测\n问: TSFM 是否仍可信?",
     PURPLE),
]
gx = Inches(0.5); gy = Inches(4.05); cw = Inches(4.05); ch = Inches(2.95)
for i, (title, sub, body, col) in enumerate(tasks):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.55), col)
    add_text(s, x, gy, cw, Inches(0.55), title, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), gy + Inches(0.7), cw - Inches(0.4), Inches(0.5),
             sub, size=12, italic=True, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), gy + Inches(1.3), cw - Inches(0.4), Inches(1.6),
             body, size=11, color=DGREY)

# ============================================================
# Slide 3 — ZJSH dataset overview
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 3, TOTAL, "ZJSH 数据集 · 我们用什么来回答这个问题",
     "目标变量 PI10102 + 15 个外生变量 + 524k 时间点 + 24 个识别出的过渡",
     section="PART II · ZJSH 实证")

add_rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(2.6), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.55), DEEP)
add_text(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.55),
         "数据规模", size=14, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(0.7), Inches(2.25), Inches(5.7), Inches(1.85), [
    [("• ", 14, ACCENT, True), ("时间点数: 524,140 (秒级 ~ 6 天连续)",
      12, NAVY, True)],
    [("• ", 14, ACCENT, True), ("变量数: 16",
      12, NAVY, True)],
    [("    包含: 流量 FIC × 7 (PV+OP), 温度 TIC, 压力 PI, 液位 LI",
      11, GREY, False)],
    [("• ", 14, ACCENT, True),
     ("识别出的过渡: 24 个 (BGMM regime + Mahalanobis 距离)",
      12, NAVY, True)],
    [("• ", 14, ACCENT, True),
     ("目标变量: PI10102 (压力指示)",
      12, NAVY, True)],
])

add_rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(2.6), WHITE, line=SOFT)
add_rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.55), DEEP)
add_text(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.55),
         "变量类型", size=14, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(7.0), Inches(2.25), Inches(5.7), Inches(1.85), [
    [("• ", 14, ACCENT, True), ("PV (Process Value):",
      12, NAVY, True), (" 被动测量, 反映系统真实状态",
      12, DGREY, False)],
    [("    ", 12, NAVY, False), ("FIC*PV, TIC*PV, LI20503, PI10102",
      11, GREY, False)],
    [("• ", 14, ACCENT, True),
     ("OP (Output / 控制器输出):", 12, NAVY, True),
     (" 主动控制信号", 12, DGREY, False)],
    [("    ", 12, NAVY, False),
     ("FIC*OP, TIC*OP — 这些是 ", 11, GREY, False),
     ("人为施加的扰动", 11, RED, True)],
    [("→ ", 12, ACCENT, True),
     ("OP 是 leading 信号源, PV 是 lagging 响应",
      12, NAVY, True)],
])

# Bottom: BGMM-detected transitions visualization
add_rect(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7), LIGHT)
add_text(s, Inches(0.85), Inches(4.45), Inches(11.8), Inches(0.4),
         "过渡检测原理 (BGMM + Mahalanobis tail)",
         size=14, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(4.85), Inches(11.8), Inches(2.0), [
    [("Step 1   ", 12, ACCENT, True),
     ("BGMM (变分贝叶斯高斯混合) 在 16 维变量空间识别 regime / 工况簇",
      12, DGREY, False)],
    [("Step 2   ", 12, ACCENT, True),
     ("regime 标签的 change-point + Mahalanobis 距离 > P97 标记『候选过渡区间』",
      12, DGREY, False)],
    [("Step 3   ", 12, ACCENT, True),
     ("中值滤波 (W=301) 后合并近邻区间, 过滤过短段 (< 301)",
      12, DGREY, False)],
    [("Step 4   ", 12, ACCENT, True),
     ("只保留两端 regime 不同的区间 ⇒ 24 个真过渡, 长度 903 - 13006 步",
      12, DGREY, False)],
    [("→ ", 12, GREEN, True),
     ("这给我们一个『跨工况受控数据集』, 可以严格比较『同一目标变量在 24 种工况切换中的统计行为』",
      12, NAVY, True)],
])

# ============================================================
# Slide 4 — Headline OOD finding 1: target Δ marginal
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 4, TOTAL, "OOD 证据 1 · 目标变量 Δ-空间 跨工况近似 i.i.d.",
     "PI10102 的 Δx 边际分布 — 跨 24 个过渡几乎不可分",
     section="PART II · ZJSH 实证")

# big number panel
add_rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(2.6), DEEP)
add_text(s, Inches(0.7), Inches(1.7), Inches(5.6), Inches(0.4),
         "Pairwise KS 检验 · PI10102 的 Δx",
         size=13, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(2.2), Inches(5.6), Inches(0.5),
         "n_pairs = 276 (24 transitions × C(2,2))",
         size=11, italic=True, color=SOFT)

# Big number
add_text(s, Inches(0.7), Inches(2.85), Inches(5.6), Inches(0.6),
         "FDR-BH 拒绝率: ",
         size=14, color=WHITE)
add_text(s, Inches(2.45), Inches(2.65), Inches(3.8), Inches(0.95),
         "18.84%",
         size=44, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(3.65), Inches(5.6), Inches(0.4),
         "(原始 raw=27.54%; 5% 显著性下随机基线 ~5%)",
         size=11, italic=True, color=SOFT)

add_rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(2.6), WHITE, line=SOFT)
add_rect(s, Inches(6.8), Inches(1.55), Inches(0.18), Inches(2.6), ACCENT)
add_text(s, Inches(7.15), Inches(1.7), Inches(5.6), Inches(0.4),
         "解读", size=14, bold=True, color=NAVY)
rich(s, Inches(7.15), Inches(2.1), Inches(5.6), Inches(2.0), [
    [("• ", 12, ACCENT, True),
     ("FDR=18.84% ", 12, NAVY, True),
     ("意味着只有 ~19% 的过渡对在 Δ 边际上 ",
      11, DGREY, False),
     ("可被区分", 11, NAVY, True), (".", 11, DGREY, False)],
    [("• 大部分过渡对的 Δx 边际看起来 ",
      11, DGREY, False),
     ("来自同一分布", 11, NAVY, True), (".", 11, DGREY, False)],
    [("• 这是 TSFM 的 ", 11, DGREY, False),
     ("『梦境数据』", 11, ACCENT, True),
     (" — RevIN 假设的边际不变性近似成立.",
      11, DGREY, False)],
    [("• 类比: 你看 24 段不同的语音, 词频分布几乎相同 ⇒ TSFM 的语言模型直接通用.",
      11, GREY, False)],
])

# Supporting evidence — multi-test consistency
add_rect(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7), LIGHT)
add_text(s, Inches(0.85), Inches(4.45), Inches(11.8), Inches(0.4),
         "其他独立检验的一致性证据 (PI10102 在 Δ 空间)",
         size=14, bold=True, color=NAVY)

# Two columns of evidences
rich(s, Inches(0.85), Inches(4.9), Inches(5.8), Inches(2.0), [
    [("• ", 13, ACCENT, True),
     ("ADF + KPSS:", 13, NAVY, True),
     (" 24/24 工况 都判定 stationary",
      12, DGREY, False)],
    [("    (ADF p < 0.05  AND  KPSS p > 0.05 一致)",
      11, GREY, False)],
    [("• ", 13, ACCENT, True),
     ("Permutation Entropy:", 13, NAVY, True),
     (" PE(m=4) 全部 ≥ 0.98",
      12, DGREY, False)],
    [("    序列在『序号模式』上接近随机 ⇒ 复杂度 saturated",
      11, GREY, False)],
    [("• ", 13, ACCENT, True),
     ("Sign-change rate:", 13, NAVY, True),
     (" 0.59 - 0.65 across all 24",
      12, DGREY, False)],
    [("    随机游走基线为 0.5; 略高反映 mild mean-reversion",
      11, GREY, False)],
])
rich(s, Inches(7.0), Inches(4.9), Inches(5.8), Inches(2.0), [
    [("• ", 13, ACCENT, True),
     ("Hurst 指数 H:", 13, NAVY, True),
     (" 多数 0.40-0.50 (random walk / 弱反持续)",
      12, DGREY, False)],
    [("    极端: T2/T3/T23 微 persistent (H ≈ 0.6)",
      11, GREY, False)],
    [("• ", 13, ACCENT, True),
     ("Variance Ratio VR(60):", 13, NAVY, True),
     (" 多在 0.18 - 0.46",
      12, DGREY, False)],
    [("    显著小于 1 ⇒ 长程反持续 (mean-reverting)",
      11, GREY, False)],
    [("→ ", 13, GREEN, True),
     ("综合判定: 在目标变量 Δ 空间, 跨工况几乎是『同一种动力学』",
      12, NAVY, True)],
])

# ============================================================
# Slide 5 — Headline OOD finding 2: controller OP variables
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 5, TOTAL, "OOD 证据 2 · 控制器 OP 变量则跨工况差异巨大",
     "OP 才是过渡的真正『驱动信号』 — KS 拒绝率高达 75 - 87%",
     section="PART II · ZJSH 实证")

# Bar chart of KS rejection rates
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(3.5), LIGHT)
add_text(s, Inches(0.85), Inches(1.7), Inches(11.8), Inches(0.4),
         "Pairwise KS · FDR-BH 拒绝率排序 (Δx 边际)",
         size=14, bold=True, color=NAVY)

# Bar data: variable, rejection rate, color (red for OP, green for PV/target)
bars = [
    ("FIC10103OP", 86.23, RED),
    ("FIC10105OP", 84.06, RED),
    ("FIC10103",   84.78, RED),
    ("FIC10101OP", 80.07, RED),
    ("FIC10102OP", 78.26, RED),
    ("FIC10102",   76.81, RED),
    ("FIC10104OP", 73.55, RED),
    ("FIC20901OP", 62.32, ORANGE),
    ("FIC10401",   51.09, ORANGE),
    ("FIC10401OP", 46.01, ORANGE),
    ("FIC10402OP", 40.58, ORANGE),
    ("TIC10204OP", 38.41, ORANGE),
    ("TIC10204",   36.59, ORANGE),
    ("LI20503",    19.20, GREEN),
    ("PI10102",    18.84, GREEN),
    ("FIC20901",   13.77, GREEN),
]
gx = Inches(0.85); gy = Inches(2.1)
bar_h = Inches(0.16)
bar_full = Inches(8.5)
for i, (name, val, col) in enumerate(bars):
    y = gy + i * (bar_h + Emu(20000))
    add_text(s, gx, y, Inches(2.0), bar_h, name, size=10, bold=True, color=NAVY,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    bx = gx + Inches(2.1)
    add_rect(s, bx, y, bar_full, bar_h, WHITE, line=SOFT)
    add_rect(s, bx, y, int(bar_full * val / 100), bar_h, col)
    add_text(s, bx + bar_full + Inches(0.1), y, Inches(1.2), bar_h,
             f"{val:.2f}%", size=10, bold=True, color=col,
             anchor=MSO_ANCHOR.MIDDLE)

# Bottom interpretation
add_rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.85), DEEP)
rich(s, Inches(0.85), Inches(5.35), Inches(11.8), Inches(1.65), [
    [("→ ", 14, ACCENT, True),
     ("控制器 OP 在 Δ 空间跨工况 显著不同 (KS FDR 75-87%)",
      14, ACCENT, True)],
    [("    OP 在不同工况下被设定为不同的目标值, 导致 Δ-空间动力学不同 — 这是预期的, 因为 OP 反映 ",
      12, SOFT, False),
     ("人为意图", 12, ACCENT, True), (".", 12, SOFT, False)],
    [("→ ", 14, ACCENT, True),
     ("PV 类变量 (PI10102, LI20503, FIC20901) 跨工况几乎不可区分 (KS FDR 14-19%)",
      14, ACCENT, True)],
    [("    这告诉我们: 物理系统对不同 OP 的响应在 Δ 空间『被压缩』到了相似的统计行为. ",
      12, SOFT, False),
     ("OP 是因, PV 是果", 12, ACCENT, True), (".", 12, SOFT, False)],
])

# ============================================================
# Slide 6 — What this means for TSFM (mechanism-level)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 6, TOTAL, "这两个观察对 TSFM 意味着什么?",
     "把 ZJSH 的统计结构精确映射到 TSFM 的工作前提 — 是否吻合?",
     section="PART II · ZJSH 实证")

# Two columns: TSFM assumption vs ZJSH evidence
items = [
    ("RevIN / InstanceNorm 的『边际不变』前提",
     "TSFM 在每个 instance 内做 z-score; 隐含假设是 不同 instance 的 (z-score 后) 边际近似一致.",
     "ZJSH 证据: PI10102 的 Δ 边际 KS=18.84%, ADF/KPSS 一致 stationary, PE 一致饱和. ",
     "→ 高度吻合: RevIN 是合适的归一化策略.",
     GREEN),
    ("Patching 的『局部模式可复用』前提",
     "假设邻近时间点构成的 patch 的语义 在不同序列间可以共享; sin / 阶跃 / 振荡等 motif 是通用词汇.",
     "ZJSH 证据: PE(m=4)≈0.98, 序号模式分布稳定 ⇒ Δ 空间的 patch 语义跨工况近似复用.",
     "→ 高度吻合: Patching token 在不同过渡间可重用.",
     GREEN),
    ("ICL / Induction Heads 的『context 可推断 regime』前提",
     "假设给定足够长 context, 模型可以从中推断当前任务参数.",
     "ZJSH 证据: 24 个工况下 Hurst 在 0.4-0.6 间小变动, VR(60) 也仅在 0.2-0.5 间, ",
     "→ 部分吻合: regime 信号微弱但存在; ICL 可识别但可能需较长 context.",
     ORANGE),
    ("跨工况共享一种动力学的隐含假设",
     "单变量 TSFM 默认: 同一个 f_θ 适用于所有过渡 / 工况, 无需切换.",
     "ZJSH 证据: PI10102 的 ADF+KPSS 全 stationary, KS 弱拒绝 ⇒ 在 Δ 空间确实近似单一动力学.",
     "→ 高度吻合: 单 model 多工况 是合理的.",
     GREEN),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(12.3); ch = Inches(1.27)
for i, (head, prem, evid, verdict, col) in enumerate(items):
    y = gy + i * (ch + Inches(0.05))
    add_rect(s, gx, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, gx, y, Inches(0.18), ch, col)
    add_text(s, gx + Inches(0.35), y + Inches(0.05), Inches(11.8), Inches(0.32),
             head, size=12, bold=True, color=NAVY)
    add_text(s, gx + Inches(0.35), y + Inches(0.4), Inches(11.8), Inches(0.3),
             "TSFM 前提:  " + prem, size=10, color=GREY)
    add_text(s, gx + Inches(0.35), y + Inches(0.7), Inches(11.8), Inches(0.3),
             "ZJSH 证据:  " + evid, size=10, color=DGREY)
    add_text(s, gx + Inches(0.35), y + Inches(0.97), Inches(11.8), Inches(0.3),
             verdict, size=10, bold=True, color=col)

# ============================================================
# Slide 7 — Headline answer for Task A (within-transition)
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 7, TOTAL, "任务 A · 过渡内部轨迹预测",
     "已经进入过渡, 给定早期 60-300 步, 预测后续 30-120 步",
     section="PART III · 三任务分别讨论")

# Verdict
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.8), GREEN)
add_text(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.8),
         "✓  单变量 TSFM (univariate-to-univariate) 在 ZJSH 上理论上可行, 且很可能最优",
         size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Three reasons
add_rect(s, Inches(0.5), Inches(2.55), Inches(12.3), Inches(1.4), LIGHT)
add_text(s, Inches(0.85), Inches(2.7), Inches(11.8), Inches(0.4),
         "三个理由 — 都来自 ZJSH 实证", size=14, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(3.1), Inches(11.8), Inches(0.85), [
    [("R1 ", 13, ACCENT, True),
     ("Δx 边际跨工况近似 i.i.d. (KS=18.84%) ⇒ TSFM 在所有 24 个过渡上看到的『局部统计』几乎相同",
      12, NAVY, True)],
    [("R2 ", 13, ACCENT, True),
     ("PE+Hurst+VR 一致 ⇒ 局部动力学在 Δ 空间几乎共享 ⇒ 一个 TSFM 模型即可覆盖",
      12, NAVY, True)],
    [("R3 ", 13, ACCENT, True),
     ("ICL 通过 context (过渡早期 60-300 步) 自动判断当前是哪种过渡的中段, ",
      12, NAVY, True),
     ("不需要外部 regime 标签", 12, ACCENT, True)],
])

# Concrete recommendations
add_rect(s, Inches(0.5), Inches(4.1), Inches(6.0), Inches(2.95), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(4.1), Inches(6.0), Inches(0.55), DEEP)
add_text(s, Inches(0.5), Inches(4.1), Inches(6.0), Inches(0.55),
         "推荐做法 (PI10102)", size=14, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(0.7), Inches(4.8), Inches(5.7), Inches(2.2), [
    [("• 模型: ", 12, ACCENT, True),
     ("Timer-XL / TimesFM-200M / Chronos-Bolt-Base",
      12, NAVY, True)],
    [("• 范式: ", 12, ACCENT, True),
     ("zero-shot 优先 → 不行再 LoRA",
      12, NAVY, True)],
    [("• 输入: ", 12, ACCENT, True),
     ("过渡早期 96-512 步的 PI10102",
      12, NAVY, True)],
    [("• 输出: ", 12, ACCENT, True),
     ("未来 24-96 步 + P10/P50/P90",
      12, NAVY, True)],
    [("• 数据: ", 12, ACCENT, True),
     ("以 24 个过渡作为 24 个『同质样本』",
      12, NAVY, True)],
    [("• 校准: ", 12, ACCENT, True),
     ("Conformal post-hoc, hold-out 4 个过渡",
      12, NAVY, True)],
])

add_rect(s, Inches(6.8), Inches(4.1), Inches(6.0), Inches(2.95), WHITE, line=SOFT)
add_rect(s, Inches(6.8), Inches(4.1), Inches(6.0), Inches(0.55), DEEP)
add_text(s, Inches(6.8), Inches(4.1), Inches(6.0), Inches(0.55),
         "为什么单变量在这里『正好够』", size=14, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(7.0), Inches(4.8), Inches(5.7), Inches(2.2), [
    [("• 进入过渡『后』, 上游 OP 变化已经在 PI10102 历史中留下印迹 — ",
      11, DGREY, False)],
    [("  ICL 可以从 context 反向推断 regime.",
      11, DGREY, False)],
    [("• 多变量带来的 ", 11, DGREY, False),
     ("额外信息在 ICL 下被显著稀释",
      11, ACCENT, True), (".", 11, DGREY, False)],
    [("• Channel-Independence + 单变量 ", 11, DGREY, False),
     ("放大有效数据 16×",
      11, ACCENT, True),
     (" (16 个变量都可独立训练).",
      11, DGREY, False)],
    [("• 实证: 在 GIFT-Eval / LOTSA, 单变量 TSFM 在 stationary-Δ 数据上 ",
      11, DGREY, False),
     ("匹配或胜过", 11, GREEN, True),
     (" multivariate.", 11, DGREY, False)],
])

# ============================================================
# Slide 8 — Caveats for Task A
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 8, TOTAL, "任务 A 的边界条件",
     "『单变量够用』不是无条件的 — 列出三种例外",
     section="PART III · 三任务分别讨论")

caveats = [
    ("Caveat 1 · 过渡持续期发生 第二次 OP 干预",
     "如果在我们预测的 horizon 内, controller 又施加了 新的 setpoint 改变,\n这个新干预不在 PI10102 历史中, 单变量 TSFM 无能为力.",
     "→ 解决: 缩短 horizon 到 OP 平均稳定窗口内; 或加 OP_future 预定义信号作为外生输入.",
     RED),
    ("Caveat 2 · 持续过渡 (T2/T3/T23) — Hurst > 0.55",
     "这 3 个过渡的 Δx 在长程上 持续, 不是 mean-reverting.\n单变量 TSFM 在 stationary-Δ 假设下可能产生轨迹漂移.",
     "→ 解决: 这些过渡上加 trend dummy 或用 Moirai/Lag-Llama 等带漂移建模的 backbone.",
     ORANGE),
    ("Caveat 3 · 极短过渡 (< 300 步)",
     "T20 (903 步), T9 (1095 步), T12 (1031 步), T19 (1171 步), T23 (1133 步).\nContext 不足以让 ICL 收敛 ⇒ 模型可能 fall back 到 prior.",
     "→ 解决: 早期检测过渡 + 切换到『短窗口预测模式』 (输入 32-64 步, 输出 8-16 步).",
     PURPLE),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(12.3); ch = Inches(1.78)
for i, (title, body, fix, col) in enumerate(caveats):
    y = gy + i * (ch + Inches(0.1))
    add_rect(s, gx, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, gx, y, Inches(0.2), ch, col)
    add_text(s, gx + Inches(0.4), y + Inches(0.1), Inches(11.7), Inches(0.4),
             title, size=13, bold=True, color=col)
    add_text(s, gx + Inches(0.4), y + Inches(0.55), Inches(11.7), Inches(0.7),
             body, size=11, color=DGREY)
    add_text(s, gx + Inches(0.4), y + Inches(1.32), Inches(11.7), Inches(0.4),
             fix, size=11, italic=True, bold=True, color=NAVY)

add_text(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
         "→ 排除以上三种边角后, 单变量 TSFM 在 ZJSH 内部轨迹任务上是 sweet spot.",
         size=12, italic=True, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 9 — Task B: early warning
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 9, TOTAL, "任务 B · 过渡早期预警 (leading prediction)",
     "稳态运行中, 提前 5-30 分钟报告: 即将发生过渡的概率",
     section="PART III · 三任务分别讨论")

add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.8), RED)
add_text(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.8),
         "✗  单变量 TSFM 在这里 不足以 — 这正好对应『反应器场景』PPT 中的论证",
         size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Why not univariate
add_rect(s, Inches(0.5), Inches(2.55), Inches(6.0), Inches(2.45), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(2.55), Inches(0.18), Inches(2.45), RED)
add_text(s, Inches(0.85), Inches(2.65), Inches(5.6), Inches(0.4),
         "为什么单变量做不到", size=14, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(3.05), Inches(5.6), Inches(1.95), [
    [("• 过渡的 ", 12, DGREY, False),
     ("触发因", 12, RED, True),
     (" = controller OP 的人为变化", 12, DGREY, False)],
    [("• 在 PI10102 (PV) 上, OP 变化的影响 ",
      12, DGREY, False),
     ("有滞后", 12, RED, True),
     (" (若干秒到数分钟).", 12, DGREY, False)],
    [("• 单变量 TSFM 只能看 PI10102 历史 ⇒ 在 OP 改变到 PV 显著响应之间的 ",
      12, DGREY, False),
     ("延迟期", 12, RED, True),
     (", 它没有任何信号.", 12, DGREY, False)],
    [("• 结果: 单变量预警 ", 12, DGREY, False),
     ("最早只能在 PV 已变化时报警", 12, RED, True),
     (" — 这不是『预警』, 是『报告』.",
      12, DGREY, False)],
])

# What works
add_rect(s, Inches(6.8), Inches(2.55), Inches(6.0), Inches(2.45), WHITE, line=SOFT)
add_rect(s, Inches(6.8), Inches(2.55), Inches(0.18), Inches(2.45), GREEN)
add_text(s, Inches(7.15), Inches(2.65), Inches(5.6), Inches(0.4),
         "什么有效", size=14, bold=True, color=NAVY)
rich(s, Inches(7.15), Inches(3.05), Inches(5.6), Inches(1.95), [
    [("• 把 ", 12, DGREY, False),
     ("FIC*OP / TIC*OP", 12, GREEN, True),
     (" 作为外生协变量输入.", 12, DGREY, False)],
    [("• ZJSH 证据: 这些 OP 在 Δ 空间跨工况 ", 12, DGREY, False),
     ("KS 差异 75-87%", 12, GREEN, True),
     (" — 它们 携带 regime 切换的明确信号.",
      12, DGREY, False)],
    [("• 架构: ", 12, DGREY, False),
     ("TimeXer", 12, GREEN, True),
     (" (endogenous + exogenous path) 或 ",
      12, DGREY, False),
     ("Moirai", 12, GREEN, True),
     (" (multivariate, 任意 channel 数).",
      12, DGREY, False)],
    [("• 训练: 有标签的『过渡发生 / 未发生』分类头, 共享 TSFM 主干.",
      12, DGREY, False)],
])

# Mechanism diagram (verbal)
add_rect(s, Inches(0.5), Inches(5.15), Inches(12.3), Inches(1.85), DEEP)
add_text(s, Inches(0.85), Inches(5.3), Inches(11.8), Inches(0.4),
         "因果链 (重要!)", size=14, bold=True, color=ACCENT)
rich(s, Inches(0.85), Inches(5.7), Inches(11.8), Inches(1.25), [
    [("操作员决定 →  ", 13, SOFT, False),
     ("控制器 SP/OP 变化", 13, ACCENT, True),
     ("  →  阀位变化  →  流量 PV 变化  →  ",
      13, SOFT, False),
     ("PI10102 (压力) 响应", 13, ACCENT, True),
     ("  →  ...新稳态.", 13, SOFT, False)],
    [("", 6, SOFT, False)],
    [("→ ", 13, ACCENT, True),
     ("OP 是因果链的 上游, PI10102 是下游. ", 13, WHITE, True),
     ("单变量预测预警 = 只看下游 = 必然滞后.",
      13, ACCENT, True)],
    [("→ ", 13, ACCENT, True),
     ("提前 5-30 分钟预警 = 必须把上游 OP 纳入观测 = ",
      12, SOFT, False),
     ("multivariate 是必要 不是 nice-to-have", 12, ACCENT, True), (".",
      12, SOFT, False)],
])

# ============================================================
# Slide 10 — Task B: leading indicators which to use?
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 10, TOTAL, "任务 B · 哪些变量值得作为 leading indicator?",
     "用 ZJSH 的 KS 差异分数 + 因果分析筛选",
     section="PART III · 三任务分别讨论")

# Tier-based ranking
tiers = [
    ("Tier 1 · 必选", GREEN,
     ["FIC10103OP (KS=86%)", "FIC10105OP (KS=84%)", "FIC10101OP (KS=80%)",
      "FIC10102OP (KS=78%)", "FIC10104OP (KS=74%)"],
     "OP 系列控制器输出 — 跨工况差异最大, 因果链最上游, 物理意义清晰"),
    ("Tier 2 · 推荐",
     ORANGE,
     ["FIC10103 (KS=85%)", "FIC10102 (KS=77%)",
      "FIC20901OP (KS=62%)", "FIC10401 (KS=51%)"],
     "PV / 部分 OP — 跨工况中等差异, 反映系统响应; 与 Tier 1 一起使用"),
    ("Tier 3 · 谨慎",
     PURPLE,
     ["TIC10204OP (KS=38%)", "TIC10204 (KS=37%)",
      "FIC10402OP (KS=41%)", "FIC10401OP (KS=46%)"],
     "温度回路 — KS 中等, 信号相对弱; 视计算预算决定是否纳入"),
    ("Tier 4 · 不必要",
     GREY,
     ["LI20503 (KS=19%)", "PI10102 (KS=19%)", "FIC20901 (KS=14%)"],
     "PV-类, 跨工况几乎不可分; 本身就是要预测的目标"),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(12.3); ch = Inches(1.3)
for i, (title, col, vars_, desc) in enumerate(tiers):
    y = gy + i * (ch + Inches(0.07))
    add_rect(s, gx, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, gx, y, Inches(2.0), ch, col)
    add_text(s, gx, y, Inches(2.0), ch, title, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, gx + Inches(2.2), y + Inches(0.15), Inches(10.0), Inches(0.4),
             "  ·  ".join(vars_), size=11, bold=True, color=NAVY)
    add_text(s, gx + Inches(2.2), y + Inches(0.6), Inches(10.0), Inches(0.6),
             desc, size=11, italic=True, color=DGREY)

add_text(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
         "→ 推荐先用 Tier 1+2 (9 个变量) 作为 multivariate 输入; 后续按贡献度 ablation 取舍.",
         size=12, italic=True, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 11 — Task C: never-seen-before regime
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 11, TOTAL, "任务 C · 真 OOD · 训练数据中未出现过的工况",
     "工厂启动新生产配方 / 设备故障导致工况漂移 / 进料组分变化",
     section="PART III · 三任务分别讨论")

# Verdict
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.8), ORANGE)
add_text(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.8),
         "△  TSFM 比传统监督方法有韧性, 但需要 OOD 检测 + fallback 层",
         size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Three subpoints
points = [
    ("为什么 TSFM 有相对韧性?",
     "TSFM 在 LOTSA / GIFT 上见过 9 大域 × 各种频率 × 各种过渡形态.\n新工况若是『见过的形态加新参数』, ICL 可以拟合.",
     "证据: TSFM zero-shot 在 GIFT-Eval 的 OOD split 平均 MAE 比 PatchTST/TFT 低 8-15%.",
     GREEN),
    ("但不是无限制的『有效』",
     "如果新工况引入了 训练数据中未出现的物理机制 (例: 新副反应, 新反馈回路),\nTSFM 同样 over-extrapolate.",
     "ZJSH 的局限: 我们的 24 个过渡都来自类似工艺; 真正『新』的工况需要外部判断.",
     RED),
    ("正确架构: TSFM + OOD 监控",
     "层次: \n  • 第 1 层 — TSFM 预测 + 概率区间 (P10/P90)\n  • 第 2 层 — 残差监控 + Mahalanobis 距离 (像 ZJSH 检测过渡的方式)\n  • 第 3 层 — 残差超阈值 ⇒ 触发 fallback (人工 / 物理模型 / 保守预测)",
     "→ 这正好把 ZJSH 的 BGMM 检测逻辑用到了在线场景",
     ACCENT),
]
gx = Inches(0.5); gy = Inches(2.55); cw = Inches(12.3); ch = Inches(1.45)
for i, (title, body, foot, col) in enumerate(points):
    y = gy + i * (ch + Inches(0.05))
    add_rect(s, gx, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, gx, y, Inches(0.18), ch, col)
    add_text(s, gx + Inches(0.35), y + Inches(0.1), Inches(11.8), Inches(0.4),
             title, size=13, bold=True, color=col)
    add_text(s, gx + Inches(0.35), y + Inches(0.5), Inches(11.8), Inches(0.65),
             body, size=11, color=DGREY)
    add_text(s, gx + Inches(0.35), y + Inches(1.13), Inches(11.8), Inches(0.3),
             foot, size=10, italic=True, bold=True, color=NAVY)

# ============================================================
# Slide 12 — Three-task summary table
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 12, TOTAL, "三任务总结 · 一张表回答用户的两个问题",
     "TSFM 在哪些子任务有潜力? 单变量在哪些子任务够用?",
     section="PART IV · 综合判断")

# Header
gx = Inches(0.5); gy = Inches(1.55)
hdr_cols = ["任务", "TSFM 是否有潜力", "单变量够吗", "推荐范式", "ZJSH 证据要点"]
hdr_w = [Inches(2.6), Inches(1.9), Inches(1.7), Inches(2.7), Inches(3.4)]
add_rect(s, gx, gy, sum(hdr_w, Inches(0)), Inches(0.6), DEEP)
x_acc = gx
for hd, w in zip(hdr_cols, hdr_w):
    add_text(s, x_acc + Inches(0.1), gy, w - Inches(0.1), Inches(0.6),
             hd, size=12, bold=True, color=ACCENT,
             anchor=MSO_ANCHOR.MIDDLE)
    x_acc += w

# Rows
rows = [
    ("Task A · 内部轨迹预测",
     ("✓ 高", GREEN),
     ("✓ 是", GREEN),
     "单变量 TSFM zero-shot 或 LoRA",
     "Δ 边际 KS 18.84%; ADF/KPSS 全 stationary; ICL 可识别 regime"),
    ("Task B · 过渡早期预警",
     ("✓ 高", GREEN),
     ("✗ 否", RED),
     "Multivariate (TimeXer / Moirai)\n+ Tier 1+2 OP 协变量",
     "OP KS 75-87% 表明 leading 信号集中在 OP; PV 是 lagging"),
    ("Task C · 新工况 (OOD)",
     ("△ 中", ORANGE),
     ("△ 部分", ORANGE),
     "TSFM + 残差监控 + fallback\n(像 BGMM 用到在线)",
     "训练数据中未见过的工况; TSFM 比 supervised 韧性强, 但需 OOD 层"),
]
y = gy + Inches(0.6)
rh = Inches(1.6)
for i, (task, (potential, p_col), (uni, u_col), pard, ev) in enumerate(rows):
    add_rect(s, gx, y, sum(hdr_w, Inches(0)), rh,
             LIGHT if i % 2 == 0 else WHITE, line=SOFT)
    x_acc = gx
    add_text(s, x_acc + Inches(0.1), y + Inches(0.25), hdr_w[0] - Inches(0.1), Inches(1.1),
             task, size=12, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
    x_acc += hdr_w[0]
    add_text(s, x_acc, y + Inches(0.25), hdr_w[1], Inches(1.1),
             potential, size=14, bold=True, color=p_col,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x_acc += hdr_w[1]
    add_text(s, x_acc, y + Inches(0.25), hdr_w[2], Inches(1.1),
             uni, size=14, bold=True, color=u_col,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x_acc += hdr_w[2]
    add_text(s, x_acc + Inches(0.1), y + Inches(0.15), hdr_w[3] - Inches(0.1), rh - Inches(0.3),
             pard, size=11, color=DGREY,
             anchor=MSO_ANCHOR.MIDDLE)
    x_acc += hdr_w[3]
    add_text(s, x_acc + Inches(0.1), y + Inches(0.15), hdr_w[4] - Inches(0.2), rh - Inches(0.3),
             ev, size=10, italic=True, color=GREY,
             anchor=MSO_ANCHOR.MIDDLE)
    y = y + rh

# ============================================================
# Slide 13 — The integrated architecture
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 13, TOTAL, "整合架构 · 同时回应三个任务",
     "一个共享 TSFM backbone, 三种使用模式",
     section="PART V · 推荐架构")

# Architecture diagram (block-based)
# Backbone
add_rect(s, Inches(4.5), Inches(1.55), Inches(4.3), Inches(1.3), DEEP)
add_text(s, Inches(4.5), Inches(1.55), Inches(4.3), Inches(0.5),
         "共享 backbone", size=12, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(4.5), Inches(1.95), Inches(4.3), Inches(0.4),
         "Timer-XL / TimesFM-200M",
         size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(4.5), Inches(2.35), Inches(4.3), Inches(0.4),
         "(预训练权重, 共享 attention/patching)",
         size=10, italic=True, color=SOFT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Three heads
heads = [
    ("Head A · 单变量轨迹",
     "PI10102 univariate\n96-512 → 24-96 步\n+ Conformal CI",
     ACCENT, GREEN, "Task A"),
    ("Head B · 多变量预警",
     "Tier 1+2 OP 输入\nendo: PI10102\nexo: 5-9 个 OP\n→ 概率: P(transition in next k)",
     ACCENT, RED, "Task B"),
    ("Head C · OOD 检测",
     "残差 + Mahalanobis\n阈值: > P95 触发警报\nfallback: 人工 / 物理模型",
     ACCENT, ORANGE, "Task C"),
]
gx = Inches(0.5); gy = Inches(3.4); cw = Inches(4.05); ch = Inches(2.7)
for i, (title, body, hcol, ecol, tag) in enumerate(heads):
    x = gx + i * (cw + Inches(0.1))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(0.55), ecol)
    add_text(s, x, gy, cw, Inches(0.55), title, size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x + cw - Inches(0.85), gy + Inches(0.1), Inches(0.75), Inches(0.32), WHITE)
    add_text(s, x + cw - Inches(0.85), gy + Inches(0.1), Inches(0.75), Inches(0.32),
             tag, size=10, bold=True, color=ecol,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), gy + Inches(0.7), cw - Inches(0.4), Inches(2.0),
             body, size=11, color=DGREY)
    # arrow from backbone
    arrow(s, Inches(6.65), Inches(2.85), x + cw // 2, gy, color=ecol, w=1.8)

# Bottom note
add_rect(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.65), LIGHT)
rich(s, Inches(0.85), Inches(6.45), Inches(11.8), Inches(0.55), [
    [("→ ", 13, ACCENT, True),
     ("好处: 三任务共享 backbone ⇒ 训练效率高 + ", 12, NAVY, True),
     ("内部 representation 一致", 12, ACCENT, True),
     (" ⇒ 减少模型间不一致带来的工程成本.", 12, NAVY, True)],
    [("    LoRA adapter 可针对每个 head 独立微调而不影响其他 head.",
      11, GREY, False)],
])

# ============================================================
# Slide 14 — Empirical-grounded benchmark plan
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 14, TOTAL, "建议的实验对照 · 用 ZJSH 验证三个判断",
     "怎么用现有 24 过渡 + 524k 数据严格证伪 / 证实",
     section="PART VI · 验证方案")

# Three experiments
exps = [
    ("Exp 1 · 验证 Task A 的『单变量够用』",
     "Hold out 6 个过渡 (T2 / T3 / T4 / T17 / T18 / T23 — 含长短和持续性多样).\n比较: (a) Univariate TSFM zero-shot, (b) Multivariate TSFM (Moirai/TimeXer), (c) 监督 Transformer (PatchTST).",
     "假设证伪条件: 若 (b) MAE 比 (a) 低 ≥ 15%, 则单变量不够 — 需重新评估.\n预期: (a) ≈ (b), 都好于 (c).",
     ACCENT),
    ("Exp 2 · 验证 Task B 必须 multivariate",
     "构造预警标签: t_i = 1 如果 30 步内会发生过渡, 否则 0.\n比较 (a) PI10102-only TSFM (b) PI10102 + Tier1 OP (c) PI10102 + Tier1+2.",
     "假设证伪条件: 若 (a) 的 ROC-AUC ≥ (b), 则单变量足够 — 但根据 KS 差异预测应有 ≥ 15-25 点 AUC 提升.\n预期: AUC(a) ≪ AUC(b) < AUC(c).",
     RED),
    ("Exp 3 · 验证 Task C 的 OOD 检测必要性",
     "人工注入: 在 hold-out 过渡上 添加 噪声 / 漂移 模拟新工况 (mean shift 2σ).\n比较: 有 vs 无 残差监控层在 false-positive vs detection-delay 上的权衡.",
     "假设证伪条件: 若纯 TSFM 在注入数据上 MAE 不变 + 不报警, 说明它已经把异常误以为是 prior — 需要监控层.",
     ORANGE),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(12.3); ch = Inches(1.78)
for i, (title, body, falsify, col) in enumerate(exps):
    y = gy + i * (ch + Inches(0.08))
    add_rect(s, gx, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, gx, y, Inches(0.18), ch, col)
    add_text(s, gx + Inches(0.35), y + Inches(0.1), Inches(11.8), Inches(0.4),
             title, size=13, bold=True, color=col)
    add_text(s, gx + Inches(0.35), y + Inches(0.5), Inches(11.8), Inches(0.65),
             body, size=11, color=DGREY)
    add_text(s, gx + Inches(0.35), y + Inches(1.2), Inches(11.8), Inches(0.55),
             falsify, size=11, italic=True, bold=True, color=NAVY)

# ============================================================
# Slide 15 — Why ZJSH is "ideal" data for this question
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 15, TOTAL, "为什么 ZJSH 是回答这个问题的『理想』数据",
     "它同时满足『跨工况 marginal 稳定』和『有明确 leading driver』",
     section="PART VII · 数据的内在含义")

reasons = [
    ("跨工况覆盖足",
     "24 个不同过渡, 长度 903 - 13006 步 ⇒ ICL context 在所有规模上都能验证.",
     "→ 这给 TSFM『多模式样本量』预测提供了样本空间.",
     ACCENT),
    ("OP / PV 同步采集",
     "16 个变量同时采集, OP (因) 和 PV (果) 都齐全.",
     "→ 可以直接做有 / 无 OP 的对照 (任务 B 的 hold-out).",
     GREEN),
    ("Δ-边际近似 i.i.d. 这一现象本身有理论意义",
     "工业受控系统, 闭环控制把 PV 拉回 setpoint ⇒ Δ 在长程上是 mean-reverting ⇒ KS 差异小.",
     "→ 这是 控制工程引入的 Δ-stationarity, 不是巧合.",
     PURPLE),
    ("Permutation entropy ≈ 0.98 暗示 Δ 接近 maximum entropy",
     "PE 饱和 ⇒ Δ 序列的『序号信息』几乎没有可压缩性.",
     "→ TSFM 在 Δ 上学不到太多 — 它学的是 patch-level 的尺度+趋势, 不是 ordinal.",
     TEAL),
    ("Hurst 0.4-0.55 + VR(60) < 0.5 ⇒ 长程反持续",
     "Mean-reverting 行为是 RevIN 假设最符合的设定.",
     "→ 工业闭环回路天然产生这种结构, 这就是 TSFM 在工业 TS 上 work 的内在原因.",
     ORANGE),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(12.3); ch = Inches(1.05)
for i, (title, body1, body2, col) in enumerate(reasons):
    y = gy + i * (ch + Inches(0.05))
    add_rect(s, gx, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, gx, y, Inches(0.2), ch, col)
    add_text(s, gx + Inches(0.35), y + Inches(0.08), Inches(3.6), Inches(0.4),
             title, size=12, bold=True, color=col)
    add_text(s, gx + Inches(0.35), y + Inches(0.5), Inches(11.7), Inches(0.3),
             body1, size=11, color=DGREY)
    add_text(s, gx + Inches(0.35), y + Inches(0.78), Inches(11.7), Inches(0.3),
             body2, size=11, italic=True, color=NAVY)

# ============================================================
# Slide 16 — Counterintuitive insight
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 16, TOTAL, "反直觉洞察 · 闭环控制 让 TSFM 单变量更强",
     "为什么? — 一个工业 TSFM 工作的隐藏机制",
     section="PART VII · 数据的内在含义")

# Top: claim
add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.9), DEEP)
add_text(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.9),
         "闭环 PID/控制器  ⇒  把 PV 拉回 SP  ⇒  PV 在 Δ 空间『被人工 stationary 化』",
         size=15, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Mechanism breakdown
add_rect(s, Inches(0.5), Inches(2.65), Inches(6.0), Inches(2.4), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(2.65), Inches(6.0), Inches(0.55), GREY)
add_text(s, Inches(0.5), Inches(2.65), Inches(6.0), Inches(0.55),
         "开环 (无控制) 系统", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(0.7), Inches(3.3), Inches(5.7), Inches(1.7), [
    [("• 输入扰动 ⇒ 输出长期漂移",
      12, DGREY, False)],
    [("• Δ 在长程上 持续 (Hurst > 0.5)",
      12, DGREY, False)],
    [("• 跨工况 边际差异巨大",
      12, DGREY, False)],
    [("• KPSS 拒绝 stationary",
      12, DGREY, False)],
    [("→ ", 12, RED, True),
     ("RevIN 假设破坏 ⇒ TSFM 单变量预测漂移失控.",
      12, RED, True)],
])

add_rect(s, Inches(6.8), Inches(2.65), Inches(6.0), Inches(2.4), WHITE, line=SOFT)
add_rect(s, Inches(6.8), Inches(2.65), Inches(6.0), Inches(0.55), GREEN)
add_text(s, Inches(6.8), Inches(2.65), Inches(6.0), Inches(0.55),
         "闭环 (PID 控制) 系统 — ZJSH 实际", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(7.0), Inches(3.3), Inches(5.7), Inches(1.7), [
    [("• 控制器 把 PV 拉回 SP (闭环负反馈)",
      12, DGREY, False)],
    [("• Δ 长程上 反持续 (Hurst < 0.5)",
      12, DGREY, False)],
    [("• 跨工况边际近似一致 (KS 19%)",
      12, DGREY, False)],
    [("• ADF + KPSS 一致 stationary",
      12, DGREY, False)],
    [("→ ", 12, GREEN, True),
     ("RevIN 假设近似成立 ⇒ TSFM 单变量预测稳定.",
      12, GREEN, True)],
])

# Punchline box
add_rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.85), LIGHT)
add_rect(s, Inches(0.5), Inches(5.2), Inches(0.18), Inches(1.85), ACCENT)
add_text(s, Inches(0.85), Inches(5.3), Inches(11.8), Inches(0.4),
         "深层结论", size=14, bold=True, color=NAVY)
rich(s, Inches(0.85), Inches(5.7), Inches(11.8), Inches(1.3), [
    [("现代化工厂的 ", 13, NAVY, False),
     ("闭环控制结构", 13, ACCENT, True),
     (" 是 TSFM 在工业 TS 上 work 的 ",
      13, NAVY, False),
     ("隐藏前提条件", 13, ACCENT, True),
     (" — 它把原本 non-stationary 的物理过程加工成了对 TSFM 友好的 stationary-Δ 数据.",
      13, NAVY, False)],
    [("", 6, NAVY, False)],
    [("→ 在 ", 12, NAVY, False),
     ("缺乏闭环控制的工艺", 12, RED, True),
     (" (例如某些化学反应器、生物反应器), TSFM 单变量的优势会大打折扣 ⇒ 必须用 multivariate + 物理模型混合.",
      12, NAVY, False)],
])

# ============================================================
# Slide 17 — Risk inventory
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 17, TOTAL, "风险清单 · 单变量 TSFM 在 ZJSH 上可能出错的地方",
     "哪怕 Task A 适合单变量, 工程上仍要警惕这些情况",
     section="PART VIII · 风险与边界")

risks = [
    ("R1 · 训练数据中过渡形态有限",
     "ZJSH 24 个过渡都来自类似工艺. 真正『新型』过渡 (例: 操作员失误 / 设备故障) 不在训练分布.",
     "缓解: 在线 Mahalanobis 残差监控 (BGMM thresh) 触发降级模式.",
     RED),
    ("R2 · OP 突变期 (Δ 不再 i.i.d.)",
     "OP 改变后的 1-3 步, PV 的 Δ 出现『突跳』, 暂时违反 marginal-stationary 假设.",
     "缓解: 检测到 OP 跳变 ⇒ 在该 horizon 上拒绝 TSFM 单变量预测, 切到 multivariate.",
     ORANGE),
    ("R3 · 测量噪声 / 数据缺失",
     "若 PI10102 传感器有间歇性故障, 单变量的 ICL 会被噪声历史误导.",
     "缓解: 采用 Robust TSFM 变体 (Lag-Llama, Chronos-Bolt) + 异常值过滤前处理.",
     PURPLE),
    ("R4 · Conformal 校准的 i.i.d. 假设",
     "标准 conformal 要求 calibration 与 test 同分布; 跨过渡 marginal 弱差异可能轻微违反.",
     "缓解: 用 Mondrian conformal (按工况分组) 或 adaptive conformal (在线更新).",
     TEAL),
    ("R5 · 长 horizon 预测的复合误差",
     "TSFM 自回归预测在 horizon > 96 步时误差累积 ⇒ Hurst 略 > 0.5 工况上更明显.",
     "缓解: 限制 horizon ≤ 60 步; 长 horizon 切换到分级预测 (粗-细 cascade).",
     ACCENT),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(12.3); ch = Inches(1.05)
for i, (title, body, fix, col) in enumerate(risks):
    y = gy + i * (ch + Inches(0.05))
    add_rect(s, gx, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, gx, y, Inches(0.2), ch, col)
    add_text(s, gx + Inches(0.35), y + Inches(0.08), Inches(11.7), Inches(0.4),
             title, size=12, bold=True, color=col)
    add_text(s, gx + Inches(0.35), y + Inches(0.45), Inches(11.7), Inches(0.32),
             body, size=11, color=DGREY)
    add_text(s, gx + Inches(0.35), y + Inches(0.78), Inches(11.7), Inches(0.3),
             fix, size=11, italic=True, bold=True, color=NAVY)

# ============================================================
# Slide 18 — Compare to previous reactor PPT
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 18, TOTAL, "和『反应器场景』PPT 的关系 · 看似矛盾, 实则统一",
     "前 PPT 说『单变量不够』, 现在说『单变量够用』 — 是不是矛盾?",
     section="PART VIII · 风险与边界")

add_rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(2.95), WHITE, line=SOFT)
add_rect(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.55), DEEP)
add_text(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.55),
         "反应器场景 PPT (前)", size=14, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(0.7), Inches(2.25), Inches(5.7), Inches(2.2), [
    [("场景: 反应器, 进料组分 c(t) 突变.",
      12, NAVY, True)],
    [("观测: 仅温度 T(t).",
      11, DGREY, False)],
    [("结论: ", 12, NAVY, True),
     ("Stark 1999 — c(t) 不可识别", 12, RED, True),
     (".", 11, DGREY, False)],
    [("→ 单变量不够 是 ", 11, DGREY, False),
     ("数学上必然", 11, RED, True), (".", 11, DGREY, False)],
    [("", 6, NAVY, False)],
    [("关键差异: ", 12, NAVY, True),
     ("反应器场景 隐含 假设组分是 不可测的 外生扰动",
      11, DGREY, False)],
    [("此时 σ²(T(t+h) | T_past) ≫ 0 (硬下界).",
      11, GREY, False)],
])

add_rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(2.95), WHITE, line=SOFT)
add_rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.55), DEEP)
add_text(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.55),
         "ZJSH 多工况过渡 (本 PPT)", size=14, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rich(s, Inches(7.0), Inches(2.25), Inches(5.7), Inches(2.2), [
    [("场景: 工厂, OP 变化驱动工况切换.",
      12, NAVY, True)],
    [("观测: 16 变量 (PV + OP).",
      11, DGREY, False)],
    [("结论: 任务 A (内部轨迹) ", 12, NAVY, True),
     ("单变量够用", 12, GREEN, True),
     (";  任务 B (预警) 必须 multivariate.",
      11, DGREY, False)],
    [("→ 任务 A 在 Task A 上 ", 11, DGREY, False),
     ("单变量统计上充分", 11, GREEN, True),
     (".", 11, DGREY, False)],
    [("", 6, NAVY, False)],
    [("关键差异: ", 12, NAVY, True),
     ("ZJSH 中 已经进入过渡 后, ", 11, DGREY, False),
     ("OP 干预的影响已编码在 PV 历史里",
      11, ACCENT, True)],
    [("ICL 可以从 PI10102 context 反推 OP 信号.",
      11, GREY, False)],
])

# Synthesis
add_rect(s, Inches(0.5), Inches(4.65), Inches(12.3), Inches(2.4), DEEP)
add_text(s, Inches(0.85), Inches(4.8), Inches(11.8), Inches(0.4),
         "统一原理: 单变量是否够用 = 因 (上游驱动) 是否在 PV 历史中已经留痕",
         size=14, bold=True, color=ACCENT)
rich(s, Inches(0.85), Inches(5.25), Inches(11.8), Inches(1.7), [
    [("• ", 13, ACCENT, True),
     ("过渡 内部 (PV 已响应):", 13, WHITE, True),
     (" 上游 OP 已经把信号写进 PI10102 历史 ⇒ 单变量充分",
      12, SOFT, False)],
    [("• ", 13, ACCENT, True),
     ("过渡 之前 (PV 未响应) 的预警:", 13, WHITE, True),
     (" 上游 OP 信号尚未传到 PV ⇒ 单变量不充分, 必须看 OP",
      12, SOFT, False)],
    [("• ", 13, ACCENT, True),
     ("反应器场景 (组分不可测):", 13, WHITE, True),
     (" 上游驱动 物理上 不存在于观测中 ⇒ 任何时候单变量都不充分",
      12, SOFT, False)],
    [("", 6, WHITE, False)],
    [("→ 三个 PPT 的结论可以用 ", 13, SOFT, False),
     ("一个公式", 13, ACCENT, True),
     (" 统一: 看『因果链上游』是否进入了观测序列.",
      13, SOFT, False)],
])

# ============================================================
# Slide 19 — Engineering deployment plan
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 19, TOTAL, "工程部署 pipeline · 从 ZJSH 到生产线",
     "6 阶段, 8-10 周首版上线",
     section="PART IX · 工程")

phases = [
    ("Phase 1\n数据 audit",
     "1 周",
     "• 确认 24 过渡的标签质量\n• OP 与 PV 的对齐 / 时延校准\n• 划分 hold-out 6 过渡",
     RGBColor(0x2E, 0x60, 0xA8)),
    ("Phase 2\nbaseline",
     "1 周",
     "• 单变量 zero-shot TSFM\n• ARIMA / 上一值\n• 输出 MAE / coverage 基线",
     RGBColor(0xC0, 0x6A, 0x1F)),
    ("Phase 3\nTask A 头",
     "2 周",
     "• 单变量 TSFM + LoRA\n• Conformal 校准\n• Hold-out 验证",
     RGBColor(0x4F, 0x8A, 0x44)),
    ("Phase 4\nTask B 头",
     "2 周",
     "• TimeXer + Tier1+2 OP\n• 二分类 head\n• ROC-AUC / lead time",
     RGBColor(0x8E, 0x44, 0xAD)),
    ("Phase 5\nTask C 监控",
     "1 周",
     "• 残差 + Mahalanobis\n• 阈值校准\n• Fallback 触发",
     RGBColor(0x16, 0x7E, 0x8A)),
    ("Phase 6\n上线 + 监控",
     "持续",
     "• 实时残差告警\n• 每月增量 LoRA fine-tune\n• 性能 dashboard",
     RGBColor(0xC0, 0x39, 0x2B)),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(2.05); ch = Inches(4.5)
for i, (name, time_, body, color) in enumerate(phases):
    x = gx + i * (cw + Inches(0.05))
    add_rect(s, x, gy, cw, ch, WHITE, line=SOFT)
    add_rect(s, x, gy, cw, Inches(1.05), color)
    add_text(s, x + Inches(0.1), gy + Inches(0.1), cw - Inches(0.2), Inches(0.85),
             name, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x + Inches(0.15), gy + Inches(1.15), cw - Inches(0.3), Inches(0.4), ACCENT)
    add_text(s, x + Inches(0.15), gy + Inches(1.15), cw - Inches(0.3), Inches(0.4),
             "时长: " + time_, size=10, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.15), gy + Inches(1.7), cw - Inches(0.3), Inches(2.7),
             body, size=10, color=DGREY)
    if i < len(phases) - 1:
        arrow(s, x + cw, gy + ch // 2, x + cw + Inches(0.05), gy + ch // 2,
              color=ACCENT, w=2.0)

add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
         "总周期 ~ 8-10 周到首版上线; 三任务 head 共享 backbone, 训练成本可控.",
         size=12, italic=True, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 20 — Expected gains
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 20, TOTAL, "预期效果 · 基于 ZJSH 实证 + 文献基准",
     "三个任务上, 相对 baseline 的合理改善幅度",
     section="PART IX · 工程")

gains = [
    ("Task A · 内部轨迹",
     "vs ARIMA / persistence", "MAE ↓ 35-55%",
     "vs supervised PatchTST", "MAE ↓ 5-15%",
     "覆盖率 P10/P90", "0.85 - 0.92",
     GREEN),
    ("Task B · 早期预警",
     "vs PV-only baseline", "AUC: 0.62 → 0.85+",
     "vs Tier1 OP only", "AUC +5-10 点",
     "Lead time", "5 - 30 分钟",
     RED),
    ("Task C · OOD 检测",
     "vs 无监控", "FP/月: ~0;  漏报率 < 5%",
     "残差阈值", "Mahalanobis P95",
     "降级触发率", "< 2% 时间",
     ORANGE),
]
gx = Inches(0.5); gy = Inches(1.55); cw = Inches(12.3); ch = Inches(1.6)
for i, (task, lab1, val1, lab2, val2, lab3, val3, col) in enumerate(gains):
    y = gy + i * (ch + Inches(0.1))
    add_rect(s, gx, y, cw, ch, WHITE, line=SOFT)
    add_rect(s, gx, y, Inches(2.7), ch, col)
    add_text(s, gx, y, Inches(2.7), ch, task, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 3 metric blocks
    metrics = [(lab1, val1), (lab2, val2), (lab3, val3)]
    mx = gx + Inches(2.85); mw = (cw - Inches(2.85) - Inches(0.1)) / 3
    for j, (lab, val) in enumerate(metrics):
        x_m = mx + j * mw + Inches(0.05)
        add_rect(s, x_m, y + Inches(0.1), mw - Inches(0.1), ch - Inches(0.2), LIGHT)
        add_text(s, x_m, y + Inches(0.2), mw - Inches(0.1), Inches(0.45),
                 lab, size=11, italic=True, color=GREY,
                 align=PP_ALIGN.CENTER)
        add_text(s, x_m, y + Inches(0.7), mw - Inches(0.1), Inches(0.7),
                 val, size=15, bold=True, color=col,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5), LIGHT)
add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5),
         "→ 数字基于 GIFT-Eval / Chronos / Moirai 公开实证;  实际值需在 ZJSH hold-out 上确认.",
         size=12, italic=True, color=NAVY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# Slide 21 — Decision flow / rule of thumb
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 21, TOTAL, "一句话决策规则 · 单变量 vs 多变量",
     "基于 ZJSH 实证总结的简单工程规则",
     section="PART X · 决策")

# Decision tree
add_rect(s, Inches(2.5), Inches(1.55), Inches(8.3), Inches(0.7), DEEP)
add_text(s, Inches(2.5), Inches(1.55), Inches(8.3), Inches(0.7),
         "Q1: 你要预测什么?",
         size=15, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Branch lines
arrow(s, Inches(4.2), Inches(2.25), Inches(2.5), Inches(3.0), color=GREEN)
arrow(s, Inches(6.65), Inches(2.25), Inches(6.65), Inches(3.0), color=ACCENT)
arrow(s, Inches(9.1), Inches(2.25), Inches(10.8), Inches(3.0), color=PURPLE)

add_text(s, Inches(2.7), Inches(2.35), Inches(2.0), Inches(0.4),
         "已进入过渡的轨迹", size=11, italic=True, bold=True, color=GREEN)
add_text(s, Inches(5.7), Inches(2.35), Inches(2.0), Inches(0.4),
         "稳态时预警", size=11, italic=True, bold=True, color=ACCENT)
add_text(s, Inches(9.4), Inches(2.35), Inches(2.5), Inches(0.4),
         "新工况 / OOD", size=11, italic=True, bold=True, color=PURPLE)

# Three answer boxes
boxes = [
    ("✓ 单变量 TSFM",
     "PI10102 univariate\nzero-shot 或 LoRA\n+ Conformal CI",
     "ZJSH KS=18.84% 支持\nADF/KPSS 全 stationary",
     GREEN, Inches(1.0)),
    ("Multivariate",
     "TimeXer / Moirai\nPI10102 + 5-9 个 OP\n二分类 head",
     "OP KS 75-87% 支持\nLeading 因果链",
     ACCENT, Inches(5.15)),
    ("TSFM + 监控",
     "Backbone + 残差监控\nMahalanobis 阈值\nFallback 机制",
     "BGMM 思路用到在线\n+ 物理模型兜底",
     PURPLE, Inches(9.3)),
]
for title, body, just, col, x in boxes:
    add_rect(s, x, Inches(3.0), Inches(3.0), Inches(2.7), WHITE, line=col, line_w=Pt(2))
    add_rect(s, x, Inches(3.0), Inches(3.0), Inches(0.55), col)
    add_text(s, x, Inches(3.0), Inches(3.0), Inches(0.55),
             title, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.15), Inches(3.7), Inches(2.7), Inches(1.2),
             body, size=11, color=NAVY, align=PP_ALIGN.CENTER)
    add_rect(s, x + Inches(0.15), Inches(5.0), Inches(2.7), Inches(0.65), LIGHT)
    add_text(s, x + Inches(0.15), Inches(5.0), Inches(2.7), Inches(0.65),
             just, size=10, italic=True, color=DGREY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Bottom rule
add_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.05), DEEP)
add_text(s, Inches(0.85), Inches(6.1), Inches(11.8), Inches(0.4),
         "Rule of thumb",
         size=14, bold=True, color=ACCENT)
rich(s, Inches(0.85), Inches(6.5), Inches(11.8), Inches(0.55), [
    [("如果 ", 13, WHITE, False),
     ("上游驱动信号已经在 PV 历史中留下印迹", 13, ACCENT, True),
     (" → 单变量 TSFM 够;",
      13, WHITE, False)],
    [("如果 ", 13, WHITE, False),
     ("上游驱动尚未传到 PV (预警场景) 或 物理上不可观测 (反应器场景)",
      13, ACCENT, True),
     (" → 必须 multivariate.",
      13, WHITE, False)],
])

# ============================================================
# Slide 22 — Final answer
# ============================================================
s = prs.slides.add_slide(BLANK)
page(s, 22, TOTAL, "最终回答你的两个问题",
     "基于 ZJSH 实证 + 三任务分解 + 因果链分析",
     section="结论")

# Three answer blocks
answers = [
    ("Q1", "TSFM 在多工况过渡过程预测预警中 是否有潜力?",
     "✓ 有, 而且潜力很大 — 但必须分解任务后看.",
     "ZJSH 实证显示: PI10102 在 Δ 空间跨工况近似 i.i.d. (KS=18.84%) — 这正好是 TSFM 的 sweet spot. " +
     "ICL + RevIN + Patching + 闭环控制带来的 stationary-Δ 协同发挥效用. " +
     "三任务: A (内部轨迹) ✓ 高潜力;  B (预警) ✓ 高潜力;  C (新工况) △ 中等潜力.",
     GREEN),
    ("Q2", "单变量 TSFM (主流模式) 此时是否更合适?",
     "答案分场景: 任务 A ✓ 是;  任务 B ✗ 否;  任务 C △ 部分.",
     "任务 A 上 单变量是最优的, 因为 OP 的影响已经写进 PV 历史 + ICL 自动识别 regime + Channel-Indep 放大有效数据 16×. " +
     "任务 B 上 必须用 multivariate (Tier 1+2 OP 协变量), 因为预警必须看上游因果信号. " +
     "任务 C 上 单变量 TSFM 比 supervised 韧性强, 但需要外加 OOD 检测层.",
     ACCENT),
    ("∑", "为什么这是合理的而不是矛盾的?",
     "因果链统一原则: 看上游驱动是否已进入观测序列.",
     "若驱动已经在 PV 历史中留下印迹 (内部轨迹) → 单变量充分;  " +
     "若驱动尚未传到 PV (预警) → multivariate 必要;  " +
     "若驱动物理上不可观测 (反应器组分) → multivariate 仍不够, 需加传感器. " +
     "三个 PPT 的结论由这一原则统一.",
     PURPLE),
]
gy = Inches(1.55)
for i, (q, sub, head, body, col) in enumerate(answers):
    y = gy + i * Inches(1.85)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.7), WHITE, line=col, line_w=Pt(2))
    add_rect(s, Inches(0.5), y, Inches(0.85), Inches(1.7), col)
    add_text(s, Inches(0.5), y, Inches(0.85), Inches(1.7),
             q, size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.5), y + Inches(0.1), Inches(11.0), Inches(0.4),
             sub, size=12, bold=True, color=NAVY, italic=True)
    add_text(s, Inches(1.5), y + Inches(0.5), Inches(11.0), Inches(0.45),
             head, size=14, bold=True, color=col)
    add_text(s, Inches(1.5), y + Inches(1.0), Inches(11.0), Inches(0.65),
             body, size=11, color=DGREY)

# Save
out = "/home/aicode/sherwin/TSFM/paper/TSFM_多工况过渡过程预测预警.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
